
# Imports
import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

base_path = r"C:\Users\ext-yaissa\OneDrive - INFOGREFFE\Documents\National"

files_needed = [
    "Modèle Stats France.pptx",
    "Immat_2024_t3.xlsx",
    "Immat_2025_t3.xlsx",
    "Pcl_2024_t3.xlsx",
    "Pcl_2025_t3.xlsx",
    "Radiation_2024_t3.xlsx",
    "Radiation_2025_t3.xlsx",
    "StockRCS_2025_t3_au15092025.xlsx",
    "InjPayer_2024_t3.xlsx",
    "InjPayer_2025_t3.xlsx",
    "TrfSiege_2024_t3.xlsx",
    "TrfSiege_2025_t3.xlsx",
    "Dpca_2024_t3.xlsx",
    "Dpca_2025_t3.xlsx",
    "SteESS_2024_t3.xlsx",
    "SteESS_2025_t3.xlsx",
    "SteMission_2024_t3.xlsx",
    "SteMission_2025_t3.xlsx",
    "Mbe_2024_t3.xlsx",
    "Mbe_2025_t3.xlsx",
    "InscrSSTR_2024_t3.xlsx",
    "InscrSSTR_2025_t3.xlsx",
    "Div_2024_t3.xlsx",
    "Div_2025_t3.xlsx",
    "Prev_2024_t3.xlsx",
    "Prev_2025_t3.xlsx",
]

# Vérification des fichiers
missing_files = [f for f in files_needed if not os.path.exists(os.path.join(base_path, f))]
if missing_files:
    print("Fichiers manquants :", missing_files)
    print("Veuillez les ajouter dans le dossier :", base_path)
else:
    print("Tous les fichiers sont présents.")

# === CHARGEMENT DES DONNÉES ===
df_immat_2024 = pd.read_excel(os.path.join(base_path, "Immat_2024_t3.xlsx"), sheet_name="Immat_2024_t3")
df_immat_2025 = pd.read_excel(os.path.join(base_path, "Immat_2025_t3.xlsx"), sheet_name="Immat_2025_t3")

df_pcl_2024 = pd.read_excel(os.path.join(base_path, "Pcl_2024_t3.xlsx"), sheet_name="Pcl_2024_t3")
df_pcl_2025 = pd.read_excel(os.path.join(base_path, "Pcl_2025_t3.xlsx"), sheet_name="Pcl_2025_t3")

df_stock_2025 = pd.read_excel(os.path.join(base_path, "StockRCS_2025_t3_au15092025.xlsx"), sheet_name="StockRCS_2025_t3_au15092025")

df_radiation_2024 = pd.read_excel(os.path.join(base_path, "Radiation_2024_t3.xlsx"), sheet_name="Radiation_2024_t3")
df_radiation_2025 = pd.read_excel(os.path.join(base_path, "Radiation_2025_t3.xlsx"), sheet_name="Radiation_2025_t3")

df_ip_2024 = pd.read_excel(os.path.join(base_path, "InjPayer_2024_t3.xlsx"), sheet_name="InjPayer_2024_t3")
df_ip_2025 = pd.read_excel(os.path.join(base_path, "InjPayer_2025_t3.xlsx"), sheet_name="InjPayer_2025_t3")

df_ts_2024 = pd.read_excel(os.path.join(base_path, "TrfSiege_2024_t3.xlsx"), sheet_name="TrfSiege_2024_t3")
df_ts_2025 = pd.read_excel(os.path.join(base_path, "TrfSiege_2025_t3.xlsx"), sheet_name="TrfSiege_2025_t3")

df_ess_2024 = pd.read_excel(os.path.join(base_path, "SteESS_2024_t3.xlsx"), sheet_name="SteESS_2024_t3")
df_ess_2025 = pd.read_excel(os.path.join(base_path, "SteESS_2025_t3.xlsx"), sheet_name="SteESS_2025_t3")

df_miss_2024 = pd.read_excel(os.path.join(base_path, "SteMission_2024_t3.xlsx"), sheet_name="SteMission_2024_t3")
df_miss_2025 = pd.read_excel(os.path.join(base_path, "SteMission_2025_t3.xlsx"), sheet_name="SteMission_2025_t3")

df_mbe_2024 = pd.read_excel(os.path.join(base_path, "Mbe_2024_t3.xlsx"), sheet_name="Mbe_2024_t3")
df_mbe_2025 = pd.read_excel(os.path.join(base_path, "Mbe_2025_t3.xlsx"), sheet_name="Mbe_2025_t3")

df_dpca_2024 = pd.read_excel(os.path.join(base_path, "Dpca_2024_t3.xlsx"), sheet_name="Dpca_2024_t3")
df_dpca_2025 = pd.read_excel(os.path.join(base_path, "Dpca_2025_t3.xlsx"), sheet_name="Dpca_2025_t3")

df_sstr_2024 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2024_t3.xlsx"), sheet_name="InscrSSTR_2024_t3")
df_sstr_2025 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2025_t3.xlsx"), sheet_name="InscrSSTR_2025_t3")

df_div_2024 = pd.read_excel(os.path.join(base_path, "Div_2024_t3.xlsx"), sheet_name="Div_2024_t3")
df_div_2025 = pd.read_excel(os.path.join(base_path, "Div_2025_t3.xlsx"), sheet_name="Div_2025_t3")

df_prev_2024 = pd.read_excel(os.path.join(base_path, "Prev_2024_t3.xlsx"), sheet_name="Prev_2024_t3")
df_prev_2025 = pd.read_excel(os.path.join(base_path, "Prev_2025_t3.xlsx"), sheet_name="Prev_2025_t3")

# Filtrage et dédoublonnements des fichiers excel
def preprocess(df):
    dedup_cols = ["GRF", "MIL", "STC", "CHRONO"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_immat_2024, df_immat_2025 = preprocess(df_immat_2024), preprocess(df_immat_2025)
df_pcl_2024, df_pcl_2025 = preprocess(df_pcl_2024), preprocess(df_pcl_2025)
df_radiation_2024, df_radiation_2025 = preprocess(df_radiation_2024), preprocess(df_radiation_2025)
df_ess_2024, df_ess_2025 = preprocess(df_ess_2024), preprocess(df_ess_2025)
df_miss_2024, df_miss_2025 = preprocess(df_miss_2024), preprocess(df_miss_2025)
df_mbe_2024, df_mbe_2025 = preprocess(df_mbe_2024), preprocess(df_mbe_2025)
df_dpca_2024, df_dpca_2025 = preprocess(df_dpca_2024), preprocess(df_dpca_2025)

def preprocess3(df):
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_stock_2025 = preprocess3(df_stock_2025)

def preprocess2(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_AFFAIRE"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_ip_2024, df_ip_2025 = preprocess2(df_ip_2024), preprocess2(df_ip_2025)

def preprocess4(df):
    dept_cols = ["DEPFROM", "REGIONTO"]
    for col in dept_cols:
        df[col] = df[col].astype(str).str.strip().str.title()
    exclure_depts = ["Bas Rhin", "Haut Rhin", "Moselle"]
    for col in dept_cols:
        df = df[~df[col].isin(exclure_depts)]
    df = df.reset_index(drop=True)
    return df

df_ts_2024, df_ts_2025 = preprocess4(df_ts_2024), preprocess4(df_ts_2025)

def preprocess6(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_sstr_2024, df_sstr_2025 = preprocess6(df_sstr_2024), preprocess6(df_sstr_2025)

def preprocess9(df):
    dept_col = "DEPARTEMENT"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_prev_2024, df_prev_2025 = preprocess9(df_prev_2024), preprocess9(df_prev_2025)

# Séparation Sécurité Sociale et Trésor - Tésor 
def preprocess5(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_tresor = df["libelle inscription"].str.contains("trésor|tresor", case=False, na=False)
    df = df[mask_tresor].reset_index(drop=True)

    return df

# Sécurité sociale 
def preprocess7(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes sécurité sociale
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_ss = df["libelle inscription"].str.contains("sécurité sociale|securite sociale|sécurité|securite", case=False, na=False)
    df = df[mask_ss].reset_index(drop=True)

    return df

df_sstr_tresor_2024 = preprocess5(df_sstr_2024)
df_sstr_tresor_2025 = preprocess5(df_sstr_2025)
df_sstr_ss_2024 = preprocess7(df_sstr_2024)
df_sstr_ss_2025 = preprocess7(df_sstr_2025)

# Ouvrir le ppt
prs = Presentation(os.path.join(base_path, "Modèle Stats France.pptx"))

#   PARTIE 1 : IMMATRICULATIONS  

# === SLIDE 3 ===
slide3 = prs.slides[2]

A_2024 = df_immat_2024["SIREN"].notna().sum()
A_2025 = df_immat_2025["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024["SIREN"].notna().sum()
C_2025 = df_pcl_2025["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024["SIREN"].notna().sum()
D_2025 = df_radiation_2025["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024["SIREN"].notna().sum()
E_2025 = df_ip_2025["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer dans la slide 3
for shape in slide3.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 4 ===
slide4 = prs.slides[3]

# A : MBE
A_2024_slide3 = df_mbe_2024["SIREN"].notna().sum()
A_2025_slide3 = df_mbe_2025["SIREN"].notna().sum()
A_val_slide3 = f"{A_2025_slide3:,}".replace(",", " ")
A_pct_slide3 = evol_percent_txt(A_2024_slide3, A_2025_slide3)

# B : DIV
B_2024_slide3 = df_div_2024["SIREN"].notna().sum()
B_2025_slide3 = df_div_2025["SIREN"].notna().sum()
B_val_slide3 = f"{B_2025_slide3:,}".replace(",", " ")
B_pct_slide3 = evol_percent_txt(B_2024_slide3, B_2025_slide3)

# C : DPCA
C_2024_slide3 = df_dpca_2024["SIREN"].notna().sum()
C_2025_slide3 = df_dpca_2025["SIREN"].notna().sum()
C_val_slide3 = f"{C_2025_slide3:,}".replace(",", " ")
C_pct_slide3 = evol_percent_txt(C_2024_slide3, C_2025_slide3)

# Remplacement du texte dans la slide 4
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            if text == "A":
                run.text = A_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER

            else:
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide3, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide3, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide3, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide3, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide3, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide3, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

def fmt_nb(val_2025):
    return f"{val_2025:,}".replace(",", " ")

val_2024 = df_immat_2024["SIREN"].notna().sum()
val_2025 = df_immat_2025["SIREN"].notna().sum()
variation = ((val_2025 - val_2024) / val_2024) * 100 if val_2024 > 0 else 0

formatted_val = fmt_nb(val_2025)
formatted_var = f"{variation:+.1f}%".replace("++", "+")

# === SLIDE 6 ===
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text in ["X%", "+X%", "-X%"]:
                    run.text = formatted_var
                    paragraph.alignment = PP_ALIGN.CENTER
                elif text == "X":
                    run.text = formatted_val
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 7 ===
slide7 = prs.slides[6]
table = None
for shape in slide7.shapes:
    if shape.has_table:
        table = shape.table
        break

if table is None:
    raise ValueError("Aucun tableau trouvé dans la slide 7.")

region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

df_non_vides = df_immat_2025[df_immat_2025["SIREN"].astype(str).str.strip() != ""]
immat_region_2025 = df_non_vides[region_col].value_counts().sort_index()
immat_region_2025 = immat_region_2025[~immat_region_2025.index.isin(regions_exclues)]
total_immat_2025 = df_non_vides.shape[0]
parts_region = (immat_region_2025 / total_immat_2025 * 100).round(1)

parts_region_str = {region.split(' (')[0]: f"{val}%" for region, val in parts_region.items()}
region_part_tuples = sorted(parts_region_str.items(), key=lambda x: -float(x[1].replace('%','')))

nb_lignes_tableau = len(table.rows)
nb_lignes_donnees = len(region_part_tuples)

if nb_lignes_donnees + 1 > nb_lignes_tableau:
    print(f"Le tableau ne contient pas assez de lignes ({nb_lignes_tableau}) pour {nb_lignes_donnees} régions.")

table.cell(0, 0).text = "Régions"
table.cell(0, 1).text = "Répartition (en %)"

for i, (region, part) in enumerate(region_part_tuples[:nb_lignes_tableau - 1]):
    table.cell(i + 1, 0).text = region
    table.cell(i + 1, 1).text = part

for row in table.rows:
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"


# === SLIDE 8 ===
slide8 = prs.slides[7]

# Calcul répartition par secteur d'activité 
secteur_col = "SECTEURLIBELLE"  

# Vérification colonnes
for df_name, df in [("df_immat_2024", df_immat_2024), ("df_immat_2025", df_immat_2025)]:
    if secteur_col not in df.columns:
        raise KeyError(f"Colonne '{secteur_col}' absente dans {df_name}")

# Normalisation
df_immat_2024[secteur_col] = df_immat_2024[secteur_col].astype(str).str.strip()
df_immat_2025[secteur_col] = df_immat_2025[secteur_col].astype(str).str.strip()

# Comptages par secteur
secteurs_2024 = df_immat_2024[secteur_col].value_counts()
secteurs_2025 = df_immat_2025[secteur_col].value_counts()

# Liste des secteurs présents dans le tableau du modèle
secteurs_selectionnes = []
for shape in slide8.shapes:
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for row_idx in range(len(table.rows)):
            secteur_nom = table.cell(row_idx, 0).text.strip()
            if secteur_nom:
                secteurs_selectionnes.append(secteur_nom)

# Dictionnaire des valeurs
secteurs_resume = {}
for secteur in secteurs_selectionnes:
    v24 = int(secteurs_2024.get(secteur, 0))
    v25 = int(secteurs_2025.get(secteur, 0))
    evol_pct = ((v25 - v24) / v24 * 100) if v24 > 0 else 0.0
    secteurs_resume[secteur] = {"2024": v24, "2025": v25, "evol_pct": evol_pct}

# Remplissage du tableau dans la slide 8 
for shape in slide8.shapes:
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for row_idx in range(len(table.rows)):
            secteur_nom = table.cell(row_idx, 0).text.strip()
            if secteur_nom in secteurs_resume:
                val2025 = secteurs_resume[secteur_nom]["2025"]
                evol_pct = secteurs_resume[secteur_nom]["evol_pct"]
                table.cell(row_idx, 1).text = f"{val2025:,}".replace(",", " ")
                table.cell(row_idx, 2).text = f"{evol_pct:+.1f}%"

# Diagramme secteur
for shape in list(slide8.shapes):
    if hasattr(shape, "has_chart") and shape.has_chart:
        slide8.shapes._spTree.remove(shape._element)

total_2025 = sum(val["2025"] for val in secteurs_resume.values())
parts_2025 = [
    (round(val["2025"] / total_2025 * 100, 1) if total_2025 > 0 else 0.0)
    for val in secteurs_resume.values()
]

chart_data = CategoryChartData()
chart_data.categories = list(secteurs_resume.keys())
chart_data.add_series("", parts_2025)

# Position et taille du graphique
x, y, cx, cy = Inches(0.4), Inches(1.0), Inches(4.8), Inches(7.3)
chart = slide8.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.bold = False
data_labels.show_category_name = True
data_labels.show_value = False
data_labels.show_percentage = True

# --- Épaisseur de l’anneau ---
try:
    plot.doughnut_hole_size = 70  # pris en charge par certaines versions de python-pptx
except AttributeError:
    pass

# Fallback XML (nécessite lxml)
try:
    from lxml import etree
    doughnut_chart = chart.plots[0]._element
    hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
    if hole_size is None:
        hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
        doughnut_chart.append(hole_size)
    hole_size.set('val', '70')  # 10..90
except Exception:
    # si lxml non dispo ou structure différente, on ignore silencieusement
    pass


# === SLIDE 9 ===
slide9 = prs.slides[8]

# --- Calcul répartition par forme juridique ---
forme_col = "FORMEJURIDIQUELIBELLEAGREGE"  

# Vérification colonnes
for df_name, df in [("df_immat_2024", df_immat_2024), ("df_immat_2025", df_immat_2025)]:
    if forme_col not in df.columns:
        raise KeyError(f"Colonne '{forme_col}' absente dans {df_name}")

# Normalisation
df_immat_2024[forme_col] = df_immat_2024[forme_col].astype(str).str.strip()
df_immat_2025[forme_col] = df_immat_2025[forme_col].astype(str).str.strip()

# Comptages par forme juridique
formes_2024 = df_immat_2024[forme_col].value_counts()
formes_2025 = df_immat_2025[forme_col].value_counts()

# Liste des formes présentes dans le tableau du modèle (colonne 1)
formes_selectionnees = []
for shape in slide9.shapes:
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for row_idx in range(len(table.rows)):
            forme_nom = table.cell(row_idx, 0).text.strip()
            if forme_nom:
                formes_selectionnees.append(forme_nom)

# Dictionnaire des valeurs
formes_resume = {}
for forme in formes_selectionnees:
    v24 = int(formes_2024.get(forme, 0))
    v25 = int(formes_2025.get(forme, 0))
    evol_pct = ((v25 - v24) / v24 * 100) if v24 > 0 else 0.0
    formes_resume[forme] = {"2024": v24, "2025": v25, "evol_pct": evol_pct}

# --- Remplissage du tableau ---
for shape in slide9.shapes:
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for row_idx in range(len(table.rows)):
            forme_nom = table.cell(row_idx, 0).text.strip()
            if forme_nom in formes_resume:
                val2025 = formes_resume[forme_nom]["2025"]
                evol_pct = formes_resume[forme_nom]["evol_pct"]
                # Colonne 2 = nombre immatriculations 2025
                table.cell(row_idx, 1).text = f"{val2025:,}".replace(",", " ")
                # Colonne 3 = évolution %
                table.cell(row_idx, 2).text = f"{evol_pct:+.1f}%"


# --- Diagramme secteur (donut) ---
# Supprimer anciens graphiques
for shape in list(slide9.shapes):
    if hasattr(shape, "has_chart") and shape.has_chart:
        slide9.shapes._spTree.remove(shape._element)

# Préparer les données pour le graphique
total_2025 = sum(val["2025"] for val in formes_resume.values())
parts_2025 = [
    (round(val["2025"] / total_2025 * 100, 1) if total_2025 > 0 else 0.0)
    for val in formes_resume.values()
]

chart_data = CategoryChartData()
chart_data.categories = list(formes_resume.keys())
chart_data.add_series("", parts_2025)

# Position et taille du graphique
x, y, cx, cy = Inches(0.4), Inches(1.3), Inches(4.8), Inches(7.3)
chart = slide9.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.bold = False
data_labels.show_category_name = True
data_labels.show_value = False
data_labels.show_percentage = True

# --- Épaisseur de l’anneau ---
try:
    plot.doughnut_hole_size = 70  # pris en charge par certaines versions de python-pptx
except AttributeError:
    pass

# Fallback XML (nécessite lxml)
try:
    from lxml import etree
    doughnut_chart = chart.plots[0]._element
    hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
    if hole_size is None:
        hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
        doughnut_chart.append(hole_size)
    hole_size.set('val', '70')  # 10..90
except Exception:
    # si lxml non dispo ou structure différente, on ignore silencieusement
    pass


#   PARTIE 2 : PROCÉDURES COL.   

# Compte uniquement les lignes non vides de la colonne SIREN
if "SIREN" not in df_pcl_2024.columns or "SIREN" not in df_pcl_2025.columns:
    raise KeyError("Colonne 'SIREN' absente de df_pcl_2024 ou df_pcl_2025.")

val_2024 = df_pcl_2024["SIREN"].notna().sum()
val_2025 = df_pcl_2025["SIREN"].notna().sum()

# Calcul de la variation en pourcentage
variation = ((val_2025 - val_2024) / val_2024 * 100) if val_2024 > 0 else 0

# Formatage
formatted_val = f"{val_2025:,}".replace(",", " ")
formatted_var = f"{variation:+.1f}%".replace("++", "+")

# === SLIDE 11 : remplir le tableau EXISTANT sans modifier la mise en forme ===
slide11 = prs.slides[10]

# Récupérer premier tableau trouvé (gauche = volumes, droit = répartition)
tables = [shape.table for shape in slide11.shapes if shape.has_table]
if len(tables) < 2:
    raise ValueError("La slide 11 doit contenir au moins 2 tableaux (gauche = volumes, droite = répartition).")
table_left, table_right = tables[0], tables[1]

# Préparation des données (à partir de df_stock_2025)
region_col = "ADRESSEFRANCEREGIONLIBELLE"
for col_required in [region_col, "INSCRIPTIONS_PRINCIPALES", "INSCRIPTIONS_SECONDAIRES"]:
    if col_required not in df_stock_2025.columns:
        raise KeyError(f"Colonne requise absente dans df_stock_2025 : {col_required}")

# Calcul des volumes et parts
counts_2025 = df_stock_2025.groupby(region_col).apply(
    lambda g: g["INSCRIPTIONS_PRINCIPALES"].fillna(0).sum() + g["INSCRIPTIONS_SECONDAIRES"].fillna(0).sum()
)
# Filtrer les index vides (sécurité si NaN/vides)
counts_2025 = counts_2025[counts_2025.index.astype(str).str.strip() != ""]
total_2025 = float(counts_2025.sum())
parts_region = (counts_2025 / total_2025 * 100).round(1) if total_2025 > 0 else counts_2025.copy() * 0

# Mapping "NOM TABLE PPT" -> (volume_str, pct_str)
lookup = {}
for r in counts_2025.index:
    name_display = str(r).split(" (")[0]  # correspondance avec noms de tableau
    vol_str = fmt_nb(counts_2025.loc[r])
    pct_str = f"{parts_region.loc[r]}%"
    lookup[name_display.strip().upper()] = (vol_str, pct_str)

# Utilitaire : remplacer le texte en conservant le format du 1er run + police 12
def set_cell_text_preserve_style(cell, new_text):
    p = cell.text_frame.paragraphs[0]
    runs = p.runs
    if runs:
        runs[0].text = new_text
        runs[0].font.size = Pt(12)  # forcer police 12
        for rr in runs[1:]:
            rr.text = ""
    else:
        p.text = new_text
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

# Remplir le tableau GAUCHE (volumes) — la colonne 0 (régions) reste comme dans PPT
n_rows_left = len(table_left.rows)
for row_idx in range(1, n_rows_left):  # ne pas toucher l'en-tête
    region_cell = table_left.cell(row_idx, 0)
    region_name = region_cell.text.strip().upper()
    if region_name == "":
        break
    if region_name in lookup:
        vol_str, _ = lookup[region_name]
        set_cell_text_preserve_style(table_left.cell(row_idx, 1), vol_str)
    else:
        set_cell_text_preserve_style(table_left.cell(row_idx, 1), "")

# Remplir le tableau DROIT (répartition %)
n_rows_right = len(table_right.rows)
for row_idx in range(1, n_rows_right):
    region_cell = table_right.cell(row_idx, 0)
    region_name = region_cell.text.strip().upper()
    if region_name == "":
        break
    if region_name in lookup:
        _, pct_str = lookup[region_name]
        set_cell_text_preserve_style(table_right.cell(row_idx, 1), pct_str)
    else:
        set_cell_text_preserve_style(table_right.cell(row_idx, 1), "")

# Mettre le total global "Y" si présent
formatted_total = fmt_nb(total_2025)
for shape in slide11.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() == "Y":
                    run.text = formatted_total
                    run.font.size = Pt(12)
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 13 : remplacer X / X% par les valeurs PCL
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text in ["X%", "+X%", "-X%"]:
                    run.text = run.text.replace("X%", formatted_var)
                    paragraph.alignment = PP_ALIGN.CENTER
                elif text == "X":
                    run.text = formatted_val
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 14 : tableau des évolutions T3 2024 vs T3 2025
slide14 = prs.slides[13]
table_shapes = [shape for shape in slide14.shapes if shape.has_table]

if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 14 — vérifie que le modèle contient bien un tableau.")
else:
    table = table_shapes[0].table

# Préparation des données
regions_exclues = ["Territoires d'outre-mer (TOM)"]
if region_col not in df_pcl_2024.columns or region_col not in df_pcl_2025.columns:
    raise KeyError(f"Colonne {region_col} absente des données PCL.")

pcl_region_2024 = (
    df_pcl_2024[df_pcl_2024["SIREN"].astype(str).str.strip() != ""]
    .groupby(region_col)["SIREN"].count()
)
pcl_region_2025 = (
    df_pcl_2025[df_pcl_2025["SIREN"].astype(str).str.strip() != ""]
    .groupby(region_col)["SIREN"].count()
)

# Nettoyage
pcl_region_2024 = pcl_region_2024[~pcl_region_2024.index.isin(regions_exclues)]
pcl_region_2025 = pcl_region_2025[~pcl_region_2025.index.isin(regions_exclues)]

# Fusion + évolution
regions = sorted(set(pcl_region_2024.index).union(set(pcl_region_2025.index)))
evol_regions = {}
for r in regions:
    v24 = int(pcl_region_2024.get(r, 0))
    v25 = int(pcl_region_2025.get(r, 0))
    evol = ( (v25 - v24) / v24 * 100 ) if v24 > 0 else 0.0
    evol_regions[str(r).split(" (")[0]] = round(evol, 1)

region_evol_tuples = sorted(evol_regions.items(), key=lambda x: -x[1])

# Remplissage du tableau
table.cell(0, 0).text = "Régions"
table.cell(0, 1).text = "Évolution (en %) T3 2024 vs T3 2025"

# Mise en forme en-têtes
for j in range(2):
    cell = table.cell(0, j)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 32, 96)
            run.font.bold = True

# Données
for i, (region, evol) in enumerate(region_evol_tuples[:len(table.rows) - 1]):
    table.cell(i + 1, 0).text = region
    table.cell(i + 1, 1).text = f"{evol:+.1f}%"

# Vider le surplus de lignes si besoin
for i in range(len(region_evol_tuples) + 1, len(table.rows)):
    table.cell(i, 0).text = ""
    table.cell(i, 1).text = ""


# Mise en forme globale
for row_idx, row in enumerate(table.rows):
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            if row_idx == 0:
                # Ligne des titres : style spécifique
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 32, 96)  # Bleu titres
                    run.font.name = "Aptos"
            else:
                # Lignes de données
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.name = "Aptos"

# === SLIDE 15 ===
slide15 = prs.slides[14]

# --- Colonnes & sélection ---
secteur_col = "SECTEURLIBELLE"
secteurs_selectionnes = [
    "ACTIVITÉS IMMOBILIÈRES",
    "COMMERCE , RÉPARATION D'AUTOMOBILES ET DE MOTOCYCLES",
    "ACTIVITÉS SPÉCIALISÉES, SCIENTIFIQUES ET TECHNIQUES",
    "TRANSPORTS ET ENTREPOSAGE",
    "ACTIVITÉS DE SERVICES ADMINISTRATIFS ET DE SOUTIEN",
    "HÉBERGEMENT ET RESTAURATION",
    "CONSTRUCTION",
]

# --- Vérifications colonnes ---
for df_name, df in [("df_pcl_2024", df_pcl_2024), ("df_pcl_2025", df_pcl_2025)]:
    if secteur_col not in df.columns:
        raise KeyError(f"Colonne '{secteur_col}' absente dans {df_name}")

# --- Normalisation (UPPER + strip) ---
df_pcl_2024[secteur_col] = df_pcl_2024[secteur_col].astype(str).str.strip().str.upper()
df_pcl_2025[secteur_col] = df_pcl_2025[secteur_col].astype(str).str.strip().str.upper()

# --- Comptages par secteur ---
pcl_2024 = df_pcl_2024[secteur_col].value_counts()
pcl_2025 = df_pcl_2025[secteur_col].value_counts()

# --- Constitution des catégories + "AUTRES" ---
secteurs = {}
for sec in secteurs_selectionnes:
    secteurs[sec] = {"2024": int(pcl_2024.get(sec, 0)), "2025": int(pcl_2025.get(sec, 0))}

autres_2024 = int(pcl_2024[~pcl_2024.index.isin(secteurs_selectionnes)].sum())
autres_2025 = int(pcl_2025[~pcl_2025.index.isin(secteurs_selectionnes)].sum())
secteurs["AUTRES"] = {"2024": autres_2024, "2025": autres_2025}

# --- Évolutions en % ---
for sec, vals in secteurs.items():
    v24 = vals["2024"]
    v25 = vals["2025"]
    evol = ((v25 - v24) / v24 * 100) if v24 > 0 else 0.0
    secteurs[sec]["evol_num"] = round(evol, 1)

# --- Suppression d'éventuels anciens graphiques ---
for shape in list(slide15.shapes):
    if hasattr(shape, "has_chart") and shape.has_chart:
        slide15.shapes._spTree.remove(shape._element)

# --- Création graphique en barres ---
chart_data = CategoryChartData()
chart_data.categories = list(secteurs.keys())
chart_data.add_series(
    "", [secteurs[s]["evol_num"] for s in secteurs]
)

x, y, cx, cy = Inches(2), Inches(1), Inches(8), Inches(5)
chart = slide15.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

# --- Mise en forme ---
chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True

# Libellés hors barres, formatés
for i, point in enumerate(plot.series[0].points):
    val = secteurs[list(secteurs.keys())[i]]["evol_num"]
    lbl = point.data_label
    lbl.has_text_frame = True
    lbl.text_frame.clear()
    p = lbl.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{val:.1f}%"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 32, 96)  # bleu #002060
    lbl.position = XL_LABEL_POSITION.OUTSIDE_END

# Axe des valeurs en pourcentage
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.number_format = '0.0"%"'
value_axis.tick_labels.font.size = Pt(10)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True

# Couleurs par point (rouge si >=0, vert si <0)
for i, point in enumerate(plot.series[0].points):
    val = secteurs[list(secteurs.keys())[i]]["evol_num"]
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 0, 0) if val >= 0 else RGBColor(0, 176, 80)


# === SLIDE 16 ===
slide16 = prs.slides[15]

# --- Calcul répartition par type de jugement ---
type_col = "TYPEJUGEMENTOUVERTURELIBELLE"

# Vérif colonnes
for df_name, df in [("df_pcl_2024", df_pcl_2024), ("df_pcl_2025", df_pcl_2025)]:
    if type_col not in df.columns:
        raise KeyError(f"Colonne '{type_col}' absente dans {df_name}")

# Normalisation
df_pcl_2024[type_col] = df_pcl_2024[type_col].astype(str).str.strip()
df_pcl_2025[type_col] = df_pcl_2025[type_col].astype(str).str.strip()

types_selectionnes = [
    "LIQUIDATION JUDICIAIRE",
    "LIQUIDATION JUDICIAIRE SIMPLIFIEE",
    "REDRESSEMENT JUDICIAIRE",
    "REDRESSEMENT JUDICIAIRE GENERAL",
    "REDRESSEMENT JUDICIAIRE SIMPLIFIE",
    "PROCEDURE DE SAUVEGARDE",
    "PROCÉDURE DE SAUVEGARDE ACCÉLÉRÉE",
]

# Comptages bruts
types_2024 = df_pcl_2024[type_col].str.upper().value_counts()
types_2025 = df_pcl_2025[type_col].str.upper().value_counts()

# Dictionnaire global par sous-type
types = {}
for f in types_selectionnes:
    f_up = f.upper()
    types[f_up] = {
        "2024": int(types_2024.get(f_up, 0)),
        "2025": int(types_2025.get(f_up, 0)),
    }

# Catégorie AUTRES
autres_2024 = int(types_2024[~types_2024.index.isin([t.upper() for t in types_selectionnes])].sum())
autres_2025 = int(types_2025[~types_2025.index.isin([t.upper() for t in types_selectionnes])].sum())
types["AUTRES TYPES JUGEMENT"] = {"2024": autres_2024, "2025": autres_2025}

# Agrégation par familles principales
regroupements = {
    "LIQUIDATION JUDICIAIRE": [
        "LIQUIDATION JUDICIAIRE", "LIQUIDATION JUDICIAIRE SIMPLIFIEE"
    ],
    "REDRESSEMENT JUDICIAIRE": [
        "REDRESSEMENT JUDICIAIRE", "REDRESSEMENT JUDICIAIRE GENERAL", "REDRESSEMENT JUDICIAIRE SIMPLIFIE"
    ],
    "PROCEDURE DE SAUVEGARDE": [
        "PROCEDURE DE SAUVEGARDE", "PROCÉDURE DE SAUVEGARDE ACCÉLÉRÉE"
    ],
}

# Résumé (pour tableau et donut)
types_resume = {}
for cat, sous_types in regroupements.items():
    v24 = sum(types.get(st.upper(), {"2024": 0})["2024"] for st in sous_types)
    v25 = sum(types.get(st.upper(), {"2025": 0})["2025"] for st in sous_types)
    diff = v25 - v24
    evol = (diff / v24 * 100) if v24 > 0 else 0.0
    types_resume[cat] = {"2024": v24, "2025": v25, "diff": diff, "evol_pct": evol}

# --- Remplissage du tableau (celui de droite dans la diapo) ---
for shape in slide16.shapes:
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for row_idx in range(len(table.rows)):
            type_nom = table.cell(row_idx, 0).text.strip().upper()
            if type_nom in types_resume:
                val2025 = types_resume[type_nom]["2025"]
                evol_pct = types_resume[type_nom]["evol_pct"]
                diff = types_resume[type_nom]["diff"]
                # Colonne 2 = nombre ouvertures 2025
                table.cell(row_idx, 1).text = fmt_nb(val2025)
                # Colonne 3 = évolution % (avec diff en volume)
                table.cell(row_idx, 2).text = f"{evol_pct:+.1f}% ({diff:+d})"

# --- Mise à jour diagramme (donut) ---
# Supprimer d'anciens graphiques
for shape in list(slide16.shapes):
    if hasattr(shape, "has_chart") and shape.has_chart:
        slide16.shapes._spTree.remove(shape._element)

total_2025 = sum(val["2025"] for val in types_resume.values())
parts_2025 = [
    (round(val["2025"] / total_2025 * 100, 1) if total_2025 > 0 else 0.0)
    for val in types_resume.values()
]

chart_data = CategoryChartData()
chart_data.categories = list(types_resume.keys())
chart_data.add_series("", parts_2025)

x, y, cx, cy = Inches(0.4), Inches(1.0), Inches(4.8), Inches(7.3)
chart = slide16.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.bold = False
data_labels.show_category_name = True
data_labels.show_value = True
data_labels.show_percentage = False

# --- Épaisseur de l’anneau ---
try:
    plot.doughnut_hole_size = 70  # pris en charge par certaines versions de python-pptx
except AttributeError:
    pass

# Fallback XML (nécessite lxml)
try:
    from lxml import etree
    doughnut_chart = chart.plots[0]._element
    hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
    if hole_size is None:
        hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
        doughnut_chart.append(hole_size)
    hole_size.set('val', '70')  # 10..90
except Exception:
    # si lxml non dispo ou structure différente, on ignore silencieusement
    pass



# === SLIDE 17 ===
slide17 = prs.slides[16]  # index 0-based : vérifie que c’est bien la bonne diapo

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- Ne supprimer que les anciens graphiques (laisser tableaux et zones de texte) ---
for shp in list(slide17.shapes):
    if hasattr(shp, "has_chart") and shp.has_chart:
        slide17.shapes._spTree.remove(shp._element)

# --- Fonction de regroupement par familles principales ---
def regrouper_statuts(df):
    type_col = "TYPEJUGEMENTOUVERTURELIBELLE"
    if type_col not in df.columns:
        raise KeyError(f"Colonne '{type_col}' absente des données.")
    # Normalisation avec accents
    series = (
        df[type_col].astype(str)
        .str.strip()
        .str.upper()
        .str.replace("ACCELEREE", "ACCÉLÉRÉE", regex=False)
        .str.replace("SAUVEGARDE", "SAUVEGARDE", regex=False)
    )

    regroupements = {
        "LIQUIDATION JUDICIAIRE": [
            "LIQUIDATION JUDICIAIRE", "LIQUIDATION JUDICIAIRE SIMPLIFIEE"
        ],
        "REDRESSEMENT JUDICIAIRE": [
            "REDRESSEMENT JUDICIAIRE",
            "REDRESSEMENT JUDICIAIRE GENERAL",
            "REDRESSEMENT JUDICIAIRE SIMPLIFIE"
        ],
        "PROCÉDURE DE SAUVEGARDE": [
            "PROCEDURE DE SAUVEGARDE",
            "PROCÉDURE DE SAUVEGARDE ACCÉLÉRÉE"  # avec accents
        ],
    }

    counts = {}
    for cat, sous_types in regroupements.items():
        counts[cat] = series.isin(sous_types).sum()
    return counts

# --- Filtrage forme SARL ---
forme_col = "FORMEJURIDIQUELIBELLE"
if forme_col not in df_pcl_2025.columns:
    raise KeyError(f"Colonne '{forme_col}' absente des données PCL.")

df_pcl_2025_norm = df_pcl_2025.copy()
df_pcl_2025_norm[forme_col] = df_pcl_2025_norm[forme_col].astype(str).str.strip()

df_sarl = df_pcl_2025_norm[
    df_pcl_2025_norm[forme_col].str.contains("responsabilité limitée", case=False, na=False)
]
counts_sarl = regrouper_statuts(df_sarl)

# --- Filtrage forme SAS ---
df_sas = df_pcl_2025_norm[
    df_pcl_2025_norm[forme_col].str.contains("actions simplifi", case=False, na=False)
]
counts_sas = regrouper_statuts(df_sas)

# --- Fonction utilitaire pour créer un donut propre ---
def add_donut(slide, title, counts, x, y, cx, cy, colors=None):
    # Si toutes les valeurs sont nulles, on met un seul segment "Aucun"
    values = list(counts.values())
    cats = list(counts.keys())
    if sum(values) == 0:
        cats = ["Aucun dossier"]
        values = [1]

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series(title, values)

    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart

    # Data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_category_name = True
    dl.show_value = True
    dl.show_percentage = False
    dl.font.size = Pt(9)

    chart.has_legend = False

    # Épaisseur anneau
    try:
        plot.doughnut_hole_size = 65
    except AttributeError:
        # fallback silencieux si version ne supporte pas
        pass

    # Palette de couleurs (optionnelle)
    if colors:
        # Applique les couleurs dans l’ordre des points
        for s in chart.series:
            for idx, pt in enumerate(s.points):
                rgb = colors[idx % len(colors)]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = RGBColor(*rgb)

    return chart

# --- Créer les deux donuts avec positions ajustées ---
# SARL
colors_sarl = [
    (0, 122, 204),   # bleu
    (0, 158, 73),    # vert
    (255, 140, 0),   # orange
]
chart_sarl = add_donut(
    slide17,
    "SARL",
    counts_sarl,
    x=Inches(0.7), y=Inches(1.4), cx=Inches(4.0), cy=Inches(4.0),
    colors=colors_sarl
)

# SAS
colors_sas = [
    (0, 122, 204),    # bleu
    (186, 85, 211),   # violet
    (220, 20, 60),    # rouge
]
chart_sas = add_donut(
    slide17,
    "SAS",
    counts_sas,
    x=Inches(7.0), y=Inches(1.4), cx=Inches(4.0), cy=Inches(4.0),
    colors=colors_sas
)

# --- Titres (placés au-dessus de chaque donut) ---
title_sarl = slide17.shapes.add_textbox(Inches(1.5), Inches(0.6), Inches(4.3), Inches(0.6))
p1 = title_sarl.text_frame.paragraphs[0]
p1.text = "Sociétés à responsabilité limitée (SARL)"
p1.alignment = PP_ALIGN.CENTER
for run in p1.runs:
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0, 32, 96)

title_sas = slide17.shapes.add_textbox(Inches(7.8), Inches(0.6), Inches(4.3), Inches(0.6))
p2 = title_sas.text_frame.paragraphs[0]
p2.text = "Sociétés par actions simplifiées (SAS)"
p2.alignment = PP_ALIGN.CENTER
for run in p2.runs:
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0, 32, 96)


# PARTIE 3 : RADIATIONS

val_2024 = df_radiation_2024["SIREN"].notna().sum()
val_2025 = df_radiation_2025["SIREN"].notna().sum()

variation = ((val_2025 - val_2024) / val_2024) * 100 if val_2024 > 0 else 0

formatted_val = f"{val_2025:,}".replace(",", " ")
formatted_var = f"{variation:+.1f}%".replace("++", "+")

# === SLIDE 19 ===
slide19 = prs.slides[18]
for shape in slide19.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text in ["X%", "+X%", "-X%"]:
                    run.text = run.text.replace("X%", formatted_var)
                    paragraph.alignment = PP_ALIGN.CENTER
                elif text == "X":
                    run.text = formatted_val
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 20 ===
slide20 = prs.slides[19]
table_shapes = [shape for shape in slide20.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 20")
table = table_shapes[0].table

region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

radiation_region_2024 = (
    df_radiation_2024[df_radiation_2024["SIREN"].astype(str).str.strip() != ""]
    .groupby(region_col)["SIREN"].count()
)
radiation_region_2025 = (
    df_radiation_2025[df_radiation_2025["SIREN"].astype(str).str.strip() != ""]
    .groupby(region_col)["SIREN"].count()
)

radiation_region_2024 = radiation_region_2024[~radiation_region_2024.index.isin(regions_exclues)]
radiation_region_2025 = radiation_region_2025[~radiation_region_2025.index.isin(regions_exclues)]

regions = sorted(set(radiation_region_2024.index).union(set(radiation_region_2025.index)))
evol_regions = {}
for region in regions:
    v2024 = radiation_region_2024.get(region, 0)
    v2025 = radiation_region_2025.get(region, 0)
    evol = ((v2025 - v2024) / v2024 * 100) if v2024 > 0 else 0
    evol_regions[region.split(" (")[0]] = round(evol, 1)

region_evol_tuples = sorted(evol_regions.items(), key=lambda x: -x[1])

# Remplissage du tableau
table.cell(0, 0).text = "Régions"
table.cell(0, 1).text = "Évolution 2025 vs 2024 (en %)"
for j in range(2):
    cell = table.cell(0, j)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 32, 96)
            run.font.bold = True
            run.font.size = Pt(11)

for i, (region, evol) in enumerate(region_evol_tuples[:len(table.rows) - 1]):
    table.cell(i + 1, 0).text = region
    table.cell(i + 1, 1).text = f"{evol:+.1f}%"

for i in range(len(region_evol_tuples) + 1, len(table.rows)):
    table.cell(i, 0).text = ""
    table.cell(i, 1).text = ""

# Mise en forme des lignes sauf la première
for row_idx, row in enumerate(table.rows):
    if row_idx == 0:  # Ignorer la ligne des titres
        continue
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"
                run.font.color.rgb = RGBColor(0, 0, 0) 


# === SLIDE 21 ===
slide21 = prs.slides[20]

# --- Colonnes & sélection ---
secteur_col = "SECTEURLIBELLE"
secteurs_selectionnes = [
    "ACTIVITÉS IMMOBILIÈRES",
    "COMMERCE , RÉPARATION D'AUTOMOBILES ET DE MOTOCYCLES",
    "ACTIVITÉS SPÉCIALISÉES, SCIENTIFIQUES ET TECHNIQUES",
    "TRANSPORTS ET ENTREPOSAGE",
    "ACTIVITÉS DE SERVICES ADMINISTRATIFS ET DE SOUTIEN",
    "HÉBERGEMENT ET RESTAURATION",
    "CONSTRUCTION"
]

# --- Vérification colonnes ---
for df_name, df in [("df_radiation_2024", df_radiation_2024), ("df_radiation_2025", df_radiation_2025)]:
    if secteur_col not in df.columns:
        raise KeyError(f"Colonne '{secteur_col}' absente dans {df_name}")

# --- Normalisation ---
df_radiation_2024[secteur_col] = df_radiation_2024[secteur_col].astype(str).str.strip().str.upper()
df_radiation_2025[secteur_col] = df_radiation_2025[secteur_col].astype(str).str.strip().str.upper()

# --- Comptages par secteur ---
tx_2024 = df_radiation_2024[secteur_col].value_counts()
tx_2025 = df_radiation_2025[secteur_col].value_counts()

# --- Constitution des catégories + AUTRES ---
secteurs = {}
for sec in secteurs_selectionnes:
    secteurs[sec] = {"2024": int(tx_2024.get(sec, 0)), "2025": int(tx_2025.get(sec, 0))}

autres_2024 = int(tx_2024[~tx_2024.index.isin(secteurs_selectionnes)].sum())
autres_2025 = int(tx_2025[~tx_2025.index.isin(secteurs_selectionnes)].sum())
secteurs["AUTRES"] = {"2024": autres_2024, "2025": autres_2025}

# --- Calcul évolutions en % ---
for sec, vals in secteurs.items():
    v24 = vals["2024"]
    v25 = vals["2025"]
    evol = ((v25 - v24) / v24 * 100) if v24 > 0 else 0.0
    secteurs[sec]["evol_num"] = round(evol, 1)

# --- Suppression anciens graphiques ---
for shape in list(slide21.shapes):
    if hasattr(shape, "has_chart") and shape.has_chart:
        slide21.shapes._spTree.remove(shape._element)

# --- Création graphique en barres ---
chart_data = CategoryChartData()
chart_data.categories = list(secteurs.keys())
chart_data.add_series("", [secteurs[s]["evol_num"] for s in secteurs])

x, y, cx, cy = Inches(1.5), Inches(1), Inches(9), Inches(5)
chart = slide21.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

# --- Mise en forme ---
chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True

# Libellés hors barres
for i, point in enumerate(plot.series[0].points):
    val = secteurs[list(secteurs.keys())[i]]["evol_num"]
    lbl = point.data_label
    lbl.has_text_frame = True
    lbl.text_frame.clear()
    p = lbl.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{val:.1f}%"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 32, 96)  # Bleu
    lbl.position = XL_LABEL_POSITION.OUTSIDE_END

# Axe des valeurs en %
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.number_format = '0.0"%"'
value_axis.tick_labels.font.size = Pt(10)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True

# Couleurs par point (rouge si hausse, vert si baisse)
for i, point in enumerate(plot.series[0].points):
    val = secteurs[list(secteurs.keys())[i]]["evol_num"]
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 0, 0) if val >= 0 else RGBColor(0, 176, 80)

# === SLIDE 22 ===
slide22 = prs.slides[21]
table_shapes = [shape for shape in slide22.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 22")
table = table_shapes[0].table

forme_col = "FORMEJURIDIQUELIBELLE"
df_2024_nv = df_radiation_2024[df_radiation_2024["SIREN"].astype(str).str.strip() != ""].copy()
df_2025_nv = df_radiation_2025[df_radiation_2025["SIREN"].astype(str).str.strip() != ""].copy()

formes_selectionnees = [
    "Société à responsabilité limitée",
    "Société civile immobilière",
    "Société par actions simplifiée"
]

formes_2024 = df_2024_nv[forme_col].astype(str).str.strip().value_counts(dropna=False)
formes_2025 = df_2025_nv[forme_col].astype(str).str.strip().value_counts(dropna=False)

formes = {}
for f in formes_selectionnees:
    formes[f] = {"2024": formes_2024.get(f, 0), "2025": formes_2025.get(f, 0)}

commercant_2024 = df_2024_nv[forme_col].isna().sum() + (df_2024_nv[forme_col].astype(str).str.strip() == "").sum()
commercant_2025 = df_2025_nv[forme_col].isna().sum() + (df_2025_nv[forme_col].astype(str).str.strip() == "").sum()
formes["Commerçant (Entreprise individuelle)"] = {"2024": commercant_2024, "2025": commercant_2025}

mask_autres_2024 = ~formes_2024.index.isin(formes_selectionnees + ["", "nan"])
mask_autres_2025 = ~formes_2025.index.isin(formes_selectionnees + ["", "nan"])
formes["Autres formes juridiques"] = {"2024": formes_2024[mask_autres_2024].sum(),
                                     "2025": formes_2025[mask_autres_2025].sum()}

total_2024 = sum(f["2024"] for f in formes.values())
total_2025 = sum(f["2025"] for f in formes.values())
formes["Total"] = {"2024": total_2024, "2025": total_2025}

for f in formes:
    v2024 = formes[f]["2024"]
    v2025 = formes[f]["2025"]
    pct = ((v2025 - v2024) / v2024 * 100) if v2024 > 0 else 0
    formes[f]["pct"] = pct

headers = ["Formes juridiques", "Valeurs T3 2025", "Évolution (%)"]
for j, h in enumerate(headers):
    table.cell(0, j).text = h
    for paragraph in table.cell(0, j).text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 32, 96)
            run.font.size = Pt(10)
            run.font.name = "Aptos"

for i, (forme, data) in enumerate(formes.items()):
    if i + 1 >= len(table.rows):
        break
    table.cell(i + 1, 0).text = forme
    table.cell(i + 1, 1).text = f"{data['2025']:,}".replace(",", " ")
    table.cell(i + 1, 2).text = f"{data['pct']:+.1f}%"

for i in range(len(formes) + 1, len(table.rows)):
    for j in range(len(table.columns)):
        table.cell(i, j).text = ""

# Mise en forme des lignes sauf la première
for row_idx, row in enumerate(table.rows):
    if row_idx == 0:  # Ignorer la ligne des titres
        continue
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"
                run.font.color.rgb = RGBColor(0, 0, 0) 

# === SLIDE 23 ===
slide23 = prs.slides[22]
table_shapes = [shape for shape in slide23.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 23")
table = table_shapes[0].table

motif_col = "RADIATIONMOTIFRADIATIONLIBELLE"

# Nettoyage des données
df_2024_nv = df_radiation_2024[df_radiation_2024["SIREN"].astype(str).str.strip() != ""].copy()
df_2025_nv = df_radiation_2025[df_radiation_2025["SIREN"].astype(str).str.strip() != ""].copy()

df_2024_nv[motif_col] = df_2024_nv[motif_col].astype(str).str.strip()
df_2025_nv[motif_col] = df_2025_nv[motif_col].astype(str).str.strip()

# Groupes de motifs
motifs_d_office = [
    "Expiration du délai de 3 mois pour mises à jour informations relatives aux BE",
    "Situation non régularisée (personne ne se trouvant plus à l'adresse indiquée)",
    "Situation non régularisée (pièce ou acte irrégulier)"
]
motifs_volontaires = ["Dissolution", "Renseignement insuffisant"]
motifs_proc_collective = ["Cessation d'activité", "Procédure collective"]

# Comptages
def count_motifs(df, motifs):
    return df[motif_col].isin(motifs).sum()

radiations = {
    "Radiations d'office": {
        "2024": count_motifs(df_2024_nv, motifs_d_office),
        "2025": count_motifs(df_2025_nv, motifs_d_office)
    },
    "Radiations volontaires": {
        "2024": count_motifs(df_2024_nv, motifs_volontaires),
        "2025": count_motifs(df_2025_nv, motifs_volontaires)
    },
    "Radiations à la suite d'une procédure collective": {
        "2024": count_motifs(df_2024_nv, motifs_proc_collective),
        "2025": count_motifs(df_2025_nv, motifs_proc_collective)
    }
}

# Autres motifs
total_2024 = len(df_2024_nv)
total_2025 = len(df_2025_nv)
autres_2024 = total_2024 - sum(r["2024"] for r in radiations.values())
autres_2025 = total_2025 - sum(r["2025"] for r in radiations.values())
radiations["Autres motifs"] = {"2024": autres_2024, "2025": autres_2025}

# Dernière ligne : % radiations d’office
pct_office_2024 = (radiations["Radiations d'office"]["2024"] / total_2024 * 100) if total_2024 > 0 else 0
pct_office_2025 = (radiations["Radiations d'office"]["2025"] / total_2025 * 100) if total_2025 > 0 else 0
evol_pct_office = ((pct_office_2025 - pct_office_2024) / pct_office_2024 * 100) if pct_office_2024 > 0 else 0

# Évolution pour chaque ligne
for key in radiations:
    v24 = radiations[key]["2024"]
    v25 = radiations[key]["2025"]
    evol = ((v25 - v24) / v24 * 100) if v24 > 0 else 0
    radiations[key]["evol"] = evol

# Remplissage du tableau
headers = ["Motifs de radiation", "T3 2025", "Évolution (%) T3 2025 vs T3 2024"]
for j, h in enumerate(headers):
    table.cell(0, j).text = h
    for paragraph in table.cell(0, j).text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 32, 96)
            run.font.size = Pt(10)
            run.font.name = "Aptos"

# Lignes principales
rows = [
    "Radiations d'office",
    "Radiations volontaires",
    "Radiations à la suite d'une procédure collective",
    "Autres motifs",
    "Pourcentage des radiations d'office par rapport au total des radiations"
]

for i, row_name in enumerate(rows):
    if i + 1 >= len(table.rows):
        break
    table.cell(i + 1, 0).text = row_name
    if row_name == "Pourcentage des radiations d'office par rapport au total des radiations":
        table.cell(i + 1, 1).text = f"{pct_office_2025:.1f}%"
        table.cell(i + 1, 2).text = f"{evol_pct_office:+.1f}%"
    else:
        table.cell(i + 1, 1).text = f"{radiations[row_name]['2025']:,}".replace(",", " ")
        table.cell(i + 1, 2).text = f"{radiations[row_name]['evol']:+.1f}%"

# Vider les lignes restantes
for i in range(len(rows) + 1, len(table.rows)):
    for j in range(len(table.columns)):
        table.cell(i, j).text = ""

# Mise en forme (sauf en-têtes)
for row_idx, row in enumerate(table.rows):
    if row_idx == 0:
        continue
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"
                run.font.color.rgb = RGBColor(0, 0, 0)

# === SLIDE 25 ===

slide25 = prs.slides[24]

# Supprimer ancien graphique s'il existe
for shape in list(slide25.shapes):
    if shape.has_chart:
        slide25.shapes._spTree.remove(shape._element)

# Colonnes arrivées / départs
region_to_col = "REGIONTO"
region_from_col = "REGIONFROM"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

# Préparer séries 2024
arrivees_2024 = df_ts_2024[region_to_col].fillna("").astype(str).str.strip()
departs_2024 = df_ts_2024[region_from_col].fillna("").astype(str).str.strip()
arrivees_2024 = arrivees_2024[arrivees_2024 != ""].value_counts()
departs_2024 = departs_2024[departs_2024 != ""].value_counts()

# Préparer séries 2025
arrivees_2025 = df_ts_2025[region_to_col].fillna("").astype(str).str.strip()
departs_2025 = df_ts_2025[region_from_col].fillna("").astype(str).str.strip()
arrivees_2025 = arrivees_2025[arrivees_2025 != ""].value_counts()
departs_2025 = departs_2025[departs_2025 != ""].value_counts()

# Exclure régions non voulues
arrivees_2024 = arrivees_2024[~arrivees_2024.index.isin(regions_exclues)]
departs_2024 = departs_2024[~departs_2024.index.isin(regions_exclues)]
arrivees_2025 = arrivees_2025[~arrivees_2025.index.isin(regions_exclues)]
departs_2025 = departs_2025[~departs_2025.index.isin(regions_exclues)]

# Liste des régions
regions_full = sorted(
    set(arrivees_2024.index).union(set(departs_2024.index))
    .union(set(arrivees_2025.index)).union(set(departs_2025.index))
)
regions_labels = [r.split(" (")[0] for r in regions_full]

# Calcul solde = arrivées - départs
solde_2024 = [arrivees_2024.get(r, 0) - departs_2024.get(r, 0) for r in regions_full]
solde_2025 = [arrivees_2025.get(r, 0) - departs_2025.get(r, 0) for r in regions_full]

# Création du graphique groupé (barres horizontales)
chart_data = CategoryChartData()
chart_data.categories = regions_labels
chart_data.add_series("T3 2024", solde_2024)
chart_data.add_series("T3 2025", solde_2025)

x, y, cx, cy = Inches(2.3), Inches(1.0), Inches(9), Inches(6)
chart = slide25.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

chart.has_legend = True
chart.legend.include_in_layout = False
chart.legend.position = XL_LEGEND_POSITION.BOTTOM

plot = chart.plots[0]
plot.has_data_labels = True

serie_colors = [RGBColor(0, 32, 96), RGBColor(91, 155, 213)]
solde_lists = [solde_2024, solde_2025]

for s_idx, series in enumerate(plot.series):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = serie_colors[s_idx]
    series.invert_if_negative = False

    for p_idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = serie_colors[s_idx]

        dLbl = point.data_label
        dLbl.has_text_frame = True
        dLbl.text_frame.clear()
        p = dLbl.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{solde_lists[s_idx][p_idx]:+d}"
        run.font.size = Pt(9)
        run.font.bold = True
        dLbl.position = XL_LABEL_POSITION.OUTSIDE_END

# Axe des valeurs
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.font.size = Pt(9)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True
category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

# === SLIDE 26 ===
slide26 = prs.slides[25]
table_shapes = [shape for shape in slide26.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 26")
table = table_shapes[0].table

region_to_col = "REGIONTO"
region_from_col = "REGIONFROM"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

arrivees = df_ts_2025[region_to_col].value_counts()
departs = df_ts_2025[region_from_col].value_counts()

arrivees = arrivees[~arrivees.index.isin(regions_exclues)]
departs = departs[~departs.index.isin(regions_exclues)]

regions = sorted(set(arrivees.index).union(set(departs.index)))

solde_regions = {}
for region in regions:
    val_in = arrivees.get(region, 0)
    val_out = departs.get(region, 0)
    solde_regions[region.split(" (")[0]] = val_in - val_out

region_solde_tuples = sorted(solde_regions.items(), key=lambda x: -x[1])

headers = ["Régions", "Solde net T3 2025"]
for j, h in enumerate(headers):
    table.cell(0, j).text = h
    for paragraph in table.cell(0, j).text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 32, 96)
            run.font.size = Pt(10)
            run.font.name = "Aptos"

for i, (region, solde) in enumerate(region_solde_tuples):
    if i + 1 >= len(table.rows):
        break
    table.cell(i + 1, 0).text = region
    table.cell(i + 1, 1).text = f"{solde:+d}"

for i in range(len(region_solde_tuples) + 1, len(table.rows)):
    for j in range(2):
        table.cell(i, j).text = ""

# Mise en forme des lignes sauf la première
for row_idx, row in enumerate(table.rows):
    if row_idx == 0:  # Ignorer la ligne des titres
        continue
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"
                run.font.color.rgb = RGBColor(0, 0, 0) 

### === SLIDE 28 ===
slide28 = prs.slides[27]

# --- Comptage global (nb lignes avec SIREN non vide) ---
val_2024 = df_ip_2024["SIREN"].dropna().astype(str).str.strip().count()
val_2025 = df_ip_2025["SIREN"].dropna().astype(str).str.strip().count()
variation = ((val_2025 - val_2024) / val_2024 * 100) if val_2024 > 0 else 0

# Formatage global
formatted_val = f"{val_2025:,}".replace(",", " ")
formatted_var = f"{variation:+.1f}%"

# --- Récupération du tableau existant ---
table_shapes = [shape for shape in slide28.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 28 — vérifie que le modèle contient bien un tableau.")
else:
    table = table_shapes[0].table

# --- Comptage par région (SIREN non vides) ---
region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

ip_region_2024 = (
    df_ip_2024[df_ip_2024["SIREN"].astype(str).str.strip() != ""]
    .groupby(region_col)["SIREN"].count()
    .drop(labels=regions_exclues, errors="ignore")
)
ip_region_2025 = (
    df_ip_2025[df_ip_2025["SIREN"].astype(str).str.strip() != ""]
    .groupby(region_col)["SIREN"].count()
    .drop(labels=regions_exclues, errors="ignore")
)

# --- Calcul évolutions % ---
regions = sorted(set(ip_region_2024.index).union(ip_region_2025.index))
region_stats = []
for region in regions:
    val_2024_r = ip_region_2024.get(region, 0)
    val_2025_r = ip_region_2025.get(region, 0)
    evol = ((val_2025_r - val_2024_r) / val_2024_r * 100) if val_2024_r > 0 else 0
    region_stats.append((region.split(" (")[0], val_2025_r, round(evol, 1)))

# --- Tri décroissant par évolution ---
region_stats = sorted(region_stats, key=lambda x: -x[2])

# --- En-têtes ---
headers = ["Régions", "Nombre T3 2025", "Évolution (%) T3 2024 vs T3 2025"]
for j, h in enumerate(headers):
    table.cell(0, j).text = h
    for paragraph in table.cell(0, j).text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 32, 96)
            run.font.size = Pt(9)
            run.font.name = "Aptos"

# --- Remplissage du tableau ---
for i, (region, val_2025_r, evol) in enumerate(region_stats):
    if i + 1 >= len(table.rows):
        break
    table.cell(i + 1, 0).text = region
    table.cell(i + 1, 1).text = f"{val_2025_r:,}".replace(",", " ")
    table.cell(i + 1, 2).text = f"{evol:+.1f}%"

# --- Vider les lignes restantes s'il y en a ---
for i in range(len(region_stats) + 1, len(table.rows)):
    for j in range(3):
        table.cell(i, j).text = ""


# --- Mise en forme générale du tableau ---
for i, row in enumerate(table.rows):
    for j, cell in enumerate(row.cells):
        p = cell.text_frame.paragraphs[0]
        if i == 0:
            # Ligne des titres : on garde l'alignement centré et on ne touche pas la couleur
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)  # Taille plus grande pour titres
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 32, 96)  # Bleu titres
                run.font.name = "Aptos"
        else:
            # Lignes de données
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(9)
                run.font.name = "Aptos"
                run.font.color.rgb = RGBColor(0, 0, 0)  # Noir pour données

# --- Remplacer "Y" par total global (corrigé) ---
for shape in slide28.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() == "Y":
                    run.text = formatted_val
                    run.font.size = Pt(22)
                    run.font.name = "Aptos"
                    paragraph.alignment = PP_ALIGN.CENTER
                    break


### === SLIDE 29 ===
slide29 = prs.slides[28]

# Colonnes et secteurs sélectionnés
secteur_col = "SECTEURLIBELLE"
secteurs_selectionnes = [
    "ACTIVITÉS IMMOBILIÈRES",
    "COMMERCE , RÉPARATION D'AUTOMOBILES ET DE MOTOCYCLES",
    "ACTIVITÉS SPÉCIALISÉES, SCIENTIFIQUES ET TECHNIQUES",
    "TRANSPORTS ET ENTREPOSAGE",
    "ACTIVITÉS DE SERVICES ADMINISTRATIFS ET DE SOUTIEN",
    "HÉBERGEMENT ET RESTAURATION",
    "CONSTRUCTION"
]

# Filtrer lignes SIREN non vides et harmoniser texte
df_ip_2025_non_vides = df_ip_2025[df_ip_2025["SIREN"].astype(str).str.strip() != ""].copy()
df_ip_2025_non_vides[secteur_col] = df_ip_2025_non_vides[secteur_col].astype(str).str.strip().str.upper()
secteurs_selectionnes_upper = [s.upper() for s in secteurs_selectionnes]

# Comptage par secteur
ip_2025 = df_ip_2025_non_vides[secteur_col].value_counts()
secteurs = {sec: {"2025": ip_2025.get(sec, 0)} for sec in secteurs_selectionnes_upper}
secteurs["AUTRES"] = {"2025": ip_2025[~ip_2025.index.isin(secteurs_selectionnes_upper)].sum()}

# Total exact de SIREN non vides calculé sur slide 29
total_2025 = val_2025
parts_2025 = [round(d["2025"] / total_2025 * 100, 1) for d in secteurs.values()]

# Supprimer ancien graphique
for shape in list(slide29.shapes):
    if shape.has_chart:
        slide29.shapes._spTree.remove(shape._element)

# Création graphique
chart_data = CategoryChartData()
chart_data.categories = list(secteurs.keys())
chart_data.add_series("", parts_2025)

x, y, cx, cy = Inches(4.2), Inches(0.3), Inches(4.8), Inches(7.3)
chart = slide29.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.bold = False
data_labels.show_category_name = True
data_labels.show_value = True
data_labels.show_percentage = False

try:
    plot.doughnut_hole_size = 70
except AttributeError:
    pass

# Force via XML si nécessaire
doughnut_chart = chart.plots[0]._element
hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
if hole_size is None:
    hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
    doughnut_chart.append(hole_size)
hole_size.set('val', '70')

### === SLIDE 30 ===
slide30 = prs.slides[29]

# Colonnes et formes sélectionnées
forme_col = "FORMEJURIDIQUELIBELLE"
formes_selectionnees = [
    "Société à responsabilité limitée",
    "Société par actions simplifiée",
]

# Filtrer lignes SIREN non vides
df_ip_2025_non_vides = df_ip_2025[df_ip_2025["SIREN"].astype(str).str.strip() != ""].copy()
df_ip_2025_non_vides[forme_col] = df_ip_2025_non_vides[forme_col].astype(str).str.strip()

# Comptage par forme juridique
formes_2025 = df_ip_2025_non_vides[forme_col].value_counts()
formes = {f: {"2025": formes_2025.get(f, 0)} for f in formes_selectionnees}
formes["Autres formes juridiques"] = {"2025": formes_2025[~formes_2025.index.isin(formes_selectionnees)].sum()}

# Total exact de SIREN non vides (slide 30)
total_formes_2025 = val_2025
parts_formes_2025 = [round(d["2025"] / total_formes_2025 * 100, 1) for d in formes.values()]

# Supprimer ancien graphique
for shape in list(slide30.shapes):
    if shape.has_chart:
        slide30.shapes._spTree.remove(shape._element)

# Création graphique
chart_data = CategoryChartData()
chart_data.categories = list(formes.keys())
chart_data.add_series("", parts_formes_2025)

x, y, cx, cy = Inches(4.2), Inches(0.3), Inches(4.8), Inches(7.3)
chart = slide30.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.bold = False
data_labels.show_category_name = True
data_labels.show_value = True
data_labels.show_percentage = False

try:
    plot.doughnut_hole_size = 70
except AttributeError:
    pass

# Force via XML si nécessaire
doughnut_chart = chart.plots[0]._element
hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
if hole_size is None:
    hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
    doughnut_chart.append(hole_size)
hole_size.set('val', '70')

# === SLIDE 32 ===
slide32 = prs.slides[31]

# Filtrage robustifié sur "trésor"
df_2024_f = df_sstr_2024[
    df_sstr_2024["libelle inscription"].str.strip().str.lower() == "Privilège(s) du Trésor".lower()
]
df_2025_f = df_sstr_2025[
    df_sstr_2025["libelle inscription"].str.strip().str.lower() == "Privilège(s) du Trésor".lower()
]

# --- Filtrer lignes SIREN non vides ---
df_2024_f_non_vides = df_2024_f[df_2024_f["SIREN"].astype(str).str.strip() != ""]
df_2025_f_non_vides = df_2025_f[df_2025_f["SIREN"].astype(str).str.strip() != ""]

# Comptage INSCRIPTIONS
n_2024 = len(df_2024_f_non_vides)
n_2025 = len(df_2025_f_non_vides)
var_n = ((n_2025 - n_2024)/n_2024*100) if n_2024 > 0 else 0

# Montant moyen
m_2024 = df_2024_f_non_vides["montant"].mean() if not df_2024_f_non_vides.empty else 0
m_2025 = df_2025_f_non_vides["montant"].mean() if not df_2025_f_non_vides.empty else 0
var_m = ((m_2025 - m_2024)/m_2024*100) if m_2024 > 0 else 0

# Formatage
f_n2025 = f"{n_2025:,}".replace(",", " ")
f_n2024 = f"{n_2024:,}".replace(",", " ")
f_var_n = f"{var_n:+.1f}%"

f_m2025 = f"{m_2025:,.0f} €".replace(",", " ")
f_m2024 = f"{m_2024:,.0f} €".replace(",", " ")
f_var_m = f"{var_m:+.1f}%"

# Remplacement dans slide 32
for shape in slide32.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text == "A":
                    run.text = f_n2025
                    run.font.size = Pt(20)
                elif text == "B":
                    run.text = f_n2024
                    run.font.size = Pt(20)
                elif text == "X":  # pourcentage nombre
                    run.text = f_var_n
                    run.font.size = Pt(12)
                elif text == "C":
                    run.text = f_m2025
                    run.font.size = Pt(20)
                elif text == "D":
                    run.text = f_m2024
                    run.font.size = Pt(20)
                elif text == "Y":  # pourcentage montant
                    run.text = f_var_m
                    run.font.size = Pt(12)

                paragraph.alignment = PP_ALIGN.CENTER
                run.font.name = "Aptos"


# === SLIDE 34 ===
slide34 = prs.slides[33]

# Filtrage robustifié sur "sécurité sociale"
df_2024_f = df_sstr_2024[
    df_sstr_2024["libelle inscription"].str.strip().str.lower() == "Privilège(s) sécurité sociale, régimes complémentaires".lower()
]
df_2025_f = df_sstr_2025[
    df_sstr_2025["libelle inscription"].str.strip().str.lower() == "Privilège(s) sécurité sociale, régimes complémentaires".lower()
]

# --- Comptage des lignes non vides dans SIREN ---
df_2024_f_non_vides = df_2024_f[df_2024_f["SIREN"].astype(str).str.strip() != ""]
df_2025_f_non_vides = df_2025_f[df_2025_f["SIREN"].astype(str).str.strip() != ""]

n_2024 = len(df_2024_f_non_vides)
n_2025 = len(df_2025_f_non_vides)
var_n = ((n_2025 - n_2024)/n_2024*100) if n_2024 > 0 else 0

# Montant moyen
m_2024 = df_2024_f_non_vides["montant"].mean() if not df_2024_f_non_vides.empty else 0
m_2025 = df_2025_f_non_vides["montant"].mean() if not df_2025_f_non_vides.empty else 0
var_m = ((m_2025 - m_2024)/m_2024*100) if m_2024 > 0 else 0

# Formatage des valeurs
f_n2025 = f"{n_2025:,}".replace(",", " ")
f_n2024 = f"{n_2024:,}".replace(",", " ")
f_var_n = f"{var_n:+.1f}%"

f_m2025 = f"{m_2025:,.0f} €".replace(",", " ")
f_m2024 = f"{m_2024:,.0f} €".replace(",", " ")
f_var_m = f"{var_m:+.1f}%"

# Remplacement des valeurs dans la slide
for shape in slide34.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text == "A":
                    run.text = f_n2025
                    run.font.size = Pt(20)
                elif text == "B":
                    run.text = f_n2024
                    run.font.size = Pt(20)
                elif text == "X":  # pourcentage nombre
                    run.text = f_var_n
                    run.font.size = Pt(12)
                elif text == "C":
                    run.text = f_m2025
                    run.font.size = Pt(20)
                elif text == "D":
                    run.text = f_m2024
                    run.font.size = Pt(20)
                elif text == "Y":  # pourcentage montant
                    run.text = f_var_m
                    run.font.size = Pt(12)

                paragraph.alignment = PP_ALIGN.CENTER
                run.font.name = "Aptos"

# === SLIDE 36 ===
slide36 = prs.slides[35]
type_col = "TYPE_DEPOT_BE"
types_selectionnes = ["DEPOT INITIAL", "RECTIFICATIF", "REGULARISATION"]

# --- Vérification de la colonne ---
for df in [df_mbe_2024, df_mbe_2025]:
    if type_col not in df.columns:
        raise KeyError(f"Colonne {type_col} introuvable dans le fichier")

# --- Copie et nettoyage ---
df_mbe_2024 = df_mbe_2024.copy()
df_mbe_2025 = df_mbe_2025.copy()
df_mbe_2024[type_col] = df_mbe_2024[type_col].astype(str).str.strip().str.upper()
df_mbe_2025[type_col] = df_mbe_2025[type_col].astype(str).str.strip().str.upper()

# --- Filtrer SIREN non vides ---
df_mbe_2024_nv = df_mbe_2024[df_mbe_2024["SIREN"].astype(str).str.strip() != ""]
df_mbe_2025_nv = df_mbe_2025[df_mbe_2025["SIREN"].astype(str).str.strip() != ""]

# --- Comptage par type ---
mbe_2024 = df_mbe_2024_nv.groupby(type_col)["SIREN"].count()
mbe_2025 = df_mbe_2025_nv.groupby(type_col)["SIREN"].count()

# --- Calcul des répartitions ---
total_2025 = mbe_2025.sum()
repartition_2025 = {}
for t in types_selectionnes:
    val = mbe_2025.get(t, 0)
    pct = round(val / total_2025 * 100, 1) if total_2025 > 0 else 0
    repartition_2025[t] = {"val": val, "pct": pct}

# --- Récupération du tableau existant ---
table_shapes = [s for s in slide36.shapes if s.has_table]
if not table_shapes:
    print("Aucun tableau trouvé sur la slide 36 — tableau ignoré")
else:
    table = table_shapes[0].table
    n_rows = len(table.rows)
    n_cols = len(table.columns)

    # Remplir le tableau sans changer sa structure
    for i, t in enumerate(types_selectionnes):
        if i + 1 >= n_rows:
            break
        table.cell(i + 1, 0).text = t.title()
        table.cell(i + 1, 1).text = f"{repartition_2025[t]['val']:,}".replace(",", " ")
        table.cell(i + 1, 2).text = f"{repartition_2025[t]['pct']:.1f}%"

    # Vider les lignes restantes
    for i in range(len(types_selectionnes) + 1, n_rows):
        for j in range(n_cols):
            table.cell(i, j).text = ""

# --- Graphique circulaire (mise à jour complète) ---
chart_shapes = [s for s in slide36.shapes if hasattr(s, "has_chart") and s.has_chart]
if chart_shapes:
    chart = chart_shapes[0].chart
    chart_data = CategoryChartData()

    # Ajouter les catégories et valeurs selon les données 2025
    chart_data.categories = [t.title() for t in types_selectionnes]
    chart_data.add_series("Répartition T3 2025", [repartition_2025[t]["val"] for t in types_selectionnes])

    # Remplacer les données du graphique
    chart.replace_data(chart_data)

    # Personnaliser les étiquettes
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_category_name = True
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.font.size = Pt(9)
    data_labels.font.bold = False
else:
    print("Aucun graphique modifiable détecté sur la slide 36 — vérifie que c’est bien un camembert PowerPoint")

# === SLIDE 37 ===
slide37 = prs.slides[36]

# --- Calcul évolution par forme juridique ---
forme_col = "FORMEJURIDIQUELIBELLEAGREGE"
formes_selectionnees = [
    "Société à responsabilité limitée",
    "Société civile immobilière",
    "Société par actions simplifiée"
]

# Harmonisation texte
df_mbe_2024[forme_col] = df_mbe_2024[forme_col].astype(str).str.strip()
df_mbe_2025[forme_col] = df_mbe_2025[forme_col].astype(str).str.strip()

# Filtrer SIREN non vides
df_mbe_2024_nv = df_mbe_2024[df_mbe_2024["SIREN"].astype(str).str.strip() != ""].copy()
df_mbe_2025_nv = df_mbe_2025[df_mbe_2025["SIREN"].astype(str).str.strip() != ""].copy()

# Comptage par forme juridique basé sur SIREN non vides
formes_2024 = df_mbe_2024_nv.groupby(forme_col)["SIREN"].count()
formes_2025 = df_mbe_2025_nv.groupby(forme_col)["SIREN"].count()

# Construire dictionnaire
formes = {}
for f in formes_selectionnees:
    formes[f] = {
        "2024": formes_2024.get(f, 0),
        "2025": formes_2025.get(f, 0)
    }

# Catégorie AUTRES
autres_2024 = formes_2024[~formes_2024.index.isin(formes_selectionnees)].sum()
autres_2025 = formes_2025[~formes_2025.index.isin(formes_selectionnees)].sum()
formes["Autres formes juridiques"] = {"2024": autres_2024, "2025": autres_2025}

# Calcul évolutions %
for f in formes:
    val2024 = formes[f]["2024"]
    val2025 = formes[f]["2025"]
    evol = ((val2025 - val2024) / val2024 * 100) if val2024 > 0 else 0
    formes[f]["evol"] = f"{evol:+.1f}%"


# --- Remplissage tableau ---
for shape in slide37.shapes:
    if shape.has_table:
        table = shape.table
        # Mise à jour des en-têtes
        table.cell(0, 0).text = "Nombre de dépôts par catégorie T3 2025"
        table.cell(0, 1).text = "Nombre T3 2025"
        table.cell(0, 2).text = "Évolution (%) T3 2024 vs T3 2025"

        # Format des en-têtes (ligne 0)
        for col_idx in range(3):
            p = table.cell(0, col_idx).text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 32, 96)  # Bleu titres
                run.font.name = "Aptos"

        # Remplissage des données à partir de la 2e ligne
        for row_idx in range(1, len(table.rows)):  # 1 = ignore la ligne des titres
            forme_nom = table.cell(row_idx, 0).text.strip().replace('\u200b', '')
            if forme_nom in formes:
                table.cell(row_idx, 1).text = f"{formes[forme_nom]['2025']:,}".replace(",", " ")
                table.cell(row_idx, 2).text = formes[forme_nom]["evol"]

                # Mise en forme des lignes de données
                for col_idx in range(3):
                    p = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.size = Pt(9)
                        run.font.name = "Aptos"
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Noir pour données

# Mise à jour diagramme
for shape in list(slide37.shapes):
    if shape.has_chart:
        slide37.shapes._spTree.remove(shape._element)

total_2025 = sum([formes[f]["2025"] for f in formes])
parts_2025 = [round(formes[f]["2025"] / total_2025 * 100, 1) for f in formes]

chart_data = CategoryChartData()
chart_data.categories = list(formes.keys())
chart_data.add_series("", parts_2025)

x, y, cx, cy = Inches(8.5), Inches(0.1), Inches(4.8), Inches(7.3)
chart = slide37.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

data_labels.font.size = Pt(9)
data_labels.font.bold = False
data_labels.show_category_name = True
data_labels.show_value = True
data_labels.show_percentage = False

# Correction épaisseur anneau
try:
    plot.doughnut_hole_size = 70
except AttributeError:
    pass

from lxml import etree
doughnut_chart = chart.plots[0]._element
hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
if hole_size is None:
    hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
    doughnut_chart.append(hole_size)
hole_size.set('val', '70')

### === SLIDE 39 ===
slide39 = prs.slides[38]

# --- Colonne "Nature de la demande" ---
nature_col = "Nature de la demande"

# Filtrer lignes avec SIREN non vide
df_div_2024_nv = df_div_2024[df_div_2024["SIREN"].astype(str).str.strip() != ""].copy()
df_div_2025_nv = df_div_2025[df_div_2025["SIREN"].astype(str).str.strip() != ""].copy()

# Comptage brut par nature
nature_2024 = df_div_2024_nv[nature_col].astype(str).str.strip().value_counts()
nature_2025 = df_div_2025_nv[nature_col].astype(str).str.strip().value_counts()

# Union des catégories
categories = sorted(set(nature_2024.index).union(nature_2025.index))

# Construction dictionnaire avec évolutions
natures = {}
for cat in categories:
    val_2024 = nature_2024.get(cat, 0)
    val_2025 = nature_2025.get(cat, 0)
    evol = ((val_2025 - val_2024) / val_2024 * 100) if val_2024 > 0 else 0
    natures[cat] = {
        "2024": val_2024,
        "2025": val_2025,
        "evol": f"{evol:+.1f}%"
    }

# Totaux
total_2024 = sum(v["2024"] for v in natures.values())
total_2025 = sum(v["2025"] for v in natures.values())
total_evol = ((total_2025 - total_2024) / total_2024 * 100) if total_2024 > 0 else 0
natures["Total général"] = {
    "2024": total_2024,
    "2025": total_2025,
    "evol": f"{total_evol:+.1f}%"
}

# === TABLEAU EXISTANT ===
table_shapes = [shape for shape in slide39.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("Aucun tableau trouvé sur la slide 39 — vérifie que le modèle contient bien un tableau.")
else:
    table = table_shapes[0].table

# --- En-têtes ---
headers = ["Nature de la demande", "T3 2025", "Évolution (%) T3 2024 vs T3 2025"]
for j, h in enumerate(headers):
    table.cell(0, j).text = h
    for paragraph in table.cell(0, j).text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.bold = True
            run.font.size = Pt(10)
            run.font.name = "Aptos"
    # Fond bleu sur la ligne d’en-tête
    fill = table.cell(0, j).fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 32, 96)

# --- Remplissage du tableau ---
for i, (nature, data) in enumerate(natures.items()):
    if i + 1 >= len(table.rows):
        break
    table.cell(i + 1, 0).text = nature
    table.cell(i + 1, 1).text = f"{data['2025']:,}".replace(",", " ")
    table.cell(i + 1, 2).text = data["evol"]

# Efface les lignes restantes si le tableau est plus grand
for i in range(len(natures) + 1, len(table.rows)):
    for j in range(3):
        table.cell(i, j).text = ""

# --- Mise en forme globale ---
for i, row in enumerate(table.rows):
    for j, cell in enumerate(row.cells):
        p = cell.text_frame.paragraphs[0]
        # Alignement : 1re colonne à gauche, les autres centrées
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(9)
            run.font.name = "Aptos"


### === SLIDE 41 ===
slide41 = prs.slides[40]

val_2025 = df_dpca_2025.shape[0]

# Formatage
formatted_val = f"{val_2025:,}".replace(",", " ")  # espace insécable

# Supprimer tableau existant
for shape in list(slide41.shapes):
    if shape.has_table:
        slide41.shapes._spTree.remove(shape._element)

# Colonnes
region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

# Comptage par région
dpca_region_2025 = df_dpca_2025[region_col].value_counts().sort_index()
dpca_region_2025 = dpca_region_2025[~dpca_region_2025.index.isin(regions_exclues)]
total_dpca_2025 = dpca_region_2025.sum()
parts_region = (dpca_region_2025 / total_dpca_2025 * 100).round(1)

# Simplifier noms ("Occitanie (Languedoc-Roussillon-Midi-Pyrénées)" etc.)
parts_region_tuples = [
    (region.split(" (")[0], f"{pct:.1f}%")
    for region, pct in parts_region.items()
]

# Tri décroissant par pourcentage
parts_region_tuples = sorted(parts_region_tuples, key=lambda x: -float(x[1].replace("%", "")))

# Création tableau
rows = len(parts_region_tuples) + 1
cols = 2
left = Inches(0.7)
top = Inches(1)
width = Inches(5.5)
height = Inches(0.25 * rows + 0.5)

table_shape = slide41.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.cell(0, 0).text = "Régions"
table.cell(0, 1).text = "Répartition (en %)"

# Remplissage
for i, (region, pct) in enumerate(parts_region_tuples):
    table.cell(i+1, 0).text = region
    table.cell(i+1, 1).text = pct


# Mise en forme
for row_idx, row in enumerate(table.rows):
    for cell in row.cells:
        p = cell.text_frame.paragraphs[0]
        if row_idx == 0:
            # Ligne des titres : on garde centré et applique style spécifique
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 32, 96)  # Bleu titres
                run.font.name = "Aptos"
        else:
            # Lignes de données
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(9)
                run.font.name = "Aptos"
                run.font.color.rgb = RGBColor(0, 0, 0)  # Noir pour données

# --- Remplacer "Y" par total global (corrigé) ---
for shape in slide41.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() == "Y":
                    run.text = formatted_val
                    run.font.size = Pt(24)
                    run.font.name = "Aptos"
                    paragraph.alignment = PP_ALIGN.CENTER
                    break



# === SLIDE 42 ===
slide42 = prs.slides[41]

# --- Nettoyage : supprimer anciens graphiques ET tableaux (pas de tableau requis sur cette slide) ---
for shape in list(slide42.shapes):
    if (hasattr(shape, "has_chart") and shape.has_chart) or (hasattr(shape, "has_table") and shape.has_table):
        slide42.shapes._spTree.remove(shape._element)

# --- Données ---
secteur_col = "SECTEURLIBELLE"
secteurs_selectionnes = [
    "ACTIVITÉS IMMOBILIÈRES",
    "COMMERCE , RÉPARATION D'AUTOMOBILES ET DE MOTOCYCLES",
    "ACTIVITÉS SPÉCIALISÉES, SCIENTIFIQUES ET TECHNIQUES",
    "TRANSPORTS ET ENTREPOSAGE",
    "ACTIVITÉS DE SERVICES ADMINISTRATIFS ET DE SOUTIEN",
    "HÉBERGEMENT ET RESTAURATION",
    "CONSTRUCTION"
]

# Normalisation
df_dpca_2024[secteur_col] = df_dpca_2024[secteur_col].astype(str).str.strip().str.upper()
df_dpca_2025[secteur_col] = df_dpca_2025[secteur_col].astype(str).str.strip().str.upper()

# Comptages
dpca_2024 = df_dpca_2024[secteur_col].value_counts()
dpca_2025 = df_dpca_2025[secteur_col].value_counts()

# Dictionnaire secteurs + AUTRES
secteurs = {}
for sec in secteurs_selectionnes:
    secteurs[sec] = {"2024": int(dpca_2024.get(sec, 0)), "2025": int(dpca_2025.get(sec, 0))}

autres_2024 = int(dpca_2024[~dpca_2024.index.isin(secteurs_selectionnes)].sum())
autres_2025 = int(dpca_2025[~dpca_2025.index.isin(secteurs_selectionnes)].sum())
secteurs["AUTRES"] = {"2024": autres_2024, "2025": autres_2025}

# Évolution %
for sec in list(secteurs.keys()):
    v24 = secteurs[sec]["2024"]
    v25 = secteurs[sec]["2025"]
    evol = ((v25 - v24) / v24 * 100) if v24 > 0 else 0.0
    secteurs[sec]["evol_num"] = round(evol, 1)


total_2025 = sum(secteurs[s]["2025"] for s in secteurs)
parts_2025 = [
    (round(secteurs[s]["2025"] / total_2025 * 100, 1) if total_2025 > 0 else 0.0)
    for s in secteurs
]

chart_data_left = CategoryChartData()
chart_data_left.categories = list(secteurs.keys())
chart_data_left.add_series("Répartition 2025", parts_2025)

x, y, cx, cy = Inches(0.6), Inches(0.6), Inches(6.0), Inches(6.5)
chart_left = slide42.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_left).chart

chart_left.has_legend = False
plot_left = chart_left.plots[0]
plot_left.has_data_labels = True
dl_left = plot_left.data_labels
dl_left.font.size = Pt(9)
dl_left.show_category_name = True
dl_left.show_percentage = True
dl_left.show_value = False

try:
    plot_left.doughnut_hole_size = 65
except AttributeError:
    pass

chart_data_right = CategoryChartData()
chart_data_right.categories = list(secteurs.keys())
chart_data_right.add_series("Évolution (%)", [secteurs[s]["evol_num"] for s in secteurs])

x2, y2, cx2, cy2 = Inches(7.0), Inches(0.6), Inches(6.0), Inches(6.5)
chart_right = slide42.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x2, y2, cx2, cy2, chart_data_right).chart

chart_right.has_legend = False
plot_right = chart_right.plots[0]
plot_right.has_data_labels = True

# Labels hors barres + couleurs conditionnelles par point
ordered_keys = list(secteurs.keys())
for i, point in enumerate(plot_right.series[0].points):
    val = secteurs[ordered_keys[i]]["evol_num"]

    # Label
    lbl = point.data_label
    lbl.has_text_frame = True
    lbl.text_frame.clear()
    p = lbl.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{val:+.1f}%"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0) if val >= 0 else RGBColor(0, 176, 80)
    lbl.position = XL_LABEL_POSITION.OUTSIDE_END

    # Couleur barre
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = RGBColor(255, 0, 0) if val >= 0 else RGBColor(0, 176, 80)

# Axes
value_axis = chart_right.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.number_format = '0.0"%"'
value_axis.tick_labels.font.size = Pt(9)

category_axis = chart_right.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True

# === SLIDE 43 ===
slide43 = prs.slides[42]

# --- Nettoyage : supprimer anciens graphiques et tableaux ---
for shape in list(slide43.shapes):
    if (hasattr(shape, "has_chart") and shape.has_chart) or (hasattr(shape, "has_table") and shape.has_table):
        slide43.shapes._spTree.remove(shape._element)

# --- Préparer les données ---
forme_col = "FORMEJURIDIQUELIBELLE"
df_dpca_2024[forme_col] = df_dpca_2024[forme_col].astype(str).str.strip()
df_dpca_2025[forme_col] = df_dpca_2025[forme_col].astype(str).str.strip()

formes_selectionnees = [
    "Société à responsabilité limitée",
    "Société par actions simplifiée"
]

formes_2024 = df_dpca_2024[forme_col].value_counts()
formes_2025 = df_dpca_2025[forme_col].value_counts()

formes = {}
for f in formes_selectionnees:
    formes[f] = {
        "2024": formes_2024.get(f, 0),
        "2025": formes_2025.get(f, 0)
    }

# Catégorie AUTRES
exclus = formes_selectionnees + ["Commerçant", "Société civile immobilière"]
autres_2024 = formes_2024[~formes_2024.index.isin(exclus)].sum()
autres_2025 = formes_2025[~formes_2025.index.isin(exclus)].sum()
formes["Autres formes juridiques"] = {"2024": autres_2024, "2025": autres_2025}

# Calcul évolutions
for f in formes:
    v2024 = formes[f]["2024"]
    v2025 = formes[f]["2025"]
    evol = ((v2025 - v2024) / v2024 * 100) if v2024 > 0 else 0
    formes[f]["evol_num"] = round(evol, 1)

total_2025 = sum(formes[f]["2025"] for f in formes)
parts_2025 = [
    (round(formes[f]["2025"] / total_2025 * 100, 1) if total_2025 > 0 else 0.0)
    for f in formes
]

chart_data_left = CategoryChartData()
chart_data_left.categories = list(formes.keys())
chart_data_left.add_series("Répartition 2025", parts_2025)

x, y, cx, cy = Inches(0.6), Inches(0.6), Inches(6.0), Inches(6.5)
chart_left = slide43.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_left).chart

chart_left.has_legend = False
plot_left = chart_left.plots[0]
plot_left.has_data_labels = True
dl_left = plot_left.data_labels
dl_left.font.size = Pt(9)
dl_left.show_category_name = True
dl_left.show_percentage = True
dl_left.show_value = False

try:
    plot_left.doughnut_hole_size = 65
except AttributeError:
    pass

chart_data_right = CategoryChartData()
chart_data_right.categories = list(formes.keys())
chart_data_right.add_series("Évolution (%)", [formes[f]["evol_num"] for f in formes])

x2, y2, cx2, cy2 = Inches(7.0), Inches(0.6), Inches(6.0), Inches(6.5)
chart_right = slide43.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x2, y2, cx2, cy2, chart_data_right).chart

chart_right.has_legend = False
plot_right = chart_right.plots[0]
plot_right.has_data_labels = True

# Labels hors barres + couleurs conditionnelles
ordered_keys = list(formes.keys())
for i, point in enumerate(plot_right.series[0].points):
    val = formes[ordered_keys[i]]["evol_num"]

    # Label
    lbl = point.data_label
    lbl.has_text_frame = True
    lbl.text_frame.clear()
    p = lbl.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{val:+.1f}%"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0) if val >= 0 else RGBColor(0, 176, 80)
    lbl.position = XL_LABEL_POSITION.OUTSIDE_END

    # Couleur barre
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = RGBColor(255, 0, 0) if val >= 0 else RGBColor(0, 176, 80)

# Axes
value_axis = chart_right.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.number_format = '0.0"%"'
value_axis.tick_labels.font.size = Pt(9)

category_axis = chart_right.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True

### === SLIDE 44 ===
slide44 = prs.slides[43]

# Supprimer ancien graphique s'il existe
for shape in list(slide44.shapes):
    if shape.has_chart:
        slide44.shapes._spTree.remove(shape._element)

# Colonne région
region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

# Comptages 2024 et 2025
regions_2024 = df_dpca_2024[region_col].fillna("").astype(str).str.strip()
regions_2025 = df_dpca_2025[region_col].fillna("").astype(str).str.strip()

regions_2024 = regions_2024[~regions_2024.isin([""] + regions_exclues)].value_counts()
regions_2025 = regions_2025[~regions_2025.isin([""] + regions_exclues)].value_counts()

# Liste complète des régions (présentes en 2024 ou 2025)
regions_full = sorted(set(regions_2024.index).union(set(regions_2025.index)))
regions_labels = [r.split(" (")[0] for r in regions_full]

# Calcul % évolution
evolutions = []
for r in regions_full:
    val2024 = regions_2024.get(r, 0)
    val2025 = regions_2025.get(r, 0)
    evol = ((val2025 - val2024) / val2024 * 100) if val2024 > 0 else 0
    evolutions.append(evol)

# --- Création graphique en barres ---
chart_data = CategoryChartData()
chart_data.categories = regions_labels
chart_data.add_series("", evolutions)

x, y, cx, cy = Inches(2.5), Inches(1.0), Inches(8.5), Inches(5.5)
chart = slide44.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

# Pas de légende
chart.has_legend = False

plot = chart.plots[0]
plot.has_data_labels = True

# Couleur unique
for series in plot.series:
    for p_idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(146, 208, 80)

        # Label affichant le % arrondi
        dLbl = point.data_label
        dLbl.has_text_frame = True
        dLbl.text_frame.clear()
        p = dLbl.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{evolutions[p_idx]:+.1f}%"
        run.font.size = Pt(9)
        run.font.bold = True
        dLbl.position = XL_LABEL_POSITION.OUTSIDE_END

# Axe des valeurs (%)
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.font.size = Pt(9)
value_axis.tick_labels.number_format = "0%"

# Axe des catégories (régions alignées à gauche)
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True
category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW


# ===== SLIDE 46 =====
slide46 = prs.slides[45]

# --- Récupérer le tableau existant ---
table_shapes = [shape for shape in slide46.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 46 — vérifie le modèle PowerPoint.")
else:
    table = table_shapes[0].table

# --- Préparation des données ---
region_col = "REGION"
val_col = "NOMBRE"
stat_col = "STATISTIQUE"
valeur_filtre = "Nombre total de dossiers détectés"

# Filtrage
df_prev_2025_filt = df_prev_2025[df_prev_2025[stat_col] == valeur_filtre].copy()

# Agrégation par région
prev_region_2025 = df_prev_2025_filt.groupby(region_col)[val_col].sum()

# Total global
total_2025 = prev_region_2025.sum()
formatted_total = f"{total_2025:,.0f}".replace(",", " ")

# Calcul des pourcentages
region_pct = (prev_region_2025 / total_2025 * 100).round(1).sort_values(ascending=False)

# --- Remplissage du tableau ---
# On garde la première ligne (titres)
for i, (region, pct) in enumerate(region_pct.items(), start=1):
    if i < len(table.rows):
        table.cell(i, 0).text = region
        table.cell(i, 1).text = f"{pct:.1f} %"

# Vider les lignes restantes
for i in range(len(region_pct) + 1, len(table.rows)):
    for j in range(2):
        table.cell(i, j).text = ""

# Mise en forme
for row in table.rows:
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"

# --- Remplacer "Y" (valeur totale) ---
for shape in slide46.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() == "Y":
                    run.text = formatted_total
                    run.font.size = Pt(22)
                    run.font.name = "Aptos"
                    paragraph.alignment = PP_ALIGN.CENTER
                    break

# ===== SLIDE 47 =====

slide47 = prs.slides[46]

# --- Récupérer le tableau existant ---
table_shapes = [shape for shape in slide47.shapes if shape.has_table]
if not table_shapes:
    raise ValueError("⚠️ Aucun tableau trouvé sur la slide 47 — vérifie le modèle PowerPoint.")
else:
    table = table_shapes[0].table

# --- Préparation des données ---
region_col = "REGION"
val_col = "NOMBRE"
stat_col = "STATISTIQUE"
valeur_filtre = "Nombre total de dossiers de prévention étudiés"

# Filtrage
df_prev_2025_filt = df_prev_2025[df_prev_2025[stat_col] == valeur_filtre].copy()

# Agrégation par région
prev_region_2025 = df_prev_2025_filt.groupby(region_col)[val_col].sum()

# Total global
total_2025 = prev_region_2025.sum()
formatted_total = f"{total_2025:,.0f}".replace(",", " ")

# Calcul des pourcentages
region_pct = (prev_region_2025 / total_2025 * 100).round(1).sort_values(ascending=False)

# --- Remplissage du tableau ---
# On garde la première ligne (titres)
for i, (region, pct) in enumerate(region_pct.items(), start=1):
    if i < len(table.rows):
        table.cell(i, 0).text = region
        table.cell(i, 1).text = f"{pct:.1f} %"

# Vider les lignes restantes
for i in range(len(region_pct) + 1, len(table.rows)):
    for j in range(2):
        table.cell(i, j).text = ""

# Mise en forme
for row in table.rows:
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "Aptos"

# --- Remplacer "Y" (valeur totale) ---
for shape in slide47.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() == "Y":
                    run.text = formatted_total
                    run.font.size = Pt(22)
                    run.font.name = "Aptos"
                    paragraph.alignment = PP_ALIGN.CENTER
                    break

### === SLIDE 49 ===
slide49 = prs.slides[48]

# --- Supprimer ancien graphique ---
for shape in list(slide49.shapes):
    if shape.has_chart:
        slide49.shapes._spTree.remove(shape._element)

# Calcul nombre total entreprises Mission
total_2024 = df_miss_2024.shape[0]
total_2025 = df_miss_2025.shape[0]

# --- Création graphique vertical (2 barres) ---
chart_data = CategoryChartData()
chart_data.categories = ["T3 2024", "T3 2025"]
chart_data.add_series("", [total_2024, total_2025])

x, y, cx, cy = Inches(4.5), Inches(1.0), Inches(4.0), Inches(5.0)
chart_shape = slide49.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
chart = chart_shape.chart

# Retirer la légende
chart.has_legend = False

# Couleurs : vert clair pour 2024, bleu clair pour 2025
serie_colors = [RGBColor(0, 32, 96), RGBColor(91, 155, 213)]
for s_idx, point in enumerate(chart.plots[0].series[0].points):
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = serie_colors[s_idx]

# Valeurs au-dessus des barres
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.position = XL_LABEL_POSITION.ABOVE
for s_idx, point in enumerate(plot.series[0].points):
    dLbl = point.data_label
    dLbl.has_text_frame = True
    dLbl.text_frame.clear()
    p = dLbl.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{[total_2024, total_2025][s_idx]:d}"
    run.font.size = Pt(12)
    run.font.bold = True

# Axe des valeurs
value_axis = chart.value_axis
value_axis.minimum_scale = 0   # <- démarrage à 0
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.font.size = Pt(10)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(10)
category_axis.tick_labels.font.bold = True

### === SLIDE 50 ===
slide50 = prs.slides[49]

# Supprimer ancien graphique s'il existe
for shape in list(slide50.shapes):
    if shape.has_chart:
        slide50.shapes._spTree.remove(shape._element)

# Colonne région
region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

# Comptages entreprises Mission
miss_2024 = df_miss_2024[region_col].fillna("").astype(str).str.strip()
miss_2025 = df_miss_2025[region_col].fillna("").astype(str).str.strip()

miss_2024 = miss_2024[miss_2024 != ""].value_counts()
miss_2025 = miss_2025[miss_2025 != ""].value_counts()

# Exclure régions non voulues
miss_2024 = miss_2024[~miss_2024.index.isin(regions_exclues)]
miss_2025 = miss_2025[~miss_2025.index.isin(regions_exclues)]

# Liste des régions (union des deux années)
regions_full = sorted(set(miss_2024.index).union(set(miss_2025.index)))
regions_labels = [r.split(" (")[0] for r in regions_full]

# Valeurs par région
val_2024 = [miss_2024.get(r, 0) for r in regions_full]
val_2025 = [miss_2025.get(r, 0) for r in regions_full]

# --- Création graphique groupé (barres horizontales) ---
chart_data = CategoryChartData()
chart_data.categories = regions_labels
chart_data.add_series("T3 2024", val_2024)
chart_data.add_series("T3 2025", val_2025)

x, y, cx, cy = Inches(2.2), Inches(1.3), Inches(8), Inches(5)
chart = slide50.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

# Retirer la légende
chart.has_legend = False

plot = chart.plots[0]
plot.has_data_labels = True

# Coordonnées de départ pour la légende (sous le graphique)
legend_y = y + cy + Inches(0.2)
square_size = Inches(0.2)
text_offset = Inches(0.05)  # espace entre carré et texte
item_spacing = Inches(1.0)   # espacement entre items
legend_items = [("T3 2024", RGBColor(0, 32, 96)), ("T3 2025", RGBColor(91, 155, 213))]

# Calcul largeur totale de la légende
total_width = len(legend_items) * (square_size + text_offset + Inches(0.5)) + (len(legend_items)-1)*item_spacing

# Position de départ pour centrer
legend_x = x + (cx - total_width)/2

for i, (label, color) in enumerate(legend_items):
    x_pos = legend_x + i * (square_size + text_offset + Inches(0.5) + item_spacing)

    # Carré coloré
    square = slide50.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x_pos,
        legend_y,
        square_size,
        square_size
    )
    square.fill.solid()
    square.fill.fore_color.rgb = color
    square.line.fill.background()  # pas de contour

    # Texte à côté
    txt_box = slide50.shapes.add_textbox(
        x_pos + square_size + text_offset,
        legend_y,
        Inches(0.5),
        square_size
    )
    txt_frame = txt_box.text_frame
    txt_frame.clear()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.bold = True
    p.alignment = PP_ALIGN.LEFT

# Couleurs : vert clair (2024), bleu clair (2025)
serie_colors = [RGBColor(0, 32, 96), RGBColor(91, 155, 213)]
val_lists = [val_2024, val_2025]

# Appliquer couleurs + labels
for s_idx, series in enumerate(plot.series):
    for p_idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = serie_colors[s_idx]

        # label (nombre d'entreprises)
        dLbl = point.data_label
        dLbl.has_text_frame = True
        dLbl.text_frame.clear()
        p = dLbl.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{val_lists[s_idx][p_idx]:d}"
        run.font.size = Pt(9)
        run.font.bold = True
        dLbl.position = XL_LABEL_POSITION.OUTSIDE_END

# Axe des valeurs
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.font.size = Pt(9)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True

### === SLIDE 51 ===
slide51 = prs.slides[50]

# --- Calcul répartition par forme juridique ---
forme_col = "SECTEURLIBELLE"
df_miss_2025[forme_col] = df_miss_2025[forme_col].astype(str).str.strip()

formes_selectionnees = [
    "ACTIVITÉS IMMOBILIÈRES",
    "COMMERCE , RÉPARATION D'AUTOMOBILES ET DE MOTOCYCLES",
    "ACTIVITÉS SPÉCIALISÉES, SCIENTIFIQUES ET TECHNIQUES",
    "TRANSPORTS ET ENTREPOSAGE",
    "ACTIVITÉS DE SERVICES ADMINISTRATIFS ET DE SOUTIEN",
    "HÉBERGEMENT ET RESTAURATION",
    "CONSTRUCTION"
]

# Comptages des formes juridiques pour 2024 et 2025
formes_2025 = df_miss_2025[forme_col].value_counts()

formes = {}
for f in formes_selectionnees:
    formes[f] = {
        "2025": formes_2025.get(f, 0)
    }

# Catégorie AUTRES
autres_2025 = formes_2025[~formes_2025.index.isin(formes_selectionnees)].sum()
formes["Autres"] = {"2025": autres_2025}

# --- Mise à jour diagramme ---
for shape in list(slide51.shapes):
    if shape.has_chart:
        slide51.shapes._spTree.remove(shape._element)

total_formes_2025 = sum([formes[f]["2025"] for f in formes])
parts_formes_2025 = [round(formes[f]["2025"] / total_formes_2025 * 100, 1) for f in formes]

chart_data = CategoryChartData()
chart_data.categories = list(formes.keys())
chart_data.add_series("", parts_formes_2025)

x, y, cx, cy = Inches(3.9), Inches(0.1), Inches(4.8), Inches(7.3)
chart = slide51.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

# Modifier la taille de police des étiquettes
data_labels.font.size = Pt(9)
data_labels.font.bold = False

data_labels.show_category_name = True
data_labels.show_value = True   # % affichés
data_labels.show_percentage = False

# --- Correction épaisseur anneau ---
try:
    plot.doughnut_hole_size = 70  # valeur si supportée par ta version de python-pptx
except AttributeError:
    pass

# Force via XML si jamais l’attribut n’est pas pris en compte
doughnut_chart = chart.plots[0]._element
hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
if hole_size is None:
    from lxml import etree
    hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
    doughnut_chart.append(hole_size)
hole_size.set('val', '70')  # entre 10 et 90 (40 = anneau plus épais, 80 = plus fin)

### === SLIDE 53 ===
slide53 = prs.slides[52]

# --- Supprimer ancien graphique ---
for shape in list(slide53.shapes):
    if shape.has_chart:
        slide53.shapes._spTree.remove(shape._element)

# Calcul nombre total entreprises ESS
total_2024 = df_ess_2024.shape[0]
total_2025 = df_ess_2025.shape[0]

# --- Création graphique vertical (2 barres) ---
chart_data = CategoryChartData()
chart_data.categories = ["T3 2024", "T3 2025"]
chart_data.add_series("", [total_2024, total_2025])

x, y, cx, cy = Inches(4.5), Inches(1.0), Inches(4.0), Inches(5.0)  # largeur totale du graphique
chart_shape = slide53.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
chart = chart_shape.chart

# Retirer la légende
chart.has_legend = False

# Couleurs : vert clair pour 2024, bleu clair pour 2025
serie_colors = [RGBColor(0, 32, 96), RGBColor(91, 155, 213)]
for s_idx, point in enumerate(chart.plots[0].series[0].points):
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = serie_colors[s_idx]

# Valeurs au-dessus des barres
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.position = XL_LABEL_POSITION.ABOVE
for s_idx, point in enumerate(plot.series[0].points):
    dLbl = point.data_label
    dLbl.has_text_frame = True
    dLbl.text_frame.clear()
    p = dLbl.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{[total_2024, total_2025][s_idx]:d}"
    run.font.size = Pt(12)
    run.font.bold = True

# Axe des valeurs
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.font.size = Pt(10)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(10)
category_axis.tick_labels.font.bold = True

# --- Correction épaisseur anneau ---
try:
    plot.doughnut_hole_size = 70  # valeur si supportée par ta version de python-pptx
except AttributeError:
    pass

# Force via XML si jamais l’attribut n’est pas pris en compte
doughnut_chart = chart.plots[0]._element
hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
if hole_size is None:
    from lxml import etree
    hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
    doughnut_chart.append(hole_size)
hole_size.set('val', '70')  # entre 10 et 90 (40 = anneau plus épais, 80 = plus fin)

### === SLIDE 54 ===
slide54 = prs.slides[53]

# Supprimer ancien graphique s'il existe
for shape in list(slide54.shapes):
    if shape.has_chart:
        slide54.shapes._spTree.remove(shape._element)

# Colonne région
region_col = "ADRESSEFRANCEREGIONLIBELLE"
regions_exclues = ["Territoires d'outre-mer (TOM)"]

# Comptages entreprises ESS
ess_2024 = df_ess_2024[region_col].fillna("").astype(str).str.strip()
ess_2025 = df_ess_2025[region_col].fillna("").astype(str).str.strip()

ess_2024 = ess_2024[ess_2024 != ""].value_counts()
ess_2025 = ess_2025[ess_2025 != ""].value_counts()

# Exclure régions non voulues
ess_2024 = ess_2024[~ess_2024.index.isin(regions_exclues)]
ess_2025 = ess_2025[~ess_2025.index.isin(regions_exclues)]

# Liste des régions (union des deux années)
regions_full = sorted(set(ess_2024.index).union(set(ess_2025.index)))
regions_labels = [r.split(" (")[0] for r in regions_full]

# Valeurs par région
val_2024 = [ess_2024.get(r, 0) for r in regions_full]
val_2025 = [ess_2025.get(r, 0) for r in regions_full]

# --- Création graphique groupé (barres horizontales) ---
chart_data = CategoryChartData()
chart_data.categories = regions_labels
chart_data.add_series("T3 2024", val_2024)
chart_data.add_series("T3 2025", val_2025)

x, y, cx, cy = Inches(0.8), Inches(1.1), Inches(8), Inches(5)
chart = slide54.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

# Retirer la légende
chart.has_legend = False

plot = chart.plots[0]
plot.has_data_labels = True

from pptx.enum.shapes import MSO_SHAPE

# Coordonnées de départ pour la légende (sous le graphique)
legend_y = y + cy + Inches(0.2)
square_size = Inches(0.2)
text_offset = Inches(0.05)  # espace entre carré et texte
item_spacing = Inches(1.0)   # espacement entre items
legend_items = [("T3 2024", RGBColor(0, 32, 96)), ("T3 2025", RGBColor(91, 155, 213))]

# Calcul largeur totale de la légende
total_width = len(legend_items) * (square_size + text_offset + Inches(0.5)) + (len(legend_items)-1)*item_spacing

# Position de départ pour centrer
legend_x = x + (cx - total_width)/2

for i, (label, color) in enumerate(legend_items):
    x_pos = legend_x + i * (square_size + text_offset + Inches(0.5) + item_spacing)

    # Carré coloré
    square = slide54.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x_pos,
        legend_y,
        square_size,
        square_size
    )
    square.fill.solid()
    square.fill.fore_color.rgb = color
    square.line.fill.background()  # pas de contour

    # Texte à côté
    txt_box = slide54.shapes.add_textbox(
        x_pos + square_size + text_offset,
        legend_y,
        Inches(0.5),
        square_size
    )
    txt_frame = txt_box.text_frame
    txt_frame.clear()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.bold = True
    p.alignment = PP_ALIGN.LEFT

# Couleurs : vert clair (2024), bleu clair (2025)
serie_colors = [RGBColor(0, 32, 96), RGBColor(91, 155, 213)]
val_lists = [val_2024, val_2025]

# Appliquer couleurs + labels
for s_idx, series in enumerate(plot.series):
    for p_idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = serie_colors[s_idx]

        # label (nombre d'entreprises)
        dLbl = point.data_label
        dLbl.has_text_frame = True
        dLbl.text_frame.clear()
        p = dLbl.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{val_lists[s_idx][p_idx]:d}"
        run.font.size = Pt(9)
        run.font.bold = True
        dLbl.position = XL_LABEL_POSITION.OUTSIDE_END

# Axe des valeurs
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.width = Pt(0.5)
value_axis.tick_labels.font.size = Pt(9)

# Axe des catégories
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.bold = True

### === SLIDE 55 ===
slide55 = prs.slides[54]

# --- Calcul répartition par forme juridique ---
forme_col = "SECTEURLIBELLE"
df_ess_2025[forme_col] = df_ess_2025[forme_col].astype(str).str.strip()

formes_selectionnees = [
    "ACTIVITÉS IMMOBILIÈRES",
    "COMMERCE , RÉPARATION D'AUTOMOBILES ET DE MOTOCYCLES",
    "ACTIVITÉS SPÉCIALISÉES, SCIENTIFIQUES ET TECHNIQUES",
    "TRANSPORTS ET ENTREPOSAGE",
    "ACTIVITÉS DE SERVICES ADMINISTRATIFS ET DE SOUTIEN",
    "HÉBERGEMENT ET RESTAURATION",
    "CONSTRUCTION"
]

# Comptages des formes juridiques pour 2024 et 2025
formes_2025 = df_ess_2025[forme_col].value_counts()

formes = {}
for f in formes_selectionnees:
    formes[f] = {
        "2025": formes_2025.get(f, 0)
    }

# Catégorie AUTRES
autres_2025 = formes_2025[~formes_2025.index.isin(formes_selectionnees)].sum()
formes["Autres"] = {"2025": autres_2025}

# --- Mise à jour diagramme ---
for shape in list(slide55.shapes):
    if shape.has_chart:
        slide55.shapes._spTree.remove(shape._element)

total_formes_2025 = sum([formes[f]["2025"] for f in formes])
parts_formes_2025 = [round(formes[f]["2025"] / total_formes_2025 * 100, 1) for f in formes]

chart_data = CategoryChartData()
chart_data.categories = list(formes.keys())
chart_data.add_series("", parts_formes_2025)

x, y, cx, cy = Inches(3.9), Inches(0.3), Inches(4.8), Inches(7.3)
chart = slide55.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart

chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

# Modifier la taille de police des étiquettes
data_labels.font.size = Pt(9)
data_labels.font.bold = False

data_labels.show_category_name = True
data_labels.show_value = True   # % affichés
data_labels.show_percentage = False

# --- Correction épaisseur anneau ---
try:
    plot.doughnut_hole_size = 70  # valeur si supportée par ta version de python-pptx
except AttributeError:
    pass

# Force via XML si jamais l’attribut n’est pas pris en compte
doughnut_chart = chart.plots[0]._element
hole_size = doughnut_chart.find('.//c:holeSize', chart._element.nsmap)
if hole_size is None:
    from lxml import etree
    hole_size = etree.Element('{%s}holeSize' % chart._element.nsmap['c'])
    doughnut_chart.append(hole_size)
hole_size.set('val', '70')  # entre 10 et 90 (40 = anneau plus épais, 80 = plus fin)

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace("X", "3").replace("x", "3")
                    run.text = run.text.replace("tribunau3", "tribunaux").replace("TRIBUNAU3", "TRIBUNAUX")
                    run.text = run.text.replace("sociau3", "sociaux").replace("SOCIAU3", "SOCIAUX")
                    run.text = run.text.replace("avril", "juillet")
                    run.text = run.text.replace("juin", "septembre")

#   SAUVEGARDE
output_filename = os.path.join(base_path, "Stats_France_new.pptx")
prs.save(output_filename)
print(f"Fichier PowerPoint modifié enregistré sous : {output_filename}") 
