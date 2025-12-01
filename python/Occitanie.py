# Imports
import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

base_path = r"C:\Users\ext-yaissa\OneDrive - INFOGREFFE\Documents\Stats Occitanie"

files_needed = [
    "Modèle Stats Occitanie.pptx",
    "Immat_2024_t3.xlsx",
    "Immat_2025_t3.xlsx",
    "Pcl_2024_t3.xlsx",
    "Pcl_2025_t3.xlsx",
    "Radiation_2024_t3.xlsx",
    "Radiation_2025_t3.xlsx",
    "StockRCS_2025_t3_au15092025.xlsx",
    "InjPayer_2024_t3.xlsx",
    "InjPayer_2025_t3.xlsx",
    "TrfSiege_2024_t3.xlsx",
    "TrfSiege_2025_t3.xlsx",
    "Dpca_2024_t3.xlsx",
    "Dpca_2025_t3.xlsx",
    "SteESS_2024_t3.xlsx",
    "SteESS_2025_t3.xlsx",
    "SteMission_2024_t3.xlsx",
    "SteMission_2025_t3.xlsx",
    "Mbe_2024_t3.xlsx",
    "Mbe_2025_t3.xlsx",
    "InscrSSTR_2024_t3.xlsx",
    "InscrSSTR_2025_t3.xlsx",
    "Div_2024_t3.xlsx",
    "Div_2025_t3.xlsx",
    "Prev_2024_t3.xlsx",
    "Prev_2025_t3.xlsx",
]

# Vérification des fichiers
missing_files = [f for f in files_needed if not os.path.exists(os.path.join(base_path, f))]
if missing_files:
    print("Fichiers manquants :", missing_files)
    print("Veuillez les ajouter dans le dossier :", base_path)
else:
    print("Tous les fichiers sont présents.")

# === CHARGEMENT DES DONNÉES ===
df_immat_2024 = pd.read_excel(os.path.join(base_path, "Immat_2024_t3.xlsx"), sheet_name="Immat_2024_t3")
df_immat_2025 = pd.read_excel(os.path.join(base_path, "Immat_2025_t3.xlsx"), sheet_name="Immat_2025_t3")

df_pcl_2024 = pd.read_excel(os.path.join(base_path, "Pcl_2024_t3.xlsx"), sheet_name="Pcl_2024_t3")
df_pcl_2025 = pd.read_excel(os.path.join(base_path, "Pcl_2025_t3.xlsx"), sheet_name="Pcl_2025_t3")

df_stock_2025 = pd.read_excel(os.path.join(base_path, "StockRCS_2025_t3_au15092025.xlsx"), sheet_name="StockRCS_2025_t3_au15092025")

df_radiation_2024 = pd.read_excel(os.path.join(base_path, "Radiation_2024_t3.xlsx"), sheet_name="Radiation_2024_t3")
df_radiation_2025 = pd.read_excel(os.path.join(base_path, "Radiation_2025_t3.xlsx"), sheet_name="Radiation_2025_t3")

df_ip_2024 = pd.read_excel(os.path.join(base_path, "InjPayer_2024_t3.xlsx"), sheet_name="InjPayer_2024_t3")
df_ip_2025 = pd.read_excel(os.path.join(base_path, "InjPayer_2025_t3.xlsx"), sheet_name="InjPayer_2025_t3")

df_ts_2024 = pd.read_excel(os.path.join(base_path, "TrfSiege_2024_t3.xlsx"), sheet_name="TrfSiege_2024_t3")
df_ts_2025 = pd.read_excel(os.path.join(base_path, "TrfSiege_2025_t3.xlsx"), sheet_name="TrfSiege_2025_t3")

df_ess_2024 = pd.read_excel(os.path.join(base_path, "SteESS_2024_t3.xlsx"), sheet_name="SteESS_2024_t3")
df_ess_2025 = pd.read_excel(os.path.join(base_path, "SteESS_2025_t3.xlsx"), sheet_name="SteESS_2025_t3")

df_miss_2024 = pd.read_excel(os.path.join(base_path, "SteMission_2024_t3.xlsx"), sheet_name="SteMission_2024_t3")
df_miss_2025 = pd.read_excel(os.path.join(base_path, "SteMission_2025_t3.xlsx"), sheet_name="SteMission_2025_t3")

df_mbe_2024 = pd.read_excel(os.path.join(base_path, "Mbe_2024_t3.xlsx"), sheet_name="Mbe_2024_t3")
df_mbe_2025 = pd.read_excel(os.path.join(base_path, "Mbe_2025_t3.xlsx"), sheet_name="Mbe_2025_t3")

df_dpca_2024 = pd.read_excel(os.path.join(base_path, "Dpca_2024_t3.xlsx"), sheet_name="Dpca_2024_t3")
df_dpca_2025 = pd.read_excel(os.path.join(base_path, "Dpca_2025_t3.xlsx"), sheet_name="Dpca_2025_t3")

df_sstr_2024 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2024_t3.xlsx"), sheet_name="InscrSSTR_2024_t3")
df_sstr_2025 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2025_t3.xlsx"), sheet_name="InscrSSTR_2025_t3")

df_div_2024 = pd.read_excel(os.path.join(base_path, "Div_2024_t3.xlsx"), sheet_name="Div_2024_t3")
df_div_2025 = pd.read_excel(os.path.join(base_path, "Div_2025_t3.xlsx"), sheet_name="Div_2025_t3")

df_prev_2024 = pd.read_excel(os.path.join(base_path, "Prev_2024_t3.xlsx"), sheet_name="Prev_2024_t3")
df_prev_2025 = pd.read_excel(os.path.join(base_path, "Prev_2025_t3.xlsx"), sheet_name="Prev_2025_t3")

# === PRÉ-TRAITEMENT COMMUN ===
def preprocess(df):
    dedup_cols = ["GRF", "MIL", "STC", "CHRONO"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_immat_2024, df_immat_2025 = preprocess(df_immat_2024), preprocess(df_immat_2025)
df_pcl_2024, df_pcl_2025 = preprocess(df_pcl_2024), preprocess(df_pcl_2025)
df_radiation_2024, df_radiation_2025 = preprocess(df_radiation_2024), preprocess(df_radiation_2025)
df_ess_2024, df_ess_2025 = preprocess(df_ess_2024), preprocess(df_ess_2025)
df_miss_2024, df_miss_2025 = preprocess(df_miss_2024), preprocess(df_miss_2025)
df_mbe_2024, df_mbe_2025 = preprocess(df_mbe_2024), preprocess(df_mbe_2025)
df_dpca_2024, df_dpca_2025 = preprocess(df_dpca_2024), preprocess(df_dpca_2025)

def preprocess3(df):
    # Filtrage des départements sans tenir compte de la casse
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_stock_2025 = preprocess3(df_stock_2025)

def preprocess2(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_AFFAIRE"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Filtrage des départements sans tenir compte de la casse
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_ip_2024, df_ip_2025 = preprocess2(df_ip_2024), preprocess2(df_ip_2025)

def preprocess4(df):
    dept_cols = ["DEPFROM", "REGIONTO"]
    for col in dept_cols:
        df[col] = df[col].astype(str).str.strip().str.title()
    exclure_depts = ["Bas Rhin", "Haut Rhin", "Moselle"]
    for col in dept_cols:
        df = df[~df[col].isin(exclure_depts)]
    df = df.reset_index(drop=True)
    return df

df_ts_2024, df_ts_2025 = preprocess4(df_ts_2024), preprocess4(df_ts_2025)

def preprocess9(df):
    # Filtrage des départements sans tenir compte de la casse
    dept_col = "DEPARTEMENT"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_prev_2024, df_prev_2025 = preprocess9(df_prev_2024), preprocess9(df_prev_2025)

# === PRÉ-TRAITEMENT pour SSTR Trésor ===
def preprocess5(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes du Trésor
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_tresor = df["libelle inscription"].str.contains("trésor|tresor", case=False, na=False)
    df = df[mask_tresor].reset_index(drop=True)

    return df


# === PRÉ-TRAITEMENT pour SSTR Sécurité sociale ===
def preprocess7(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes sécurité sociale
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_ss = df["libelle inscription"].str.contains("sécurité sociale|securite sociale|sécurité|securite", case=False, na=False)
    df = df[mask_ss].reset_index(drop=True)

    return df

df_sstr_tresor_2024 = preprocess5(df_sstr_2024)
df_sstr_tresor_2025 = preprocess5(df_sstr_2025)
df_sstr_ss_2024 = preprocess7(df_sstr_2024)
df_sstr_ss_2025 = preprocess7(df_sstr_2025)

# === OUVRIR LE PPT ===
prs = Presentation(os.path.join(base_path, "Modèle Stats Occitanie.pptx"))

# === FILTRAGE ARIEGE ===
df_immat_2024_ge = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_immat_2025_ge = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_pcl_2024_ge = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_pcl_2025_ge = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_radiation_2024_ge = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_radiation_2025_ge = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_stock_2025_ge = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"].reset_index(drop=True)

df_ip_2024_ge = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_ip_2025_ge = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_sstr_tresor_2024_ge = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_sstr_tresor_2025_ge = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_sstr_ss_2024_ge = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_sstr_ss_2025_ge = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_mbe_2024_ge = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_mbe_2025_ge = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

greffes_ge = [
    "FOIX"
]

df_div_2024_ge = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_ge)].reset_index(drop=True)
df_div_2025_ge = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_ge)].reset_index(drop=True)

df_dpca_2024_ge = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]
df_dpca_2025_ge = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "ARIEGE"]

df_ts_2024_ge = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "ARIEGE") |
    (df_ts_2024["DEPTO"] == "ARIEGE")
]

df_ts_2025_ge = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "ARIEGE") |
    (df_ts_2025["DEPTO"] == "ARIEGE")
]

df_prev_2024_ge = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "ARIEGE"]
df_prev_2025_ge = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "ARIEGE"]

# === SLIDE 2 ===
slide2 = prs.slides[1]

A_2024 = df_immat_2024_ge["SIREN"].notna().sum()
A_2025 = df_immat_2025_ge["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_ge["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_ge["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_ge["SIREN"].notna().sum()
C_2025 = df_pcl_2025_ge["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_ge["SIREN"].notna().sum()
D_2025 = df_radiation_2025_ge["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_ge["SIREN"].notna().sum()
E_2025 = df_ip_2025_ge["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_ge["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_ge["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_ge["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_ge["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_ge, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_ge, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_ge, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_ge, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_ge, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_ge, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_ge, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_ge, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_ge, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_ge, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 2
for shape in slide2.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 3 ===
slide3 = prs.slides[2]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide3 = df_mbe_2024_ge["SIREN"].notna().sum()
A_2025_slide3 = df_mbe_2025_ge["SIREN"].notna().sum()
A_val_slide3 = f"{A_2025_slide3:,}".replace(",", " ")
A_pct_slide3 = evol_percent_txt(A_2024_slide3, A_2025_slide3)

# B : DIV
B_2024_slide3 = df_div_2024_ge["SIREN"].notna().sum()
B_2025_slide3 = df_div_2025_ge["SIREN"].notna().sum()
B_val_slide3 = f"{B_2025_slide3:,}".replace(",", " ")
B_pct_slide3 = evol_percent_txt(B_2024_slide3, B_2025_slide3)

# C : DPCA
C_2024_slide3 = df_dpca_2024_ge["SIREN"].notna().sum()
C_2025_slide3 = df_dpca_2025_ge["SIREN"].notna().sum()
C_val_slide3 = f"{C_2025_slide3:,}".replace(",", " ")
C_pct_slide3 = evol_percent_txt(C_2024_slide3, C_2025_slide3)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "ARIEGE").sum()
departs = (df_ts_2025["DEPFROM"] == "ARIEGE").sum()
solde = arrivees - departs

D_val_slide3 = f"{arrivees:,}".replace(",", " ")
E_val_slide3 = f"{departs:,}".replace(",", " ")
F_val_slide3 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 3 ===
for shape in slide3.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide3, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide3, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide3, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide3, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide3, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide3, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide3, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide3, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide3, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT AUDE

# === FILTRAGE AUDE ===
df_immat_2024_cvd = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_immat_2025_cvd = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_pcl_2024_cvd = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_pcl_2025_cvd = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_radiation_2024_cvd = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_radiation_2025_cvd = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_stock_2025_cvd = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"].reset_index(drop=True)

df_ip_2024_cvd = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_ip_2025_cvd = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_sstr_tresor_2024_cvd = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_sstr_tresor_2025_cvd = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_sstr_ss_2024_cvd = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_sstr_ss_2025_cvd = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_mbe_2024_cvd = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_mbe_2025_cvd = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

greffes_cvd = [
    "CARCASSONNE", "NARBONNE"
]

df_div_2024_cvd = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_cvd)].reset_index(drop=True)
df_div_2025_cvd = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_cvd)].reset_index(drop=True)

df_dpca_2024_cvd = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]
df_dpca_2025_cvd = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AUDE"]

df_ts_2024_cvd = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "AUDE") |
    (df_ts_2024["DEPTO"] == "AUDE")
]

df_ts_2025_cvd = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "AUDE") |
    (df_ts_2025["DEPTO"] == "AUDE")
]

df_prev_2024_cvd = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "AUDE"]
df_prev_2025_cvd = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "AUDE"]

# === SLIDE 4 ===
slide4 = prs.slides[3]

A_2024 = df_immat_2024_cvd["SIREN"].notna().sum()
A_2025 = df_immat_2025_cvd["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_cvd["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_cvd["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_cvd["SIREN"].notna().sum()
C_2025 = df_pcl_2025_cvd["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_cvd["SIREN"].notna().sum()
D_2025 = df_radiation_2025_cvd["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_cvd["SIREN"].notna().sum()
E_2025 = df_ip_2025_cvd["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_cvd["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_cvd["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_cvd["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_cvd["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_cvd, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_cvd, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_cvd, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_cvd, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_cvd, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_cvd, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_cvd, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_cvd, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_cvd, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_cvd, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

# Helper: compute percent text consistently
def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 4
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 5 ===
slide5 = prs.slides[4]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide5 = df_mbe_2024_cvd["SIREN"].notna().sum()
A_2025_slide5 = df_mbe_2025_cvd["SIREN"].notna().sum()
A_val_slide5 = f"{A_2025_slide5:,}".replace(",", " ")
A_pct_slide5 = evol_percent_txt(A_2024_slide5, A_2025_slide5)

# B : DIV
B_2024_slide5 = df_div_2024_cvd["SIREN"].notna().sum()
B_2025_slide5 = df_div_2025_cvd["SIREN"].notna().sum()
B_val_slide5 = f"{B_2025_slide5:,}".replace(",", " ")
B_pct_slide5 = evol_percent_txt(B_2024_slide5, B_2025_slide5)

# C : DPCA
C_2024_slide5 = df_dpca_2024_cvd["SIREN"].notna().sum()
C_2025_slide5 = df_dpca_2025_cvd["SIREN"].notna().sum()
C_val_slide5 = f"{C_2025_slide5:,}".replace(",", " ")
C_pct_slide5 = evol_percent_txt(C_2024_slide5, C_2025_slide5)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees1 = (df_ts_2025["DEPTO"] == "AUDE").sum()
departs1 = (df_ts_2025["DEPFROM"] == "AUDE").sum()
solde1 = arrivees1 - departs1

D_val_slide5 = f"{arrivees1:,}".replace(",", " ")
E_val_slide5 = f"{departs1:,}".replace(",", " ")
F_val_slide5 = f"{solde1:,}".replace(",", " ")

# === Remplacement du texte dans la slide 5 ===
for shape in slide5.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide5, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide5, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide5, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide5, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide5, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide5, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide5, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide5, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide5, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT AVEYRON

# === FILTRAGE AVEYRON ===
df_immat_2024_co = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_immat_2025_co = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_pcl_2024_co = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_pcl_2025_co = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_radiation_2024_co = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_radiation_2025_co = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_stock_2025_co = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"].reset_index(drop=True)

df_ip_2024_co = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_ip_2025_co = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_sstr_tresor_2024_co = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_sstr_tresor_2025_co = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_sstr_ss_2024_co = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_sstr_ss_2025_co = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_mbe_2024_co = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_mbe_2025_co = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

greffes_co = [
    "RODEZ"
]

df_div_2024_co = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_co)].reset_index(drop=True)
df_div_2025_co = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_co)].reset_index(drop=True)

df_dpca_2024_co = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]
df_dpca_2025_co = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "AVEYRON"]

df_ts_2024_co = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "AVEYRON") |
    (df_ts_2024["DEPTO"] == "AVEYRON")
]

df_ts_2025_co = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "AVEYRON") |
    (df_ts_2025["DEPTO"] == "AVEYRON")
]

df_prev_2024_co = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "AVEYRON"]
df_prev_2025_co = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "AVEYRON"]

# === SLIDE 6 ===
slide6 = prs.slides[5]

A_2024 = df_immat_2024_co["SIREN"].notna().sum()
A_2025 = df_immat_2025_co["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_co["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_co["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_co["SIREN"].notna().sum()
C_2025 = df_pcl_2025_co["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_co["SIREN"].notna().sum()
D_2025 = df_radiation_2025_co["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_co["SIREN"].notna().sum()
E_2025 = df_ip_2025_co["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_co["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_co["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_co["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_co["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_co, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_co, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_co, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_co, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_co, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_co, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_co, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_co, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_co, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_co, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 6
for shape in slide6.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 7 ===
slide7 = prs.slides[6]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide7 = df_mbe_2024_co["SIREN"].notna().sum()
A_2025_slide7 = df_mbe_2025_co["SIREN"].notna().sum()
A_val_slide7 = f"{A_2025_slide7:,}".replace(",", " ")
A_pct_slide7 = evol_percent_txt(A_2024_slide7, A_2025_slide7)

# B : DIV
B_2024_slide7 = df_div_2024_co["SIREN"].notna().sum()
B_2025_slide7 = df_div_2025_co["SIREN"].notna().sum()
B_val_slide7 = f"{B_2025_slide7:,}".replace(",", " ")
B_pct_slide7 = evol_percent_txt(B_2024_slide7, B_2025_slide7)

# C : DPCA
C_2024_slide7 = df_dpca_2024_co["SIREN"].notna().sum()
C_2025_slide7 = df_dpca_2025_co["SIREN"].notna().sum()
C_val_slide7 = f"{C_2025_slide7:,}".replace(",", " ")
C_pct_slide7 = evol_percent_txt(C_2024_slide7, C_2025_slide7)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "AVEYRON").sum()
departs = (df_ts_2025["DEPFROM"] == "AVEYRON").sum()
solde = arrivees - departs

D_val_slide7 = f"{arrivees:,}".replace(",", " ")
E_val_slide7 = f"{departs:,}".replace(",", " ")
F_val_slide7 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 7 ===
for shape in slide7.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide7, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide7, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide7, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide7, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide7, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide7, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide7, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide7, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide7, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT GARD

# === FILTRAGE GARD ===
df_immat_2024_hdf = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_immat_2025_hdf = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_pcl_2024_hdf = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_pcl_2025_hdf = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_radiation_2024_hdf = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_radiation_2025_hdf = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_stock_2025_hdf = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"].reset_index(drop=True)

df_ip_2024_hdf = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_ip_2025_hdf = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_sstr_tresor_2024_hdf = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_sstr_tresor_2025_hdf = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_sstr_ss_2024_hdf = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_sstr_ss_2025_hdf = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_mbe_2024_hdf = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_mbe_2025_hdf = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

greffes_hdf = [
    "NIMES"
]

df_div_2024_hdf = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_hdf)].reset_index(drop=True)
df_div_2025_hdf = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_hdf)].reset_index(drop=True)

df_dpca_2024_hdf = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]
df_dpca_2025_hdf = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GARD"]

df_ts_2024_hdf = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "GARD") |
    (df_ts_2024["DEPTO"] == "GARD")
]

df_ts_2025_hdf = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "GARD") |
    (df_ts_2025["DEPTO"] == "GARD")
]

df_prev_2024_hdf = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "GARD"]
df_prev_2025_hdf = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "GARD"]

# === SLIDE 8 ===
slide8 = prs.slides[7]

A_2024 = df_immat_2024_hdf["SIREN"].notna().sum()
A_2025 = df_immat_2025_hdf["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_hdf["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_hdf["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_hdf["SIREN"].notna().sum()
C_2025 = df_pcl_2025_hdf["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_hdf["SIREN"].notna().sum()
D_2025 = df_radiation_2025_hdf["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_hdf["SIREN"].notna().sum()
E_2025 = df_ip_2025_hdf["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_hdf["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_hdf["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_hdf["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_hdf["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_hdf, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_hdf, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_hdf, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_hdf, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_hdf, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_hdf, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_hdf, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_hdf, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_hdf, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_hdf, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 8
for shape in slide8.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 9 ===
slide9 = prs.slides[8]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide9 = df_mbe_2024_hdf["SIREN"].notna().sum()
A_2025_slide9 = df_mbe_2025_hdf["SIREN"].notna().sum()
A_val_slide9 = f"{A_2025_slide9:,}".replace(",", " ")
A_pct_slide9 = evol_percent_txt(A_2024_slide9, A_2025_slide9)

# B : DIV
B_2024_slide9 = df_div_2024_hdf["SIREN"].notna().sum()
B_2025_slide9 = df_div_2025_hdf["SIREN"].notna().sum()
B_val_slide9 = f"{B_2025_slide9:,}".replace(",", " ")
B_pct_slide9 = evol_percent_txt(B_2024_slide9, B_2025_slide9)

# C : DPCA
C_2024_slide9 = df_dpca_2024_hdf["SIREN"].notna().sum()
C_2025_slide9 = df_dpca_2025_hdf["SIREN"].notna().sum()
C_val_slide9 = f"{C_2025_slide9:,}".replace(",", " ")
C_pct_slide9 = evol_percent_txt(C_2024_slide9, C_2025_slide9)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "GARD").sum()
departs = (df_ts_2025["DEPFROM"] == "GARD").sum()
solde = arrivees - departs

D_val_slide9 = f"{arrivees:,}".replace(",", " ")
E_val_slide9 = f"{departs:,}".replace(",", " ")
F_val_slide9 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 9 ===
for shape in slide9.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide9, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide9, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide9, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide9, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide9, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide9, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide9, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide9, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide9, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT GERS

# === FILTRAGE GERS ===
df_immat_2024_bfc = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_immat_2025_bfc = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_pcl_2024_bfc = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_pcl_2025_bfc = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_radiation_2024_bfc = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_radiation_2025_bfc = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_stock_2025_bfc = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"].reset_index(drop=True)

df_ip_2024_bfc = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_ip_2025_bfc = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_sstr_tresor_2024_bfc = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_sstr_tresor_2025_bfc = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_sstr_ss_2024_bfc = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_sstr_ss_2025_bfc = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_mbe_2024_bfc = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_mbe_2025_bfc = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

greffes_bfc = [
    "AUCH"
]

df_div_2024_bfc = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bfc)].reset_index(drop=True)
df_div_2025_bfc = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bfc)].reset_index(drop=True)

df_dpca_2024_bfc = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]
df_dpca_2025_bfc = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "GERS"]

df_ts_2024_bfc = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "GERS") |
    (df_ts_2024["DEPTO"] == "GERS")
]

df_ts_2025_bfc = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "GERS") |
    (df_ts_2025["DEPTO"] == "GERS")
]

df_prev_2024_bfc = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "GERS"]
df_prev_2025_bfc = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "GERS"]

# === SLIDE 10 ===
slide10 = prs.slides[9]

A_2024 = df_immat_2024_bfc["SIREN"].notna().sum()
A_2025 = df_immat_2025_bfc["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bfc["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bfc["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bfc["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bfc["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bfc["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bfc["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bfc["SIREN"].notna().sum()
E_2025 = df_ip_2025_bfc["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bfc["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bfc["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bfc["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bfc["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bfc, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bfc, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bfc, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bfc, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bfc, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bfc, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bfc, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bfc, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bfc, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bfc, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 10
for shape in slide10.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 11 ===
slide11 = prs.slides[10]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide11 = df_mbe_2024_bfc["SIREN"].notna().sum()
A_2025_slide11 = df_mbe_2025_bfc["SIREN"].notna().sum()
A_val_slide11 = f"{A_2025_slide11:,}".replace(",", " ")
A_pct_slide11 = evol_percent_txt(A_2024_slide11, A_2025_slide11)

# B : DIV
B_2024_slide11 = df_div_2024_bfc["SIREN"].notna().sum()
B_2025_slide11 = df_div_2025_bfc["SIREN"].notna().sum()
B_val_slide11 = f"{B_2025_slide11:,}".replace(",", " ")
B_pct_slide11 = evol_percent_txt(B_2024_slide11, B_2025_slide11)

# C : DPCA
C_2024_slide11 = df_dpca_2024_bfc["SIREN"].notna().sum()
C_2025_slide11 = df_dpca_2025_bfc["SIREN"].notna().sum()
C_val_slide11 = f"{C_2025_slide11:,}".replace(",", " ")
C_pct_slide11 = evol_percent_txt(C_2024_slide11, C_2025_slide11)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "GERS").sum()
departs = (df_ts_2025["DEPFROM"] == "GERS").sum()
solde = arrivees - departs

D_val_slide11 = f"{arrivees:,}".replace(",", " ")
E_val_slide11 = f"{departs:,}".replace(",", " ")
F_val_slide11 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 11 ===
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide11, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide11, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide11, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide11, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide11, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide11, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide11, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide11, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide11, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT HAUTE GARONNE

# === FILTRAGE HAUTE GARONNE ===
df_immat_2024_aura = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_immat_2025_aura = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_pcl_2024_aura = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_pcl_2025_aura = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_radiation_2024_aura = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_radiation_2025_aura = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_stock_2025_aura = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"].reset_index(drop=True)

df_ip_2024_aura = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_ip_2025_aura = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_sstr_tresor_2024_aura = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_sstr_tresor_2025_aura = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_sstr_ss_2024_aura = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_sstr_ss_2025_aura = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_mbe_2024_aura = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_mbe_2025_aura = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

greffes_aura = [
    "TOULOUSE"
]

df_div_2024_aura = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_aura)].reset_index(drop=True)
df_div_2025_aura = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_aura)].reset_index(drop=True)

df_dpca_2024_aura = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]
df_dpca_2025_aura = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTE GARONNE"]

df_ts_2024_aura = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "HAUTE GARONNE") |
    (df_ts_2024["DEPTO"] == "HAUTE GARONNE")
]

df_ts_2025_aura = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "HAUTE GARONNE") |
    (df_ts_2025["DEPTO"] == "HAUTE GARONNE")
]

df_prev_2024_aura = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "HAUTE GARONNE"]
df_prev_2025_aura = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "HAUTE GARONNE"]

# === SLIDE 12 ===
slide12 = prs.slides[11]

A_2024 = df_immat_2024_aura["SIREN"].notna().sum()
A_2025 = df_immat_2025_aura["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_aura["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_aura["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_aura["SIREN"].notna().sum()
C_2025 = df_pcl_2025_aura["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_aura["SIREN"].notna().sum()
D_2025 = df_radiation_2025_aura["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_aura["SIREN"].notna().sum()
E_2025 = df_ip_2025_aura["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_aura["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_aura["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_aura["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_aura["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_aura, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_aura, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_aura, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_aura, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_aura, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_aura, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_aura, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_aura, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_aura, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_aura, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 12
for shape in slide12.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 13 ===
slide13 = prs.slides[12]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide13 = df_mbe_2024_aura["SIREN"].notna().sum()
A_2025_slide13 = df_mbe_2025_aura["SIREN"].notna().sum()
A_val_slide13 = f"{A_2025_slide13:,}".replace(",", " ")
A_pct_slide13 = evol_percent_txt(A_2024_slide13, A_2025_slide13)

# B : DIV
B_2024_slide13 = df_div_2024_aura["SIREN"].notna().sum()
B_2025_slide13 = df_div_2025_aura["SIREN"].notna().sum()
B_val_slide13 = f"{B_2025_slide13:,}".replace(",", " ")
B_pct_slide13 = evol_percent_txt(B_2024_slide13, B_2025_slide13)

# C : DPCA
C_2024_slide13 = df_dpca_2024_aura["SIREN"].notna().sum()
C_2025_slide13 = df_dpca_2025_aura["SIREN"].notna().sum()
C_val_slide13 = f"{C_2025_slide13:,}".replace(",", " ")
C_pct_slide13 = evol_percent_txt(C_2024_slide13, C_2025_slide13)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "HAUTE GARONNE").sum()
departs = (df_ts_2025["DEPFROM"] == "HAUTE GARONNE").sum()
solde = arrivees - departs

D_val_slide13 = f"{arrivees:,}".replace(",", " ")
E_val_slide13 = f"{departs:,}".replace(",", " ")
F_val_slide13 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 13 ===
for shape in slide13.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide13, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide13, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide13, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide13, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide13, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide13, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide13, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide13, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide13, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT HAUTES PYRENEES

# === FILTRAGE HAUTES PYRENEES ===
df_immat_2024_idf = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_immat_2025_idf = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_pcl_2024_idf = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_pcl_2025_idf = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_radiation_2024_idf = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_radiation_2025_idf = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_stock_2025_idf = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"].reset_index(drop=True)

df_ip_2024_idf = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_ip_2025_idf = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_sstr_tresor_2024_idf = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_sstr_tresor_2025_idf = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_sstr_ss_2024_idf = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_sstr_ss_2025_idf = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_mbe_2024_idf = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_mbe_2025_idf = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

greffes_idf = [
    "TARBES"
]

df_div_2024_idf = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_idf)].reset_index(drop=True)
df_div_2025_idf = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_idf)].reset_index(drop=True)

df_dpca_2024_idf = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]
df_dpca_2025_idf = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HAUTES PYRENEES"]

df_ts_2024_idf = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "HAUTES PYRENEES") |
    (df_ts_2024["DEPTO"] == "HAUTES PYRENEES")
]

df_ts_2025_idf = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "HAUTES PYRENEES") |
    (df_ts_2025["DEPTO"] == "HAUTES PYRENEES")
]

df_prev_2024_idf = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "HAUTES PYRENEES"]
df_prev_2025_idf = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "HAUTES PYRENEES"]

# === SLIDE 14 ===
slide14 = prs.slides[13]

A_2024 = df_immat_2024_idf["SIREN"].notna().sum()
A_2025 = df_immat_2025_idf["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_idf["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_idf["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_idf["SIREN"].notna().sum()
C_2025 = df_pcl_2025_idf["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_idf["SIREN"].notna().sum()
D_2025 = df_radiation_2025_idf["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_idf["SIREN"].notna().sum()
E_2025 = df_ip_2025_idf["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_idf["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_idf["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_idf["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_idf["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_idf, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_idf, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_idf, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_idf, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_idf, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_idf, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_idf, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_idf, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_idf, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_idf, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 14
for shape in slide14.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 15 ===
slide15 = prs.slides[14]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide15 = df_mbe_2024_idf["SIREN"].notna().sum()
A_2025_slide15 = df_mbe_2025_idf["SIREN"].notna().sum()
A_val_slide15 = f"{A_2025_slide15:,}".replace(",", " ")
A_pct_slide15 = evol_percent_txt(A_2024_slide15, A_2025_slide15)

# B : DIV
B_2024_slide15 = df_div_2024_idf["SIREN"].notna().sum()
B_2025_slide15 = df_div_2025_idf["SIREN"].notna().sum()
B_val_slide15 = f"{B_2025_slide15:,}".replace(",", " ")
B_pct_slide15 = evol_percent_txt(B_2024_slide15, B_2025_slide15)

# C : DPCA
C_2024_slide15 = df_dpca_2024_idf["SIREN"].notna().sum()
C_2025_slide15 = df_dpca_2025_idf["SIREN"].notna().sum()
C_val_slide15 = f"{C_2025_slide15:,}".replace(",", " ")
C_pct_slide15 = evol_percent_txt(C_2024_slide15, C_2025_slide15)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "HAUTES PYRENEES").sum()
departs = (df_ts_2025["DEPFROM"] == "HAUTES PYRENEES").sum()
solde = arrivees - departs

D_val_slide15 = f"{arrivees:,}".replace(",", " ")
E_val_slide15 = f"{departs:,}".replace(",", " ")
F_val_slide15 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 15 ===
for shape in slide15.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide15, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide15, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide15, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide15, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide15, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide15, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide15, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide15, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide15, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT HERAULT

# === FILTRAGE HERAULT===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

greffes_bre = [
    "BEZIERS", "MONTPELLIER"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "HERAULT"]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "HERAULT") |
    (df_ts_2024["DEPTO"] == "HERAULT")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "HERAULT") |
    (df_ts_2025["DEPTO"] == "HERAULT")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "HERAULT"]
df_prev_2025_bre = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "HERAULT"]

# === SLIDE 16 ===
slide16 = prs.slides[15]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 16
for shape in slide16.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 17 ===
slide17 = prs.slides[16]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide17 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide17 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide17 = f"{A_2025_slide17:,}".replace(",", " ")
A_pct_slide17 = evol_percent_txt(A_2024_slide17, A_2025_slide17)

# B : DIV
B_2024_slide17 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide17 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide17 = f"{B_2025_slide17:,}".replace(",", " ")
B_pct_slide17 = evol_percent_txt(B_2024_slide17, B_2025_slide17)

# C : DPCA
C_2024_slide17 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide17 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide17 = f"{C_2025_slide17:,}".replace(",", " ")
C_pct_slide17 = evol_percent_txt(C_2024_slide17, C_2025_slide17)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "HERAULT").sum()
departs = (df_ts_2025["DEPFROM"] == "HERAULT").sum()
solde = arrivees - departs

D_val_slide17 = f"{arrivees:,}".replace(",", " ")
E_val_slide17 = f"{departs:,}".replace(",", " ")
F_val_slide17 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 17 ===
for shape in slide17.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide17, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide17, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide17, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide17, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide17, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide17, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide17, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide17, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide17, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT LOT

# === FILTRAGE LOT===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

greffes_bre = [
    "CAHORS"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOT"]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "LOT") |
    (df_ts_2024["DEPTO"] == "LOT")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "LOT") |
    (df_ts_2025["DEPTO"] == "LOT")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "LOT"]
df_prev_2025_bre = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "LOT"]

# === SLIDE 18 ===
slide18 = prs.slides[17]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 18
for shape in slide18.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 19 ===
slide19 = prs.slides[18]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide19 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide19 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide19 = f"{A_2025_slide19:,}".replace(",", " ")
A_pct_slide19 = evol_percent_txt(A_2024_slide19, A_2025_slide19)

# B : DIV
B_2024_slide19 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide19 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide19 = f"{B_2025_slide19:,}".replace(",", " ")
B_pct_slide19 = evol_percent_txt(B_2024_slide19, B_2025_slide19)

# C : DPCA
C_2024_slide19 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide19 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide19 = f"{C_2025_slide19:,}".replace(",", " ")
C_pct_slide19 = evol_percent_txt(C_2024_slide19, C_2025_slide19)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "LOT").sum()
departs = (df_ts_2025["DEPFROM"] == "LOT").sum()
solde = arrivees - departs

D_val_slide19 = f"{arrivees:,}".replace(",", " ")
E_val_slide19 = f"{departs:,}".replace(",", " ")
F_val_slide19 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 19 ===
for shape in slide19.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide19, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide19, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide19, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide19, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide19, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide19, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide19, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide19, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide19, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT LOZERE

# === FILTRAGE LOZERE===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

greffes_bre = [
    "MENDE"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "LOZERE"]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "LOZERE") |
    (df_ts_2024["DEPTO"] == "LOZERE")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "LOZERE") |
    (df_ts_2025["DEPTO"] == "LOZERE")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "LOZERE"]
df_prev_2025_bre = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "LOZERE"]

# === SLIDE 20 ===
slide20 = prs.slides[19]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 20
for shape in slide20.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 21 ===
slide21 = prs.slides[20]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide21 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide21 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide21 = f"{A_2025_slide21:,}".replace(",", " ")
A_pct_slide21 = evol_percent_txt(A_2024_slide21, A_2025_slide21)

# B : DIV
B_2024_slide21 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide21 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide21 = f"{B_2025_slide21:,}".replace(",", " ")
B_pct_slide21 = evol_percent_txt(B_2024_slide21, B_2025_slide21)

# C : DPCA
C_2024_slide21 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide21 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide21 = f"{C_2025_slide21:,}".replace(",", " ")
C_pct_slide21 = evol_percent_txt(C_2024_slide21, C_2025_slide21)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "LOZERE").sum()
departs = (df_ts_2025["DEPFROM"] == "LOZERE").sum()
solde = arrivees - departs

D_val_slide21 = f"{arrivees:,}".replace(",", " ")
E_val_slide21 = f"{departs:,}".replace(",", " ")
F_val_slide21 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 21 ===
for shape in slide21.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide21, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide21, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide21, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide21, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide21, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide21, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide21, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide21, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide21, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT PYRENEES ORIENTALES

# === FILTRAGE PYRENEES ORIENTALES===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

greffes_bre = [
    "PERPIGNAN"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "PYRENEES ORIENTALES"]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "PYRENEES ORIENTALES") |
    (df_ts_2024["DEPTO"] == "PYRENEES ORIENTALES")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "PYRENEES ORIENTALES") |
    (df_ts_2025["DEPTO"] == "PYRENEES ORIENTALES")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "PYRENEES ORIENTALES"]
df_prev_2025_bre = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "PYRENEES ORIENTALES"]

# === SLIDE 22 ===
slide22 = prs.slides[21]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 22
for shape in slide22.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 23 ===
slide23 = prs.slides[22]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide23 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide23 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide23 = f"{A_2025_slide23:,}".replace(",", " ")
A_pct_slide23 = evol_percent_txt(A_2024_slide23, A_2025_slide23)

# B : DIV
B_2024_slide23 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide23 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide23 = f"{B_2025_slide23:,}".replace(",", " ")
B_pct_slide23 = evol_percent_txt(B_2024_slide23, B_2025_slide23)

# C : DPCA
C_2024_slide23 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide23 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide23 = f"{C_2025_slide23:,}".replace(",", " ")
C_pct_slide23 = evol_percent_txt(C_2024_slide23, C_2025_slide23)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "PYRENEES ORIENTALES").sum()
departs = (df_ts_2025["DEPFROM"] == "PYRENEES ORIENTALES").sum()
solde = arrivees - departs

D_val_slide23 = f"{arrivees:,}".replace(",", " ")
E_val_slide23 = f"{departs:,}".replace(",", " ")
F_val_slide23 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 23 ===
for shape in slide23.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide23, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide23, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide23, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide23, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide23, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide23, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide23, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide23, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide23, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# DEPARTEMENT TARN

# === FILTRAGE TARN===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

greffes_bre = [
    "ALBI", "CASTRES"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN"]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "TARN") |
    (df_ts_2024["DEPTO"] == "TARN")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "TARN") |
    (df_ts_2025["DEPTO"] == "TARN")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "TARN"]
df_prev_2025_bre = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "TARN"]

# === SLIDE 24 ===
slide24 = prs.slides[23]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 24
for shape in slide24.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 25 ===
slide25 = prs.slides[24]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide25 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide25 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide25 = f"{A_2025_slide25:,}".replace(",", " ")
A_pct_slide25 = evol_percent_txt(A_2024_slide25, A_2025_slide25)

# B : DIV
B_2024_slide25 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide25 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide25 = f"{B_2025_slide25:,}".replace(",", " ")
B_pct_slide25 = evol_percent_txt(B_2024_slide25, B_2025_slide25)

# C : DPCA
C_2024_slide25 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide25 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide25 = f"{C_2025_slide25:,}".replace(",", " ")
C_pct_slide25 = evol_percent_txt(C_2024_slide25, C_2025_slide25)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "TARN").sum()
departs = (df_ts_2025["DEPFROM"] == "TARN").sum()
solde = arrivees - departs

D_val_slide25 = f"{arrivees:,}".replace(",", " ")
E_val_slide25 = f"{departs:,}".replace(",", " ")
F_val_slide25 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 25 ===
for shape in slide25.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide25, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide25, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide25, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide25, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide25, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide25, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide25, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide25, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide25, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER


# DEPARTEMENT TARN ET GARONNE

# === FILTRAGE TARN ET GARONNE===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

greffes_bre = [
    "MONTAUBAN"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEDEPARTEMENTLIB"] == "TARN ET GARONNE"]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["DEPFROM"] == "TARN ET GARONNE") |
    (df_ts_2024["DEPTO"] == "TARN ET GARONNE")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["DEPFROM"] == "TARN ET GARONNE") |
    (df_ts_2025["DEPTO"] == "TARN ET GARONNE")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["DEPARTEMENT"] == "TARN ET GARONNE"]
df_prev_2025_bre = df_prev_2025[df_prev_2025["DEPARTEMENT"] == "TARN ET GARONNE"]

# === SLIDE 26 ===
slide26 = prs.slides[25]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 26
for shape in slide26.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 27 ===
slide27 = prs.slides[26]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide27 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide27 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide27 = f"{A_2025_slide27:,}".replace(",", " ")
A_pct_slide27 = evol_percent_txt(A_2024_slide27, A_2025_slide27)

# B : DIV
B_2024_slide27 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide27 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide27 = f"{B_2025_slide27:,}".replace(",", " ")
B_pct_slide27 = evol_percent_txt(B_2024_slide27, B_2025_slide27)

# C : DPCA
C_2024_slide27 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide27 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide27 = f"{C_2025_slide27:,}".replace(",", " ")
C_pct_slide27 = evol_percent_txt(C_2024_slide27, C_2025_slide27)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["DEPTO"] == "TARN ET GARONNE").sum()
departs = (df_ts_2025["DEPFROM"] == "TARN ET GARONNE").sum()
solde = arrivees - departs

D_val_slide27 = f"{arrivees:,}".replace(",", " ")
E_val_slide27 = f"{departs:,}".replace(",", " ")
F_val_slide27 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 27 ===
for shape in slide27.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide27, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide27, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide27, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide27, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide27, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide27, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide27, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide27, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide27, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace("X", "3").replace("x", "3")
                    run.text = run.text.replace("tribunau3", "tribunaux").replace("TRIBUNAU3", "TRIBUNAUX")
                    run.text = run.text.replace("sociau3", "sociaux").replace("SOCIAU3", "SOCIAUX")
                    run.text = run.text.replace("avril", "juillet")
                    run.text = run.text.replace("juin", "septembre")

#   SAUVEGARDE                   
output_filename = os.path.join(base_path, "Stats_Occitanie_new.pptx")
prs.save(output_filename)
print(f"Fichier PowerPoint modifié enregistré sous : {output_filename}") 
