# Imports
import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

base_path = r"C:\Users\ext-yaissa\OneDrive - INFOGREFFE\Documents\Régional"

files_needed = [
    "Modèle Stats T2 - Copie.pptx",
    "Immat_2024_t3.xlsx",
    "Immat_2025_t3.xlsx",
    "Pcl_2024_t3.xlsx",
    "Pcl_2025_t3.xlsx",
    "Radiation_2024_t3.xlsx",
    "Radiation_2025_t3.xlsx",
    "StockRCS_2025_t3_au15092025.xlsx",
    "InjPayer_2024_t3.xlsx",
    "InjPayer_2025_t3.xlsx",
    "TrfSiege_2024_t3.xlsx",
    "TrfSiege_2025_t3.xlsx",
    "Dpca_2024_t3.xlsx",
    "Dpca_2025_t3.xlsx",
    "Mbe_2024_t3.xlsx",
    "Mbe_2025_t3.xlsx",
    "InscrSSTR_2024_t3.xlsx",
    "InscrSSTR_2025_t3.xlsx",
    "Div_2024_t3.xlsx",
    "Div_2025_t3.xlsx",
    "Prev_2024_t3.xlsx",
    "Prev_2025_t3.xlsx",
]

# Vérification des fichiers
missing_files = [f for f in files_needed if not os.path.exists(os.path.join(base_path, f))]
if missing_files:
    print("Fichiers manquants :", missing_files)
    print("Veuillez les ajouter dans le dossier :", base_path)
else:
    print("Tous les fichiers sont présents.")

# === CHARGEMENT DES DONNÉES ===
df_immat_2024 = pd.read_excel(os.path.join(base_path, "Immat_2024_t3.xlsx"), sheet_name="Immat_2024_t3")
df_immat_2025 = pd.read_excel(os.path.join(base_path, "Immat_2025_t3.xlsx"), sheet_name="Immat_2025_t3")

df_pcl_2024 = pd.read_excel(os.path.join(base_path, "Pcl_2024_t3.xlsx"), sheet_name="Pcl_2024_t3")
df_pcl_2025 = pd.read_excel(os.path.join(base_path, "Pcl_2025_t3.xlsx"), sheet_name="Pcl_2025_t3")

df_stock_2025 = pd.read_excel(os.path.join(base_path, "StockRCS_2025_t3_au15092025.xlsx"), sheet_name="StockRCS_2025_t3_au15092025")

df_radiation_2024 = pd.read_excel(os.path.join(base_path, "Radiation_2024_t3.xlsx"), sheet_name="Radiation_2024_t3")
df_radiation_2025 = pd.read_excel(os.path.join(base_path, "Radiation_2025_t3.xlsx"), sheet_name="Radiation_2025_t3")

df_ip_2024 = pd.read_excel(os.path.join(base_path, "InjPayer_2024_t3.xlsx"), sheet_name="InjPayer_2024_t3")
df_ip_2025 = pd.read_excel(os.path.join(base_path, "InjPayer_2025_t3.xlsx"), sheet_name="InjPayer_2025_t3")

df_ts_2024 = pd.read_excel(os.path.join(base_path, "TrfSiege_2024_t3.xlsx"), sheet_name="TrfSiege_2024_t3")
df_ts_2025 = pd.read_excel(os.path.join(base_path, "TrfSiege_2025_t3.xlsx"), sheet_name="TrfSiege_2025_t3")

df_mbe_2024 = pd.read_excel(os.path.join(base_path, "Mbe_2024_t3.xlsx"), sheet_name="Mbe_2024_t3")
df_mbe_2025 = pd.read_excel(os.path.join(base_path, "Mbe_2025_t3.xlsx"), sheet_name="Mbe_2025_t3")

df_dpca_2024 = pd.read_excel(os.path.join(base_path, "Dpca_2024_t3.xlsx"), sheet_name="Dpca_2024_t3")
df_dpca_2025 = pd.read_excel(os.path.join(base_path, "Dpca_2025_t3.xlsx"), sheet_name="Dpca_2025_t3")

df_sstr_2024 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2024_t3.xlsx"), sheet_name="InscrSSTR_2024_t3")
df_sstr_2025 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2025_t3.xlsx"), sheet_name="InscrSSTR_2025_t3")

df_div_2024 = pd.read_excel(os.path.join(base_path, "Div_2024_t3.xlsx"), sheet_name="Div_2024_t3")
df_div_2025 = pd.read_excel(os.path.join(base_path, "Div_2025_t3.xlsx"), sheet_name="Div_2025_t3")

df_prev_2024 = pd.read_excel(os.path.join(base_path, "Prev_2024_t3.xlsx"), sheet_name="Prev_2024_t3")
df_prev_2025 = pd.read_excel(os.path.join(base_path, "Prev_2025_t3.xlsx"), sheet_name="Prev_2025_t3")

# === PRÉ-TRAITEMENT COMMUN ===
def preprocess(df):
    dedup_cols = ["GRF", "MIL", "STC", "CHRONO"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_immat_2024, df_immat_2025 = preprocess(df_immat_2024), preprocess(df_immat_2025)
df_pcl_2024, df_pcl_2025 = preprocess(df_pcl_2024), preprocess(df_pcl_2025)
df_radiation_2024, df_radiation_2025 = preprocess(df_radiation_2024), preprocess(df_radiation_2025)
df_mbe_2024, df_mbe_2025 = preprocess(df_mbe_2024), preprocess(df_mbe_2025)
df_dpca_2024, df_dpca_2025 = preprocess(df_dpca_2024), preprocess(df_dpca_2025)

def preprocess3(df):
    # Filtrage des départements sans tenir compte de la casse
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_stock_2025 = preprocess3(df_stock_2025)

def preprocess2(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_AFFAIRE"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Filtrage des départements sans tenir compte de la casse
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_ip_2024, df_ip_2025 = preprocess2(df_ip_2024), preprocess2(df_ip_2025)

def preprocess4(df):
    dept_cols = ["DEPFROM", "REGIONTO"]
    for col in dept_cols:
        df[col] = df[col].astype(str).str.strip().str.title()
    exclure_depts = ["Bas Rhin", "Haut Rhin", "Moselle"]
    for col in dept_cols:
        df = df[~df[col].isin(exclure_depts)]
    df = df.reset_index(drop=True)
    return df

df_ts_2024, df_ts_2025 = preprocess4(df_ts_2024), preprocess4(df_ts_2025)

# === PRÉ-TRAITEMENT COMMUN ===
def preprocess9(df):
    # Filtrage des départements sans tenir compte de la casse
    dept_col = "DEPARTEMENT"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_prev_2024, df_prev_2025 = preprocess9(df_prev_2024), preprocess9(df_prev_2025)

# === PRÉ-TRAITEMENT pour SSTR Trésor ===
def preprocess5(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes du Trésor
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_tresor = df["libelle inscription"].str.contains("trésor|tresor", case=False, na=False)
    df = df[mask_tresor].reset_index(drop=True)

    return df


# === PRÉ-TRAITEMENT pour SSTR Sécurité sociale ===
def preprocess7(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes sécurité sociale
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_ss = df["libelle inscription"].str.contains("sécurité sociale|securite sociale|sécurité|securite", case=False, na=False)
    df = df[mask_ss].reset_index(drop=True)

    return df

df_sstr_tresor_2024 = preprocess5(df_sstr_2024)
df_sstr_tresor_2025 = preprocess5(df_sstr_2025)
df_sstr_ss_2024 = preprocess7(df_sstr_2024)
df_sstr_ss_2025 = preprocess7(df_sstr_2025)


# === OUVRIR LE PPT ===
prs = Presentation(os.path.join(base_path, "Modèle Stats T2 - Copie.pptx"))

# === FILTRAGE Grand Est (Alsace-Champagne-Ardenne-Lorraine) ===
df_immat_2024_ge = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_immat_2025_ge = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_pcl_2024_ge = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_pcl_2025_ge = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_radiation_2024_ge = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_radiation_2025_ge = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_stock_2025_ge = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"].reset_index(drop=True)

df_ip_2024_ge = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_ip_2025_ge = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_sstr_tresor_2024_ge = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_sstr_tresor_2025_ge = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_sstr_ss_2024_ge = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_sstr_ss_2025_ge = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_mbe_2024_ge = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_mbe_2025_ge = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

greffes_ge = [
    "STRASBOURG", "REIMS", "METZ", "NANCY", "COLMAR", "EPINAL", "CHALONS-EN-CHAMPAGNE", "SAVERNE", "TROYES", "THIONVILLE", "SENNE", "BESANCON"
]

df_div_2024_ge = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_ge)].reset_index(drop=True)
df_div_2025_ge = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_ge)].reset_index(drop=True)

df_dpca_2024_ge = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_dpca_2025_ge = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

df_ts_2024_ge = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)") |
    (df_ts_2024["REGIONTO"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)")
]

df_ts_2025_ge = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)") |
    (df_ts_2025["REGIONTO"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)")
]

df_prev_2024_ge = df_prev_2024[df_prev_2024["REGION"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]
df_prev_2025_ge = df_prev_2025[df_prev_2025["REGION"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)"]

# === SLIDE 2 ===
slide2 = prs.slides[1]

A_2024 = df_immat_2024_ge["SIREN"].notna().sum()
A_2025 = df_immat_2025_ge["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_ge["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_ge["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_ge["SIREN"].notna().sum()
C_2025 = df_pcl_2025_ge["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_ge["SIREN"].notna().sum()
D_2025 = df_radiation_2025_ge["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_ge["SIREN"].notna().sum()
E_2025 = df_ip_2025_ge["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_ge["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_ge["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_ge["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_ge["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_ge, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_ge, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_ge, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_ge, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_ge, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_ge, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_ge, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_ge, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_ge, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_ge, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 2
for shape in slide2.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 3 ===
slide3 = prs.slides[2]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide3 = df_mbe_2024_ge["SIREN"].notna().sum()
A_2025_slide3 = df_mbe_2025_ge["SIREN"].notna().sum()
A_val_slide3 = f"{A_2025_slide3:,}".replace(",", " ")
A_pct_slide3 = evol_percent_txt(A_2024_slide3, A_2025_slide3)

# B : DIV
B_2024_slide3 = df_div_2024_ge["SIREN"].notna().sum()
B_2025_slide3 = df_div_2025_ge["SIREN"].notna().sum()
B_val_slide3 = f"{B_2025_slide3:,}".replace(",", " ")
B_pct_slide3 = evol_percent_txt(B_2024_slide3, B_2025_slide3)

# C : DPCA
C_2024_slide3 = df_dpca_2024_ge["SIREN"].notna().sum()
C_2025_slide3 = df_dpca_2025_ge["SIREN"].notna().sum()
C_val_slide3 = f"{C_2025_slide3:,}".replace(",", " ")
C_pct_slide3 = evol_percent_txt(C_2024_slide3, C_2025_slide3)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)").sum()
departs = (df_ts_2025["REGIONFROM"] == "Grand Est (Alsace-Champagne-Ardenne-Lorraine)").sum()
solde = arrivees - departs

D_val_slide3 = f"{arrivees:,}".replace(",", " ")
E_val_slide3 = f"{departs:,}".replace(",", " ")
F_val_slide3 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 3 ===
for shape in slide3.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide3, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide3, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide3, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide3, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide3, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide3, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide3, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide3, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide3, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION CENTRE VAL DE LOIRE

# === FILTRAGE Centre - Val de Loire ===
df_immat_2024_cvd = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_immat_2025_cvd = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_pcl_2024_cvd = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_pcl_2025_cvd = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_radiation_2024_cvd = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_radiation_2025_cvd = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_stock_2025_cvd = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"].reset_index(drop=True)

df_ip_2024_cvd = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_ip_2025_cvd = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_sstr_tresor_2024_cvd = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_sstr_tresor_2025_cvd = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_sstr_ss_2024_cvd = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_sstr_ss_2025_cvd = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_mbe_2024_cvd = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_mbe_2025_cvd = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

greffes_cvd = [
    "ORLEANS", "BLOIS", "TOURS", "CHARTRES", "CHÂTEAUROUX", "NEVERS", "BOURGES", "VESOUL"
]

df_div_2024_cvd = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_cvd)].reset_index(drop=True)
df_div_2025_cvd = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_cvd)].reset_index(drop=True)

df_dpca_2024_cvd = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]
df_dpca_2025_cvd = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Centre - Val de Loire"]

df_ts_2024_cvd = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Centre - Val de Loire") |
    (df_ts_2024["REGIONTO"] == "Centre - Val de Loire")
]

df_ts_2025_cvd = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Centre - Val de Loire") |
    (df_ts_2025["REGIONTO"] == "Centre - Val de Loire")
]

df_prev_2024_cvd = df_prev_2024[df_prev_2024["REGION"] == "Centre - Val de Loire"]
df_prev_2025_cvd = df_prev_2025[df_prev_2025["REGION"] == "Centre - Val de Loire"]

# === SLIDE 4 ===
slide4 = prs.slides[3]

# Pour être sûr, recompute A..G,W from your DataFrames (copie directe de ta logique)
A_2024 = df_immat_2024_cvd["SIREN"].notna().sum()
A_2025 = df_immat_2025_cvd["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_cvd["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_cvd["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_cvd["SIREN"].notna().sum()
C_2025 = df_pcl_2025_cvd["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_cvd["SIREN"].notna().sum()
D_2025 = df_radiation_2025_cvd["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_cvd["SIREN"].notna().sum()
E_2025 = df_ip_2025_cvd["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_cvd["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_cvd["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_cvd["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_cvd["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_cvd, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_cvd, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_cvd, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_cvd, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_cvd, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_cvd, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_cvd, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_cvd, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_cvd, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_cvd, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 4
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 5 ===
slide5 = prs.slides[4]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide5 = df_mbe_2024_cvd["SIREN"].notna().sum()
A_2025_slide5 = df_mbe_2025_cvd["SIREN"].notna().sum()
A_val_slide5 = f"{A_2025_slide5:,}".replace(",", " ")
A_pct_slide5 = evol_percent_txt(A_2024_slide5, A_2025_slide5)

# B : DIV
B_2024_slide5 = df_div_2024_cvd["SIREN"].notna().sum()
B_2025_slide5 = df_div_2025_cvd["SIREN"].notna().sum()
B_val_slide5 = f"{B_2025_slide5:,}".replace(",", " ")
B_pct_slide5 = evol_percent_txt(B_2024_slide5, B_2025_slide5)

# C : DPCA
C_2024_slide5 = df_dpca_2024_cvd["SIREN"].notna().sum()
C_2025_slide5 = df_dpca_2025_cvd["SIREN"].notna().sum()
C_val_slide5 = f"{C_2025_slide5:,}".replace(",", " ")
C_pct_slide5 = evol_percent_txt(C_2024_slide5, C_2025_slide5)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees1 = (df_ts_2025["REGIONTO"] == "Centre - Val de Loire").sum()
departs1 = (df_ts_2025["REGIONFROM"] == "Centre - Val de Loire").sum()
solde1 = arrivees1 - departs1

D_val_slide5 = f"{arrivees1:,}".replace(",", " ")
E_val_slide5 = f"{departs1:,}".replace(",", " ")
F_val_slide5 = f"{solde1:,}".replace(",", " ")

# === Remplacement du texte dans la slide 5 ===
for shape in slide5.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide5, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide5, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide5, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide5, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide5, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide5, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide5, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide5, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide5, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION CORSE

# === FILTRAGE Corse  ===
df_immat_2024_co = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_immat_2025_co = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_pcl_2024_co = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_pcl_2025_co = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_radiation_2024_co = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_radiation_2025_co = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_stock_2025_co = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "].reset_index(drop=True)

df_ip_2024_co = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_ip_2025_co = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_sstr_tresor_2024_co = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_sstr_tresor_2025_co = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_sstr_ss_2024_co = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_sstr_ss_2025_co = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_mbe_2024_co = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_mbe_2025_co = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

greffes_co = [
    "AJACCIO", "BASTIA"
]

df_div_2024_co = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_co)].reset_index(drop=True)
df_div_2025_co = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_co)].reset_index(drop=True)

df_dpca_2024_co = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]
df_dpca_2025_co = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Corse "]

df_ts_2024_co = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Corse ") |
    (df_ts_2024["REGIONTO"] == "Corse ")
]

df_ts_2025_co = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Corse ") |
    (df_ts_2025["REGIONTO"] == "Corse ")
]

df_prev_2024_co = df_prev_2024[df_prev_2024["REGION"] == "Corse "]
df_prev_2025_co = df_prev_2025[df_prev_2025["REGION"] == "Corse "]

# === SLIDE 6 ===
slide6 = prs.slides[5]

A_2024 = df_immat_2024_co["SIREN"].notna().sum()
A_2025 = df_immat_2025_co["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_co["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_co["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_co["SIREN"].notna().sum()
C_2025 = df_pcl_2025_co["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_co["SIREN"].notna().sum()
D_2025 = df_radiation_2025_co["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_co["SIREN"].notna().sum()
E_2025 = df_ip_2025_co["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_co["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_co["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_co["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_co["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_co, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_co, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_co, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_co, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_co, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_co, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_co, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_co, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_co, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_co, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 6
for shape in slide6.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 7 ===
slide7 = prs.slides[6]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide7 = df_mbe_2024_co["SIREN"].notna().sum()
A_2025_slide7 = df_mbe_2025_co["SIREN"].notna().sum()
A_val_slide7 = f"{A_2025_slide7:,}".replace(",", " ")
A_pct_slide7 = evol_percent_txt(A_2024_slide7, A_2025_slide7)

# B : DIV
B_2024_slide7 = df_div_2024_co["SIREN"].notna().sum()
B_2025_slide7 = df_div_2025_co["SIREN"].notna().sum()
B_val_slide7 = f"{B_2025_slide7:,}".replace(",", " ")
B_pct_slide7 = evol_percent_txt(B_2024_slide7, B_2025_slide7)

# C : DPCA
C_2024_slide7 = df_dpca_2024_co["SIREN"].notna().sum()
C_2025_slide7 = df_dpca_2025_co["SIREN"].notna().sum()
C_val_slide7 = f"{C_2025_slide7:,}".replace(",", " ")
C_pct_slide7 = evol_percent_txt(C_2024_slide7, C_2025_slide7)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Corse ").sum()
departs = (df_ts_2025["REGIONFROM"] == "Corse ").sum()
solde = arrivees - departs

D_val_slide7 = f"{arrivees:,}".replace(",", " ")
E_val_slide7 = f"{departs:,}".replace(",", " ")
F_val_slide7 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 7 ===
for shape in slide7.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide7, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide7, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide7, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide7, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide7, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide7, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide7, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide7, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide7, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION HAUTS DE FRANCE

# === FILTRAGE Hauts-De-France (Nord-Pas-de-Calais-Picardie) ===
df_immat_2024_hdf = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_immat_2025_hdf = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_pcl_2024_hdf = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_pcl_2025_hdf = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_radiation_2024_hdf = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_radiation_2025_hdf = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_stock_2025_hdf = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"].reset_index(drop=True)

df_ip_2024_hdf = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_ip_2025_hdf = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_sstr_tresor_2024_hdf = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_sstr_tresor_2025_hdf = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_sstr_ss_2024_hdf = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_sstr_ss_2025_hdf = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_mbe_2024_hdf = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_mbe_2025_hdf = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

greffes_hdf = [
    "AMIENS", "ARRAS", "BOBIGNY", "DUNKERQUE", "LILLE METROPOLE", "METZ", "NEVERS", "ROUEN", "SAINT-QUENTIN", "VALENCIENNES"
]

df_div_2024_hdf = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_hdf)].reset_index(drop=True)
df_div_2025_hdf = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_hdf)].reset_index(drop=True)

df_dpca_2024_hdf = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_dpca_2025_hdf = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

df_ts_2024_hdf = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)") |
    (df_ts_2024["REGIONTO"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)")
]

df_ts_2025_hdf = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)") |
    (df_ts_2025["REGIONTO"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)")
]

df_prev_2024_hdf = df_prev_2024[df_prev_2024["REGION"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]
df_prev_2025_hdf = df_prev_2025[df_prev_2025["REGION"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)"]

# === SLIDE 8 ===
slide8 = prs.slides[7]

A_2024 = df_immat_2024_hdf["SIREN"].notna().sum()
A_2025 = df_immat_2025_hdf["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_hdf["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_hdf["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_hdf["SIREN"].notna().sum()
C_2025 = df_pcl_2025_hdf["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_hdf["SIREN"].notna().sum()
D_2025 = df_radiation_2025_hdf["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_hdf["SIREN"].notna().sum()
E_2025 = df_ip_2025_hdf["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_hdf["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_hdf["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_hdf["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_hdf["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_hdf, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_hdf, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_hdf, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_hdf, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_hdf, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_hdf, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_hdf, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_hdf, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_hdf, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_hdf, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 8
for shape in slide8.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 9 ===
slide9 = prs.slides[8]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide9 = df_mbe_2024_hdf["SIREN"].notna().sum()
A_2025_slide9 = df_mbe_2025_hdf["SIREN"].notna().sum()
A_val_slide9 = f"{A_2025_slide9:,}".replace(",", " ")
A_pct_slide9 = evol_percent_txt(A_2024_slide9, A_2025_slide9)

# B : DIV
B_2024_slide9 = df_div_2024_hdf["SIREN"].notna().sum()
B_2025_slide9 = df_div_2025_hdf["SIREN"].notna().sum()
B_val_slide9 = f"{B_2025_slide9:,}".replace(",", " ")
B_pct_slide9 = evol_percent_txt(B_2024_slide9, B_2025_slide9)

# C : DPCA
C_2024_slide9 = df_dpca_2024_hdf["SIREN"].notna().sum()
C_2025_slide9 = df_dpca_2025_hdf["SIREN"].notna().sum()
C_val_slide9 = f"{C_2025_slide9:,}".replace(",", " ")
C_pct_slide9 = evol_percent_txt(C_2024_slide9, C_2025_slide9)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)").sum()
departs = (df_ts_2025["REGIONFROM"] == "Hauts-De-France (Nord-Pas-de-Calais-Picardie)").sum()
solde = arrivees - departs

D_val_slide9 = f"{arrivees:,}".replace(",", " ")
E_val_slide9 = f"{departs:,}".replace(",", " ")
F_val_slide9 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 9 ===
for shape in slide9.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide9, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide9, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide9, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide9, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide9, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide9, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide9, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide9, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide9, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION BOURGOGNE FRANCHE COMTE

# === FILTRAGE Bourgogne-Franche-Comté ===
df_immat_2024_bfc = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_immat_2025_bfc = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_pcl_2024_bfc = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_pcl_2025_bfc = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_radiation_2024_bfc = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_radiation_2025_bfc = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_stock_2025_bfc = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"].reset_index(drop=True)

df_ip_2024_bfc = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_ip_2025_bfc = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_sstr_tresor_2024_bfc = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_sstr_tresor_2025_bfc = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_sstr_ss_2024_bfc = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_sstr_ss_2025_bfc = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_mbe_2024_bfc = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_mbe_2025_bfc = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

greffes_bfc = [
    "AUXERRE",
    "CHALON-SUR-SAONE",
    "DIJON",
    "LONS-LE-SAUNIER",
    "MACON",
    "NEVERS",
    "SENS",
    "VESOUL",
    "BELFORT",
    "BESANCON"
]

df_div_2024_bfc = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bfc)].reset_index(drop=True)
df_div_2025_bfc = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bfc)].reset_index(drop=True)

df_dpca_2024_bfc = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]
df_dpca_2025_bfc = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bourgogne-Franche-Comté"]

df_ts_2024_bfc = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Bourgogne-Franche-Comté") |
    (df_ts_2024["REGIONTO"] == "Bourgogne-Franche-Comté")
]

df_ts_2025_bfc = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Bourgogne-Franche-Comté") |
    (df_ts_2025["REGIONTO"] == "Bourgogne-Franche-Comté")
]

df_prev_2024_bfc = df_prev_2024[df_prev_2024["REGION"] == "Bourgogne-Franche-Comté"]
df_prev_2025_bfc = df_prev_2025[df_prev_2025["REGION"] == "Bourgogne-Franche-Comté"]

# === SLIDE 10 ===
slide10 = prs.slides[9]

A_2024 = df_immat_2024_bfc["SIREN"].notna().sum()
A_2025 = df_immat_2025_bfc["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bfc["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bfc["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bfc["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bfc["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bfc["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bfc["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bfc["SIREN"].notna().sum()
E_2025 = df_ip_2025_bfc["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bfc["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bfc["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bfc["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bfc["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bfc, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bfc, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bfc, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bfc, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bfc, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bfc, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bfc, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bfc, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bfc, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bfc, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 10
for shape in slide10.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 11 ===
slide11 = prs.slides[10]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide11 = df_mbe_2024_bfc["SIREN"].notna().sum()
A_2025_slide11 = df_mbe_2025_bfc["SIREN"].notna().sum()
A_val_slide11 = f"{A_2025_slide11:,}".replace(",", " ")
A_pct_slide11 = evol_percent_txt(A_2024_slide11, A_2025_slide11)

# B : DIV
B_2024_slide11 = df_div_2024_bfc["SIREN"].notna().sum()
B_2025_slide11 = df_div_2025_bfc["SIREN"].notna().sum()
B_val_slide11 = f"{B_2025_slide11:,}".replace(",", " ")
B_pct_slide11 = evol_percent_txt(B_2024_slide11, B_2025_slide11)

# C : DPCA
C_2024_slide11 = df_dpca_2024_bfc["SIREN"].notna().sum()
C_2025_slide11 = df_dpca_2025_bfc["SIREN"].notna().sum()
C_val_slide11 = f"{C_2025_slide11:,}".replace(",", " ")
C_pct_slide11 = evol_percent_txt(C_2024_slide11, C_2025_slide11)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Bourgogne-Franche-Comté").sum()
departs = (df_ts_2025["REGIONFROM"] == "Bourgogne-Franche-Comté").sum()
solde = arrivees - departs

D_val_slide11 = f"{arrivees:,}".replace(",", " ")
E_val_slide11 = f"{departs:,}".replace(",", " ")
F_val_slide11 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 11 ===
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide11, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide11, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide11, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide11, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide11, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide11, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide11, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide11, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide11, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION AUVERGNE RHÔNE ALPES

# === FILTRAGE Auvergne-Rhône-Alpes ===
df_immat_2024_aura = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_immat_2025_aura = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_pcl_2024_aura = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_pcl_2025_aura = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_radiation_2024_aura = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_radiation_2025_aura = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_stock_2025_aura = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"].reset_index(drop=True)

df_ip_2024_aura = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_ip_2025_aura = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_sstr_tresor_2024_aura = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_sstr_tresor_2025_aura = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_sstr_ss_2024_aura = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_sstr_ss_2025_aura = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_mbe_2024_aura = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_mbe_2025_aura = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

greffes_aura = [
    "AUBENAS",
    "AURILLAC",
    "BOURG-EN-BRESSE",
    "CHAMBERY",
    "CLERMONT-FERRAND",
    "CUSSET",
    "GRENOBLE",
    "LE PUY-EN-VELAY",
    "LYON",
    "ROANNE",
    "ROMANS",
    "SAINT-ETIENNE",
    "THONON-LES-BAINS",
    "VIENNE",
    "VILLEFRANCHE-TARARE"
]

df_div_2024_aura = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_aura)].reset_index(drop=True)
df_div_2025_aura = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_aura)].reset_index(drop=True)

df_dpca_2024_aura = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]
df_dpca_2025_aura = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Auvergne-Rhône-Alpes"]

df_ts_2024_aura = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Auvergne-Rhône-Alpes") |
    (df_ts_2024["REGIONTO"] == "Auvergne-Rhône-Alpes")
]

df_ts_2025_aura = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Auvergne-Rhône-Alpes") |
    (df_ts_2025["REGIONTO"] == "Auvergne-Rhône-Alpes")
]

df_prev_2024_aura = df_prev_2024[df_prev_2024["REGION"] == "Auvergne-Rhône-Alpes"]
df_prev_2025_aura = df_prev_2025[df_prev_2025["REGION"] == "Auvergne-Rhône-Alpes"]

# === SLIDE 12 ===
slide12 = prs.slides[11]

A_2024 = df_immat_2024_aura["SIREN"].notna().sum()
A_2025 = df_immat_2025_aura["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_aura["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_aura["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_aura["SIREN"].notna().sum()
C_2025 = df_pcl_2025_aura["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_aura["SIREN"].notna().sum()
D_2025 = df_radiation_2025_aura["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_aura["SIREN"].notna().sum()
E_2025 = df_ip_2025_aura["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_aura["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_aura["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_aura["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_aura["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_aura, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_aura, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_aura, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_aura, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_aura, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_aura, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_aura, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_aura, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_aura, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_aura, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 12
for shape in slide12.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 13 ===
slide13 = prs.slides[12]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide13 = df_mbe_2024_aura["SIREN"].notna().sum()
A_2025_slide13 = df_mbe_2025_aura["SIREN"].notna().sum()
A_val_slide13 = f"{A_2025_slide13:,}".replace(",", " ")
A_pct_slide13 = evol_percent_txt(A_2024_slide13, A_2025_slide13)

# B : DIV
B_2024_slide13 = df_div_2024_aura["SIREN"].notna().sum()
B_2025_slide13 = df_div_2025_aura["SIREN"].notna().sum()
B_val_slide13 = f"{B_2025_slide13:,}".replace(",", " ")
B_pct_slide13 = evol_percent_txt(B_2024_slide13, B_2025_slide13)

# C : DPCA
C_2024_slide13 = df_dpca_2024_aura["SIREN"].notna().sum()
C_2025_slide13 = df_dpca_2025_aura["SIREN"].notna().sum()
C_val_slide13 = f"{C_2025_slide13:,}".replace(",", " ")
C_pct_slide13 = evol_percent_txt(C_2024_slide13, C_2025_slide13)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Auvergne-Rhône-Alpes").sum()
departs = (df_ts_2025["REGIONFROM"] == "Auvergne-Rhône-Alpes").sum()
solde = arrivees - departs

D_val_slide13 = f"{arrivees:,}".replace(",", " ")
E_val_slide13 = f"{departs:,}".replace(",", " ")
F_val_slide13 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 3 ===
for shape in slide13.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide13
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide13, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide13, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide13, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide13, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide13, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide13, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide13, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide13, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide13, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION ILE DE FRANCE

# === FILTRAGE idfitanie (Languedoc-Roussillon-Midi-Pyrénées) ===
df_immat_2024_idf = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_immat_2025_idf = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_pcl_2024_idf = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_pcl_2025_idf = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_radiation_2024_idf = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_radiation_2025_idf = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_stock_2025_idf = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "].reset_index(drop=True)

df_ip_2024_idf = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_ip_2025_idf = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_sstr_tresor_2024_idf = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_sstr_tresor_2025_idf = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_sstr_ss_2024_idf = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_sstr_ss_2025_idf = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_mbe_2024_idf = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_mbe_2025_idf = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

greffes_idf = [
    "PARIS",
    "NANTERRE",
    "EVRY",
    "BOBIGNY",
    "CRETEIL",
    "MEAUX",
    "MELUN",
    "PONTOISE",
    "VERSAILLES"
]

df_div_2024_idf = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_idf)].reset_index(drop=True)
df_div_2025_idf = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_idf)].reset_index(drop=True)

df_dpca_2024_idf = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]
df_dpca_2025_idf = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Ile-de-France "]

df_ts_2024_idf = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Ile-de-France  ") |
    (df_ts_2024["REGIONTO"] == "Ile-de-France  ")
]

df_ts_2025_idf = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Ile-de-France ") |
    (df_ts_2025["REGIONTO"] == "Ile-de-France ")
]

df_prev_2024_idf = df_prev_2024[df_prev_2024["REGION"] == "Ile-de-France "]
df_prev_2025_idf = df_prev_2025[df_prev_2025["REGION"] == "Ile-de-France "]

# === SLIDE 14 ===
slide14 = prs.slides[13]

A_2024 = df_immat_2024_idf["SIREN"].notna().sum()
A_2025 = df_immat_2025_idf["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_idf["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_idf["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_idf["SIREN"].notna().sum()
C_2025 = df_pcl_2025_idf["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_idf["SIREN"].notna().sum()
D_2025 = df_radiation_2025_idf["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_idf["SIREN"].notna().sum()
E_2025 = df_ip_2025_idf["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_idf["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_idf["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_idf["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_idf["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_idf, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_idf, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_idf, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_idf, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_idf, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_idf, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_idf, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_idf, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_idf, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_idf, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

# Helper: compute percent text consistently
def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Replace on slide 14
for shape in slide14.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 15 ===
slide15 = prs.slides[14]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide15 = df_mbe_2024_idf["SIREN"].notna().sum()
A_2025_slide15 = df_mbe_2025_idf["SIREN"].notna().sum()
A_val_slide15 = f"{A_2025_slide15:,}".replace(",", " ")
A_pct_slide15 = evol_percent_txt(A_2024_slide15, A_2025_slide15)

# B : DIV
B_2024_slide15 = df_div_2024_idf["SIREN"].notna().sum()
B_2025_slide15 = df_div_2025_idf["SIREN"].notna().sum()
B_val_slide15 = f"{B_2025_slide15:,}".replace(",", " ")
B_pct_slide15 = evol_percent_txt(B_2024_slide15, B_2025_slide15)

# C : DPCA
C_2024_slide15 = df_dpca_2024_idf["SIREN"].notna().sum()
C_2025_slide15 = df_dpca_2025_idf["SIREN"].notna().sum()
C_val_slide15 = f"{C_2025_slide15:,}".replace(",", " ")
C_pct_slide15 = evol_percent_txt(C_2024_slide15, C_2025_slide15)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Ile-de-France ").sum()
departs = (df_ts_2025["REGIONFROM"] == "Ile-de-France ").sum()
solde = arrivees - departs

D_val_slide15 = f"{arrivees:,}".replace(",", " ")
E_val_slide15 = f"{departs:,}".replace(",", " ")
F_val_slide15 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 15 ===
for shape in slide15.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide15
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide15, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide15, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide15, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide15, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide15, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide15, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide15, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide15, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide15, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION BRETAGNE

# === FILTRAGE Bretagne ===
df_immat_2024_bre = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_immat_2025_bre = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_pcl_2024_bre = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_pcl_2025_bre = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_radiation_2024_bre = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_radiation_2025_bre = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_stock_2025_bre = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "].reset_index(drop=True)

df_ip_2024_bre = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_ip_2025_bre = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_sstr_tresor_2024_bre = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_sstr_tresor_2025_bre = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_sstr_ss_2024_bre = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_sstr_ss_2025_bre = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_mbe_2024_bre = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_mbe_2025_bre = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

greffes_bre = [
    "BREST", "QUIMPER", "RENNES", "SAINT-BRIEUC", "SAINT MALO", "VANNES"
]

df_div_2024_bre = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)
df_div_2025_bre = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bre)].reset_index(drop=True)

df_dpca_2024_bre = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]
df_dpca_2025_bre = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Bretagne "]

df_ts_2024_bre = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Bretagne ") |
    (df_ts_2024["REGIONTO"] == "Bretagne ")
]

df_ts_2025_bre = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Bretagne ") |
    (df_ts_2025["REGIONTO"] == "Bretagne ")
]

df_prev_2024_bre = df_prev_2024[df_prev_2024["REGION"] == "Bretagne "]
df_prev_2025_bre = df_prev_2025[df_prev_2025["REGION"] == "Bretagne "]

# === SLIDE 16 ===
slide16 = prs.slides[15]

A_2024 = df_immat_2024_bre["SIREN"].notna().sum()
A_2025 = df_immat_2025_bre["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bre["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bre["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bre["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bre["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bre["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bre["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bre["SIREN"].notna().sum()
E_2025 = df_ip_2025_bre["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bre["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bre["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bre["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bre["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bre, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bre, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bre, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bre, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bre, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bre, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bre, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bre, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bre, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bre, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer dans la slide 16
for shape in slide16.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 17 ===
slide17 = prs.slides[16]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide17 = df_mbe_2024_bre["SIREN"].notna().sum()
A_2025_slide17 = df_mbe_2025_bre["SIREN"].notna().sum()
A_val_slide17 = f"{A_2025_slide17:,}".replace(",", " ")
A_pct_slide17 = evol_percent_txt(A_2024_slide17, A_2025_slide17)

# B : DIV
B_2024_slide17 = df_div_2024_bre["SIREN"].notna().sum()
B_2025_slide17 = df_div_2025_bre["SIREN"].notna().sum()
B_val_slide17 = f"{B_2025_slide17:,}".replace(",", " ")
B_pct_slide17 = evol_percent_txt(B_2024_slide17, B_2025_slide17)

# C : DPCA
C_2024_slide17 = df_dpca_2024_bre["SIREN"].notna().sum()
C_2025_slide17 = df_dpca_2025_bre["SIREN"].notna().sum()
C_val_slide17 = f"{C_2025_slide17:,}".replace(",", " ")
C_pct_slide17 = evol_percent_txt(C_2024_slide17, C_2025_slide17)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Bretagne ").sum()
departs = (df_ts_2025["REGIONFROM"] == "Bretagne ").sum()
solde = arrivees - departs

D_val_slide17 = f"{arrivees:,}".replace(",", " ")
E_val_slide17 = f"{departs:,}".replace(",", " ")
F_val_slide17 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 17 ===
for shape in slide17.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide17
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide17, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide17, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide17, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide17, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide17, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide17, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide17, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide17, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide17, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION OCCITANIE

# === FILTRAGE Occitanie ===
df_immat_2024_occ = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_immat_2025_occ = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_pcl_2024_occ = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_pcl_2025_occ = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_radiation_2024_occ = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_radiation_2025_occ = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_stock_2025_occ = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"].reset_index(drop=True)

df_ip_2024_occ = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_ip_2025_occ = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_sstr_tresor_2024_occ = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_sstr_tresor_2025_occ = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_sstr_ss_2024_occ = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_sstr_ss_2025_occ = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_mbe_2024_occ = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_mbe_2025_occ = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

greffes_occ = [
    "ALBI",
    "AUCH",
    "BEZIERS",
    "CARCASSONNE",
    "CASTRES",
    "CAHORS",
    "FOIX",
    "MENDE",
    "MONTAUBAN",
    "MONTPELLIER",
    "NARBONNE",
    "NÎMES",
    "PERPIGNAN",
    "RODEZ",
    "TARBES",
    "TOULOUSE"
]

df_div_2024_occ = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_occ)].reset_index(drop=True)
df_div_2025_occ = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_occ)].reset_index(drop=True)

df_dpca_2024_occ = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_dpca_2025_occ = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

df_ts_2024_occ = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)") |
    (df_ts_2024["REGIONTO"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)")
]

df_ts_2025_occ = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)") |
    (df_ts_2025["REGIONTO"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)")
]

df_prev_2024_occ = df_prev_2024[df_prev_2024["REGION"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]
df_prev_2025_occ = df_prev_2025[df_prev_2025["REGION"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)"]

# === SLIDE 18 ===
slide18 = prs.slides[17]

A_2024 = df_immat_2024_occ["SIREN"].notna().sum()
A_2025 = df_immat_2025_occ["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_occ["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_occ["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_occ["SIREN"].notna().sum()
C_2025 = df_pcl_2025_occ["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_occ["SIREN"].notna().sum()
D_2025 = df_radiation_2025_occ["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_occ["SIREN"].notna().sum()
E_2025 = df_ip_2025_occ["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_occ["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_occ["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_occ["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_occ["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_occ, "Nomocc total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_occ, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_occ, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_occ, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_occ, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_occ, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_occ, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_occ, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_occ, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_occ, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

# Helper: compute percent text consistently
def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 18
for shape in slide18.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 19 ===
slide19 = prs.slides[18]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide19 = df_mbe_2024_occ["SIREN"].notna().sum()
A_2025_slide19 = df_mbe_2025_occ["SIREN"].notna().sum()
A_val_slide19 = f"{A_2025_slide19:,}".replace(",", " ")
A_pct_slide19 = evol_percent_txt(A_2024_slide19, A_2025_slide19)

# B : DIV
B_2024_slide19 = df_div_2024_occ["SIREN"].notna().sum()
B_2025_slide19 = df_div_2025_occ["SIREN"].notna().sum()
B_val_slide19 = f"{B_2025_slide19:,}".replace(",", " ")
B_pct_slide19 = evol_percent_txt(B_2024_slide19, B_2025_slide19)

# C : DPCA
C_2024_slide19 = df_dpca_2024_occ["SIREN"].notna().sum()
C_2025_slide19 = df_dpca_2025_occ["SIREN"].notna().sum()
C_val_slide19 = f"{C_2025_slide19:,}".replace(",", " ")
C_pct_slide19 = evol_percent_txt(C_2024_slide19, C_2025_slide19)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)").sum()
departs = (df_ts_2025["REGIONFROM"] == "Occitanie (Languedoc-Roussillon-Midi-Pyrénées)").sum()
solde = arrivees - departs

D_val_slide19 = f"{arrivees:,}".replace(",", " ")
E_val_slide19 = f"{departs:,}".replace(",", " ")
F_val_slide19 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 3 ===
for shape in slide19.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide19
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide19, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide19, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide19, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide19, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide19, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide19, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide19, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide19, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide19, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION PROVENCE ALPES COTE D'AZUR

# === FILTRAGE PACA ===
df_immat_2024_paca = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_immat_2025_paca = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_pcl_2024_paca = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_pcl_2025_paca = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_radiation_2024_paca = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_radiation_2025_paca = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_stock_2025_paca = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"].reset_index(drop=True)

df_ip_2024_paca = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_ip_2025_paca = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_sstr_tresor_2024_paca = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_sstr_tresor_2025_paca = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_sstr_ss_2024_paca = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_sstr_ss_2025_paca = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_mbe_2024_paca = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_mbe_2025_paca = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

greffes_paca = [
    "AIX-EN-PROVENCE",
    "ANTIBES",
    "CANNES",
    "DRAGUIGNAN",
    "FRÉJUS",
    "GAP",
    "GRASSE",
    "MANOSQUE",
    "MARSEILLE",
    "NICE",
    "SALON-DE-PROVENCE",
    "TARASCON",
    "TOULON"
]

df_div_2024_paca = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_paca)].reset_index(drop=True)
df_div_2025_paca = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_paca)].reset_index(drop=True)

df_dpca_2024_paca = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]
df_dpca_2025_paca = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Provence-Alpes-Côte d'Azur"]

df_ts_2024_paca = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Provence-Alpes-Côte d'Azur") |
    (df_ts_2024["REGIONTO"] == "Provence-Alpes-Côte d'Azur")
]

df_ts_2025_paca = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Provence-Alpes-Côte d'Azur") |
    (df_ts_2025["REGIONTO"] == "Provence-Alpes-Côte d'Azur")
]

df_prev_2024_paca = df_prev_2024[df_prev_2024["REGION"] == "Provence-Alpes-Côte d'Azur"]
df_prev_2025_paca = df_prev_2025[df_prev_2025["REGION"] == "Provence-Alpes-Côte d'Azur"]

# === SLIDE 20 ===
slide20 = prs.slides[19]

A_2024 = df_immat_2024_paca["SIREN"].notna().sum()
A_2025 = df_immat_2025_paca["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_paca["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_paca["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_paca["SIREN"].notna().sum()
C_2025 = df_pcl_2025_paca["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_paca["SIREN"].notna().sum()
D_2025 = df_radiation_2025_paca["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_paca["SIREN"].notna().sum()
E_2025 = df_ip_2025_paca["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_paca["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_paca["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_paca["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_paca["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_paca, "Nompaca total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_paca, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_paca, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_paca, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_paca, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_paca, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_paca, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_paca, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_paca, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_paca, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"


def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer la slide 20
for shape in slide20.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 21 ===
slide21 = prs.slides[20]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide21 = df_mbe_2024_paca["SIREN"].notna().sum()
A_2025_slide21 = df_mbe_2025_paca["SIREN"].notna().sum()
A_val_slide21 = f"{A_2025_slide21:,}".replace(",", " ")
A_pct_slide21 = evol_percent_txt(A_2024_slide21, A_2025_slide21)

# B : DIV
B_2024_slide21 = df_div_2024_paca["SIREN"].notna().sum()
B_2025_slide21 = df_div_2025_paca["SIREN"].notna().sum()
B_val_slide21 = f"{B_2025_slide21:,}".replace(",", " ")
B_pct_slide21 = evol_percent_txt(B_2024_slide21, B_2025_slide21)

# C : DPCA
C_2024_slide21 = df_dpca_2024_paca["SIREN"].notna().sum()
C_2025_slide21 = df_dpca_2025_paca["SIREN"].notna().sum()
C_val_slide21 = f"{C_2025_slide21:,}".replace(",", " ")
C_pct_slide21 = evol_percent_txt(C_2024_slide21, C_2025_slide21)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Provence-Alpes-Côte d'Azur").sum()
departs = (df_ts_2025["REGIONFROM"] == "Provence-Alpes-Côte d'Azur").sum()
solde = arrivees - departs

D_val_slide21 = f"{arrivees:,}".replace(",", " ")
E_val_slide21 = f"{departs:,}".replace(",", " ")
F_val_slide21 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 21 ===
for shape in slide21.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide21
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide21, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide21, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide21, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide21, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide21, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide21, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide21, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide21, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide21, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION Normandie

# === FILTRAGE nor ===
df_immat_2024_nor = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_immat_2025_nor = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_pcl_2024_nor = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_pcl_2025_nor = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_radiation_2024_nor = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_radiation_2025_nor = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_stock_2025_nor = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"].reset_index(drop=True)

df_ip_2024_nor = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_ip_2025_nor = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_sstr_tresor_2024_nor = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_sstr_tresor_2025_nor = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_sstr_ss_2024_nor = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_sstr_ss_2025_nor = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_mbe_2024_nor = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_mbe_2025_nor = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

greffes_nor = [
    "ALENCON",
    "BERNAY",
    "CAEN",
    "CHERBOURG",
    "COUTANCES",
    "DIEPPE",
    "EVREUX",
    "LE HAVRE",
    "LISIEUX",
    "ROUEN"
]

df_div_2024_nor = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_nor)].reset_index(drop=True)
df_div_2025_nor = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_nor)].reset_index(drop=True)

df_dpca_2024_nor = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]
df_dpca_2025_nor = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Normandie"]

df_ts_2024_nor = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Normandie") |
    (df_ts_2024["REGIONTO"] == "Normandie")
]

df_ts_2025_nor = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Normandie") |
    (df_ts_2025["REGIONTO"] == "Normandie")
]

df_prev_2024_nor = df_prev_2024[df_prev_2024["REGION"] == "Normandie"]
df_prev_2025_nor = df_prev_2025[df_prev_2025["REGION"] == "Normandie"]

# === SLIDE 22 ===
slide22 = prs.slides[21]

A_2024 = df_immat_2024_nor["SIREN"].notna().sum()
A_2025 = df_immat_2025_nor["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_nor["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_nor["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_nor["SIREN"].notna().sum()
C_2025 = df_pcl_2025_nor["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_nor["SIREN"].notna().sum()
D_2025 = df_radiation_2025_nor["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_nor["SIREN"].notna().sum()
E_2025 = df_ip_2025_nor["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_nor["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_nor["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_nor["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_nor["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_nor, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_nor, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_nor, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_nor, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_nor, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_nor, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_nor, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_nor, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_nor, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_nor, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 22
for shape in slide22.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 23 ===
slide23 = prs.slides[22]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide23 = df_mbe_2024_nor["SIREN"].notna().sum()
A_2025_slide23 = df_mbe_2025_nor["SIREN"].notna().sum()
A_val_slide23 = f"{A_2025_slide23:,}".replace(",", " ")
A_pct_slide23 = evol_percent_txt(A_2024_slide23, A_2025_slide23)

# B : DIV
B_2024_slide23 = df_div_2024_nor["SIREN"].notna().sum()
B_2025_slide23 = df_div_2025_nor["SIREN"].notna().sum()
B_val_slide23 = f"{B_2025_slide23:,}".replace(",", " ")
B_pct_slide23 = evol_percent_txt(B_2024_slide23, B_2025_slide23)

# C : DPCA
C_2024_slide23 = df_dpca_2024_nor["SIREN"].notna().sum()
C_2025_slide23 = df_dpca_2025_nor["SIREN"].notna().sum()
C_val_slide23 = f"{C_2025_slide23:,}".replace(",", " ")
C_pct_slide23 = evol_percent_txt(C_2024_slide23, C_2025_slide23)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Normandie").sum()
departs = (df_ts_2025["REGIONFROM"] == "Normandie").sum()
solde = arrivees - departs

D_val_slide23 = f"{arrivees:,}".replace(",", " ")
E_val_slide23 = f"{departs:,}".replace(",", " ")
F_val_slide23 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 23 ===
for shape in slide23.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide23
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide23, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide23, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide23, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide23, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide23, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide23, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide23, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide23, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide23, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION NOUVELLE AQUITAINE

# === FILTRAGE NOUVELLE AQUITAINE ===
df_immat_2024_nouva = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_immat_2025_nouva = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_pcl_2024_nouva = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_pcl_2025_nouva = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_radiation_2024_nouva = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_radiation_2025_nouva = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_stock_2025_nouva = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"].reset_index(drop=True)

df_ip_2024_nouva = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_ip_2025_nouva = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_sstr_tresor_2024_nouva = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_sstr_tresor_2025_nouva = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_sstr_ss_2024_nouva = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_sstr_ss_2025_nouva = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_mbe_2024_nouva = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_mbe_2025_nouva = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

greffes_nouva = [
    "AGEN",
    "ANGOULEME",
    "BAYONNE",
    "BERGERAC",
    "BORDEAUX",
    "BRIVE",
    "DAX",
    "GUERET",
    "LIBOURNE",
    "LIMOGES",
    "NIORT",
    "PAU",
    "PERIGUEUX",
    "POITIERS",
    "SAINTES"
]

df_div_2024_nouva = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_nouva)].reset_index(drop=True)
df_div_2025_nouva = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_nouva)].reset_index(drop=True)

df_dpca_2024_nouva = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_dpca_2025_nouva = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

df_ts_2024_nouva = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)") |
    (df_ts_2024["REGIONTO"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)")
]

df_ts_2025_nouva = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)") |
    (df_ts_2025["REGIONTO"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)")
]

df_prev_2024_nouva = df_prev_2024[df_prev_2024["REGION"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]
df_prev_2025_nouva = df_prev_2025[df_prev_2025["REGION"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)"]

# === SLIDE 24 ===
slide24 = prs.slides[23]

A_2024 = df_immat_2024_nouva["SIREN"].notna().sum()
A_2025 = df_immat_2025_nouva["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_nouva["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_nouva["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_nouva["SIREN"].notna().sum()
C_2025 = df_pcl_2025_nouva["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_nouva["SIREN"].notna().sum()
D_2025 = df_radiation_2025_nouva["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_nouva["SIREN"].notna().sum()
E_2025 = df_ip_2025_nouva["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_nouva["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_nouva["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_nouva["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_nouva["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_nouva, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_nouva, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_nouva, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_nouva, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_nouva, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_nouva, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_nouva, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_nouva, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_nouva, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_nouva, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 24
for shape in slide24.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 25 ===
slide25 = prs.slides[24]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide25 = df_mbe_2024_nouva["SIREN"].notna().sum()
A_2025_slide25 = df_mbe_2025_nouva["SIREN"].notna().sum()
A_val_slide25 = f"{A_2025_slide25:,}".replace(",", " ")
A_pct_slide25 = evol_percent_txt(A_2024_slide25, A_2025_slide25)

# B : DIV
B_2024_slide25 = df_div_2024_nouva["SIREN"].notna().sum()
B_2025_slide25 = df_div_2025_nouva["SIREN"].notna().sum()
B_val_slide25 = f"{B_2025_slide25:,}".replace(",", " ")
B_pct_slide25 = evol_percent_txt(B_2024_slide25, B_2025_slide25)

# C : DPCA
C_2024_slide25 = df_dpca_2024_nouva["SIREN"].notna().sum()
C_2025_slide25 = df_dpca_2025_nouva["SIREN"].notna().sum()
C_val_slide25 = f"{C_2025_slide25:,}".replace(",", " ")
C_pct_slide25 = evol_percent_txt(C_2024_slide25, C_2025_slide25)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)").sum()
departs = (df_ts_2025["REGIONFROM"] == "Nouvelle Aquitaine (Aquitaine-Limousin-Poitou-Charentes)").sum()
solde = arrivees - departs

D_val_slide25 = f"{arrivees:,}".replace(",", " ")
E_val_slide25 = f"{departs:,}".replace(",", " ")
F_val_slide25 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 25 ===
for shape in slide25.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide25
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide25, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide25, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide25, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide25, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide25, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide25, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide25, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide25, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide25, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION PAYS DE LA LOIRE

# === FILTRAGE PAYS DE LA LOIRE ===
df_immat_2024_pdl = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_immat_2025_pdl = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_pcl_2024_pdl = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_pcl_2025_pdl = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_radiation_2024_pdl = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_radiation_2025_pdl = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_stock_2025_pdl = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "].reset_index(drop=True)

df_ip_2024_pdl = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_ip_2025_pdl = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_sstr_tresor_2024_pdl = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_sstr_tresor_2025_pdl = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_sstr_ss_2024_pdl = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_sstr_ss_2025_pdl = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_mbe_2024_pdl = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_mbe_2025_pdl = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

greffes_pdl = [
    "ANGERS",
    "LAVAL",
    "LE MANS",
    "LA ROCHE-SUR-YON",
    "NANTES",
    "SAINT-NAZAIRE"
]

df_div_2024_pdl = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_pdl)].reset_index(drop=True)
df_div_2025_pdl = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_pdl)].reset_index(drop=True)

df_dpca_2024_pdl = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]
df_dpca_2025_pdl = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Pays de la Loire "]

df_ts_2024_pdl = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Pays de la Loire ") |
    (df_ts_2024["REGIONTO"] == "Pays de la Loire ")
]

df_ts_2025_pdl = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Pays de la Loire ") |
    (df_ts_2025["REGIONTO"] == "Pays de la Loire ")
]

df_prev_2024_pdl = df_prev_2024[df_prev_2024["REGION"] == "Pays de la Loire "]
df_prev_2025_pdl = df_prev_2025[df_prev_2025["REGION"] == "Pays de la Loire "]

# === SLIDE 26 ===
slide26 = prs.slides[25]

A_2024 = df_immat_2024_pdl["SIREN"].notna().sum()
A_2025 = df_immat_2025_pdl["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_pdl["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_pdl["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_pdl["SIREN"].notna().sum()
C_2025 = df_pcl_2025_pdl["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_pdl["SIREN"].notna().sum()
D_2025 = df_radiation_2025_pdl["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_pdl["SIREN"].notna().sum()
E_2025 = df_ip_2025_pdl["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_pdl["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_pdl["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_pdl["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_pdl["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_pdl, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_pdl, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_pdl, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_pdl, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_pdl, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_pdl, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_pdl, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_pdl, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_pdl, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_pdl, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 26
for shape in slide26.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 27 ===
slide27 = prs.slides[26]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide27 = df_mbe_2024_pdl["SIREN"].notna().sum()
A_2025_slide27 = df_mbe_2025_pdl["SIREN"].notna().sum()
A_val_slide27 = f"{A_2025_slide27:,}".replace(",", " ")
A_pct_slide27 = evol_percent_txt(A_2024_slide27, A_2025_slide27)

# B : DIV
B_2024_slide27 = df_div_2024_pdl["SIREN"].notna().sum()
B_2025_slide27 = df_div_2025_pdl["SIREN"].notna().sum()
B_val_slide27 = f"{B_2025_slide27:,}".replace(",", " ")
B_pct_slide27 = evol_percent_txt(B_2024_slide27, B_2025_slide27)

# C : DPCA
C_2024_slide27 = df_dpca_2024_pdl["SIREN"].notna().sum()
C_2025_slide27 = df_dpca_2025_pdl["SIREN"].notna().sum()
C_val_slide27 = f"{C_2025_slide27:,}".replace(",", " ")
C_pct_slide27 = evol_percent_txt(C_2024_slide27, C_2025_slide27)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Pays de la Loire ").sum()
departs = (df_ts_2025["REGIONFROM"] == "Pays de la Loire ").sum()
solde = arrivees - departs

D_val_slide27 = f"{arrivees:,}".replace(",", " ")
E_val_slide27 = f"{departs:,}".replace(",", " ")
F_val_slide27 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 27 ===
for shape in slide27.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide27
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide27, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide27, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide27, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide27, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide27, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide27, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide27, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide27, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide27, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace("X", "3").replace("x", "3")
                    run.text = run.text.replace("tribunau3", "tribunaux").replace("TRIBUNAU3", "TRIBUNAUX")
                    run.text = run.text.replace("sociau3", "sociaux").replace("SOCIAU3", "SOCIAUX")
                    run.text = run.text.replace("avril", "juillet")
                    run.text = run.text.replace("juin", "septembre")


#   SAUVEGARDE           
output_filename = os.path.join(base_path, "Stats_Régions_new.pptx")
prs.save(output_filename)
print(f"Fichier PowerPoint modifié enregistré sous : {output_filename}") 