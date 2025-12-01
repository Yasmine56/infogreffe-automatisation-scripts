# Imports
import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

base_path = r"C:\Users\ext-yaissa\OneDrive - INFOGREFFE\Documents\Régions DROM"

files_needed = [
    "Modèle Stats DROM.pptx",
    "Immat_2024_t3.xlsx",
    "Immat_2025_t3.xlsx",
    "Pcl_2024_t3.xlsx",
    "Pcl_2025_t3.xlsx",
    "Radiation_2024_t3.xlsx",
    "Radiation_2025_t3.xlsx",
    "StockRCS_2025_t3_au15092025.xlsx",
    "InjPayer_2024_t3.xlsx",
    "InjPayer_2025_t3.xlsx",
    "TrfSiege_2024_t3.xlsx",
    "TrfSiege_2025_t3.xlsx",
    "Dpca_2024_t3.xlsx",
    "Dpca_2025_t3.xlsx",
    "Mbe_2024_t3.xlsx",
    "Mbe_2025_t3.xlsx",
    "InscrSSTR_2024_t3.xlsx",
    "InscrSSTR_2025_t3.xlsx",
    "Div_2024_t3.xlsx",
    "Div_2025_t3.xlsx",
    "Prev_2024_t3.xlsx",
    "Prev_2025_t3.xlsx",
]

# Vérification des fichiers
missing_files = [f for f in files_needed if not os.path.exists(os.path.join(base_path, f))]
if missing_files:
    print("Fichiers manquants :", missing_files)
    print("Veuillez les ajouter dans le dossier :", base_path)
else:
    print("Tous les fichiers sont présents.")

# === CHARGEMENT DES DONNÉES ===
df_immat_2024 = pd.read_excel(os.path.join(base_path, "Immat_2024_t3.xlsx"), sheet_name="Immat_2024_t3")
df_immat_2025 = pd.read_excel(os.path.join(base_path, "Immat_2025_t3.xlsx"), sheet_name="Immat_2025_t3")

df_pcl_2024 = pd.read_excel(os.path.join(base_path, "Pcl_2024_t3.xlsx"), sheet_name="Pcl_2024_t3")
df_pcl_2025 = pd.read_excel(os.path.join(base_path, "Pcl_2025_t3.xlsx"), sheet_name="Pcl_2025_t3")

df_stock_2025 = pd.read_excel(os.path.join(base_path, "StockRCS_2025_t3_au15092025.xlsx"), sheet_name="StockRCS_2025_t3_au15092025")

df_radiation_2024 = pd.read_excel(os.path.join(base_path, "Radiation_2024_t3.xlsx"), sheet_name="Radiation_2024_t3")
df_radiation_2025 = pd.read_excel(os.path.join(base_path, "Radiation_2025_t3.xlsx"), sheet_name="Radiation_2025_t3")

df_ip_2024 = pd.read_excel(os.path.join(base_path, "InjPayer_2024_t3.xlsx"), sheet_name="InjPayer_2024_t3")
df_ip_2025 = pd.read_excel(os.path.join(base_path, "InjPayer_2025_t3.xlsx"), sheet_name="InjPayer_2025_t3")

df_ts_2024 = pd.read_excel(os.path.join(base_path, "TrfSiege_2024_t3.xlsx"), sheet_name="TrfSiege_2024_t3")
df_ts_2025 = pd.read_excel(os.path.join(base_path, "TrfSiege_2025_t3.xlsx"), sheet_name="TrfSiege_2025_t3")

df_mbe_2024 = pd.read_excel(os.path.join(base_path, "Mbe_2024_t3.xlsx"), sheet_name="Mbe_2024_t3")
df_mbe_2025 = pd.read_excel(os.path.join(base_path, "Mbe_2025_t3.xlsx"), sheet_name="Mbe_2025_t3")

df_dpca_2024 = pd.read_excel(os.path.join(base_path, "Dpca_2024_t3.xlsx"), sheet_name="Dpca_2024_t3")
df_dpca_2025 = pd.read_excel(os.path.join(base_path, "Dpca_2025_t3.xlsx"), sheet_name="Dpca_2025_t3")

df_sstr_2024 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2024_t3.xlsx"), sheet_name="InscrSSTR_2024_t3")
df_sstr_2025 = pd.read_excel(os.path.join(base_path, "InscrSSTR_2025_t3.xlsx"), sheet_name="InscrSSTR_2025_t3")

df_div_2024 = pd.read_excel(os.path.join(base_path, "Div_2024_t3.xlsx"), sheet_name="Div_2024_t3")
df_div_2025 = pd.read_excel(os.path.join(base_path, "Div_2025_t3.xlsx"), sheet_name="Div_2025_t3")

df_prev_2024 = pd.read_excel(os.path.join(base_path, "Prev_2024_t3.xlsx"), sheet_name="Prev_2024_t3")
df_prev_2025 = pd.read_excel(os.path.join(base_path, "Prev_2025_t3.xlsx"), sheet_name="Prev_2025_t3")

# === PRÉ-TRAITEMENT COMMUN ===
def preprocess(df):
    dedup_cols = ["GRF", "MIL", "STC", "CHRONO"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_immat_2024, df_immat_2025 = preprocess(df_immat_2024), preprocess(df_immat_2025)
df_pcl_2024, df_pcl_2025 = preprocess(df_pcl_2024), preprocess(df_pcl_2025)
df_radiation_2024, df_radiation_2025 = preprocess(df_radiation_2024), preprocess(df_radiation_2025)
df_mbe_2024, df_mbe_2025 = preprocess(df_mbe_2024), preprocess(df_mbe_2025)
df_dpca_2024, df_dpca_2025 = preprocess(df_dpca_2024), preprocess(df_dpca_2025)

def preprocess3(df):
    # Filtrage des départements sans tenir compte de la casse
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_stock_2025 = preprocess3(df_stock_2025)

def preprocess2(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_AFFAIRE"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Filtrage des départements sans tenir compte de la casse
    dept_col = "ADRESSEFRANCEDEPARTEMENTLIB"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_ip_2024, df_ip_2025 = preprocess2(df_ip_2024), preprocess2(df_ip_2025)

def preprocess4(df):
    dept_cols = ["DEPFROM", "REGIONTO"]
    for col in dept_cols:
        df[col] = df[col].astype(str).str.strip().str.title()
    exclure_depts = ["Bas Rhin", "Haut Rhin", "Moselle"]
    for col in dept_cols:
        df = df[~df[col].isin(exclure_depts)]
    df = df.reset_index(drop=True)
    return df

df_ts_2024, df_ts_2025 = preprocess4(df_ts_2024), preprocess4(df_ts_2025)

def preprocess9(df):
    # Filtrage des départements sans tenir compte de la casse
    dept_col = "DEPARTEMENT"
    df[dept_col] = df[dept_col].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df[dept_col].str.lower().isin(exclure_depts)].reset_index(drop=True)

    return df

df_prev_2024, df_prev_2025 = preprocess9(df_prev_2024), preprocess9(df_prev_2025)

# === PRÉ-TRAITEMENT pour SSTR Trésor ===
def preprocess5(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes du Trésor
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_tresor = df["libelle inscription"].str.contains("trésor|tresor", case=False, na=False)
    df = df[mask_tresor].reset_index(drop=True)

    return df

# === PRÉ-TRAITEMENT pour SSTR Sécurité sociale ===
def preprocess7(df):
    dedup_cols = ["CODE_GREFFE", "NUMERO_INSCRIPTION"]
    for col in dedup_cols:
        df[col] = df[col].astype(str).str.strip().str.upper()
    df = df.drop_duplicates(subset=dedup_cols).reset_index(drop=True)

    # Nettoyage
    df["ADRESSEFRANCEDEPARTEMENTLIB"] = df["ADRESSEFRANCEDEPARTEMENTLIB"].astype(str).str.strip()
    exclure_depts = ["bas rhin", "haut rhin", "moselle"]
    df = df[~df["ADRESSEFRANCEDEPARTEMENTLIB"].str.lower().isin(exclure_depts)].reset_index(drop=True)

    # Filtrer uniquement les lignes sécurité sociale
    df["libelle inscription"] = df["libelle inscription"].astype(str).str.strip()
    mask_ss = df["libelle inscription"].str.contains("sécurité sociale|securite sociale|sécurité|securite", case=False, na=False)
    df = df[mask_ss].reset_index(drop=True)

    return df

df_sstr_tresor_2024 = preprocess5(df_sstr_2024)
df_sstr_tresor_2025 = preprocess5(df_sstr_2025)
df_sstr_ss_2024 = preprocess7(df_sstr_2024)
df_sstr_ss_2025 = preprocess7(df_sstr_2025)


# === OUVRIR LE PPT ===
prs = Presentation(os.path.join(base_path, "Modèle Stats DROM.pptx"))

# === FILTRAGE Grand Est (Alsace-Champagne-Ardenne-Lorraine) ===
df_immat_2024_ge = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_immat_2025_ge = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_pcl_2024_ge = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_pcl_2025_ge = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_radiation_2024_ge = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_radiation_2025_ge = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_stock_2025_ge = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "].reset_index(drop=True)

df_ip_2024_ge = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_ip_2025_ge = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_sstr_tresor_2024_ge = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_sstr_tresor_2025_ge = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_sstr_ss_2024_ge = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_sstr_ss_2025_ge = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_mbe_2024_ge = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_mbe_2025_ge = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

greffes_ge = [
    "POINTE-A-PITRE",
    "BASSE-TERRE"
]

df_div_2024_ge = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_ge)].reset_index(drop=True)
df_div_2025_ge = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_ge)].reset_index(drop=True)

df_dpca_2024_ge = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]
df_dpca_2025_ge = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guadeloupe "]

df_ts_2024_ge = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Guadeloupe ") |
    (df_ts_2024["REGIONTO"] == "Guadeloupe ")
]

df_ts_2025_ge = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Guadeloupe ") |
    (df_ts_2025["REGIONTO"] == "Guadeloupe ")
]

df_prev_2024_ge = df_prev_2024[df_prev_2024["REGION"] == "Guadeloupe "]
df_prev_2025_ge = df_prev_2025[df_prev_2025["REGION"] == "Guadeloupe "]

# === SLIDE 2 ===
slide2 = prs.slides[1]

# Pour être sûr, recompute A..G,W from your DataFrames (copie directe de ta logique)
A_2024 = df_immat_2024_ge["SIREN"].notna().sum()
A_2025 = df_immat_2025_ge["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_ge["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_ge["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_ge["SIREN"].notna().sum()
C_2025 = df_pcl_2025_ge["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_ge["SIREN"].notna().sum()
D_2025 = df_radiation_2025_ge["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_ge["SIREN"].notna().sum()
E_2025 = df_ip_2025_ge["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_ge["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_ge["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_ge["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_ge["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_ge, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_ge, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_ge, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_ge, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_ge, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_ge, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_ge, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_ge, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_ge, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_ge, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

# Helper: compute percent text consistently
def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 2
for shape in slide2.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 3 ===
slide3 = prs.slides[2]

# A : MBE
A_2024_slide3 = df_mbe_2024_ge["SIREN"].notna().sum()
A_2025_slide3 = df_mbe_2025_ge["SIREN"].notna().sum()
A_val_slide3 = f"{A_2025_slide3:,}".replace(",", " ")
A_pct_slide3 = evol_percent_txt(A_2024_slide3, A_2025_slide3)

# B : DIV
B_2024_slide3 = df_div_2024_ge["SIREN"].notna().sum()
B_2025_slide3 = df_div_2025_ge["SIREN"].notna().sum()
B_val_slide3 = f"{B_2025_slide3:,}".replace(",", " ")
B_pct_slide3 = evol_percent_txt(B_2024_slide3, B_2025_slide3)

# C : DPCA
C_2024_slide3 = df_dpca_2024_ge["SIREN"].notna().sum()
C_2025_slide3 = df_dpca_2025_ge["SIREN"].notna().sum()
C_val_slide3 = f"{C_2025_slide3:,}".replace(",", " ")
C_pct_slide3 = evol_percent_txt(C_2024_slide3, C_2025_slide3)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Guadeloupe ").sum()
departs = (df_ts_2025["REGIONFROM"] == "Guadeloupe ").sum()
solde = arrivees - departs

D_val_slide3 = f"{arrivees:,}".replace(",", " ")
E_val_slide3 = f"{departs:,}".replace(",", " ")
F_val_slide3 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 3 ===
for shape in slide3.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide3
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide3, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide3, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide3, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide3, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide3, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide3, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide3, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide3, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide3, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION GUYANE

# === FILTRAGE Guyane  ===
df_immat_2024_cvd = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_immat_2025_cvd = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_pcl_2024_cvd = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_pcl_2025_cvd = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_radiation_2024_cvd = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_radiation_2025_cvd = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_stock_2025_cvd = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "].reset_index(drop=True)

df_ip_2024_cvd = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_ip_2025_cvd = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_sstr_tresor_2024_cvd = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_sstr_tresor_2025_cvd = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_sstr_ss_2024_cvd = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_sstr_ss_2025_cvd = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_mbe_2024_cvd = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_mbe_2025_cvd = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

greffes_cvd = [
    "CAYENNE"
]

df_div_2024_cvd = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_cvd)].reset_index(drop=True)
df_div_2025_cvd = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_cvd)].reset_index(drop=True)

df_dpca_2024_cvd = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]
df_dpca_2025_cvd = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Guyane "]

df_ts_2024_cvd = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Guyane ") |
    (df_ts_2024["REGIONTO"] == "Guyane ")
]

df_ts_2025_cvd = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Guyane ") |
    (df_ts_2025["REGIONTO"] == "Guyane ")
]

df_prev_2024_cvd = df_prev_2024[df_prev_2024["REGION"] == "Guyane "]
df_prev_2025_cvd = df_prev_2025[df_prev_2025["REGION"] == "Guyane "]

# === SLIDE 4 ===
slide4 = prs.slides[3]

A_2024 = df_immat_2024_cvd["SIREN"].notna().sum()
A_2025 = df_immat_2025_cvd["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_cvd["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_cvd["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_cvd["SIREN"].notna().sum()
C_2025 = df_pcl_2025_cvd["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_cvd["SIREN"].notna().sum()
D_2025 = df_radiation_2025_cvd["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_cvd["SIREN"].notna().sum()
E_2025 = df_ip_2025_cvd["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_cvd["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_cvd["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_cvd["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_cvd["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_cvd, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_cvd, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_cvd, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_cvd, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_cvd, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_cvd, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_cvd, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_cvd, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_cvd, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_cvd, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 4
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# --- SLIDE 5 ---
slide5 = prs.slides[4]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide5 = df_mbe_2024_cvd["SIREN"].notna().sum()
A_2025_slide5 = df_mbe_2025_cvd["SIREN"].notna().sum()
A_val_slide5 = f"{A_2025_slide5:,}".replace(",", " ")
A_pct_slide5 = evol_percent_txt(A_2024_slide5, A_2025_slide5)

# B : DIV
B_2024_slide5 = df_div_2024_cvd["SIREN"].notna().sum()
B_2025_slide5 = df_div_2025_cvd["SIREN"].notna().sum()
B_val_slide5 = f"{B_2025_slide5:,}".replace(",", " ")
B_pct_slide5 = evol_percent_txt(B_2024_slide5, B_2025_slide5)

# C : DPCA
C_2024_slide5 = df_dpca_2024_cvd["SIREN"].notna().sum()
C_2025_slide5 = df_dpca_2025_cvd["SIREN"].notna().sum()
C_val_slide5 = f"{C_2025_slide5:,}".replace(",", " ")
C_pct_slide5 = evol_percent_txt(C_2024_slide5, C_2025_slide5)

# D, E, F : arrivées, départs, solde
arrivees1 = (df_ts_2025["REGIONTO"] == "Guyane ").sum()
departs1 = (df_ts_2025["REGIONFROM"] == "Guyane ").sum()
solde1 = arrivees1 - departs1

D_val_slide5 = f"{arrivees1:,}".replace(",", " ")
E_val_slide5 = f"{departs1:,}".replace(",", " ")
F_val_slide5 = f"{solde1:,}".replace(",", " ")

# === Remplacement du texte dans la slide 5 ===
for shape in slide5.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide5
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide5, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide5, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide5, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide5, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide5, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide5, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide5, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide5, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide5, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION LA REUNION

# === FILTRAGE La Réunion  ===
df_immat_2024_co = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_immat_2025_co = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_pcl_2024_co = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_pcl_2025_co = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_radiation_2024_co = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_radiation_2025_co = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_stock_2025_co = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "].reset_index(drop=True)

df_ip_2024_co = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_ip_2025_co = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_sstr_tresor_2024_co = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_sstr_tresor_2025_co = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_sstr_ss_2024_co = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_sstr_ss_2025_co = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_mbe_2024_co = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_mbe_2025_co = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

greffes_co = [
    "SAINT-DENIS-DE-LA-REUNION",
    "SAINT-PIERRE-DE-LA-REUNION"
]

df_div_2024_co = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_co)].reset_index(drop=True)
df_div_2025_co = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_co)].reset_index(drop=True)

df_dpca_2024_co = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]
df_dpca_2025_co = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "La Réunion "]

df_ts_2024_co = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "La Réunion ") |
    (df_ts_2024["REGIONTO"] == "La Réunion ")
]

df_ts_2025_co = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "La Réunion ") |
    (df_ts_2025["REGIONTO"] == "La Réunion ")
]

df_prev_2024_co = df_prev_2024[df_prev_2024["REGION"] == "La Réunion "]
df_prev_2025_co = df_prev_2025[df_prev_2025["REGION"] == "La Réunion "]

# === SLIDE 6 ===
slide6 = prs.slides[5]

A_2024 = df_immat_2024_co["SIREN"].notna().sum()
A_2025 = df_immat_2025_co["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_co["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_co["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_co["SIREN"].notna().sum()
C_2025 = df_pcl_2025_co["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_co["SIREN"].notna().sum()
D_2025 = df_radiation_2025_co["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_co["SIREN"].notna().sum()
E_2025 = df_ip_2025_co["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_co["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_co["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_co["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_co["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_co, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_co, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_co, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_co, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_co, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_co, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_co, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_co, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_co, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_co, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 6
for shape in slide6.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 7 ===
slide7 = prs.slides[6]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide7 = df_mbe_2024_co["SIREN"].notna().sum()
A_2025_slide7 = df_mbe_2025_co["SIREN"].notna().sum()
A_val_slide7 = f"{A_2025_slide7:,}".replace(",", " ")
A_pct_slide7 = evol_percent_txt(A_2024_slide7, A_2025_slide7)

# B : DIV
B_2024_slide7 = df_div_2024_co["SIREN"].notna().sum()
B_2025_slide7 = df_div_2025_co["SIREN"].notna().sum()
B_val_slide7 = f"{B_2025_slide7:,}".replace(",", " ")
B_pct_slide7 = evol_percent_txt(B_2024_slide7, B_2025_slide7)

# C : DPCA
C_2024_slide7 = df_dpca_2024_co["SIREN"].notna().sum()
C_2025_slide7 = df_dpca_2025_co["SIREN"].notna().sum()
C_val_slide7 = f"{C_2025_slide7:,}".replace(",", " ")
C_pct_slide7 = evol_percent_txt(C_2024_slide7, C_2025_slide7)

# D, E, F : arrivées, départs, solde 
arrivees = (df_ts_2025["REGIONTO"] == "La Réunion ").sum()
departs = (df_ts_2025["REGIONFROM"] == "La Réunion ").sum()
solde = arrivees - departs

D_val_slide7 = f"{arrivees:,}".replace(",", " ")
E_val_slide7 = f"{departs:,}".replace(",", " ")
F_val_slide7 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 3 ===
for shape in slide7.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            if text == "A":
                run.text = A_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide7
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide7, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide7, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide7, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide7, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide7, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide7, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide7, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide7, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide7, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION MARTINIQUE

# === FILTRAGE Martinique  ===
df_immat_2024_hdf = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_immat_2025_hdf = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_pcl_2024_hdf = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_pcl_2025_hdf = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_radiation_2024_hdf = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_radiation_2025_hdf = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_stock_2025_hdf = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "].reset_index(drop=True)

df_ip_2024_hdf = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_ip_2025_hdf = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_sstr_tresor_2024_hdf = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_sstr_tresor_2025_hdf = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_sstr_ss_2024_hdf = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_sstr_ss_2025_hdf = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_mbe_2024_hdf = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_mbe_2025_hdf = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

greffes_hdf = [
    "FORT-DE-FRANCE"
]

df_div_2024_hdf = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_hdf)].reset_index(drop=True)
df_div_2025_hdf = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_hdf)].reset_index(drop=True)

df_dpca_2024_hdf = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]
df_dpca_2025_hdf = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Martinique "]

df_ts_2024_hdf = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Martinique ") |
    (df_ts_2024["REGIONTO"] == "Martinique ")
]

df_ts_2025_hdf = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Martinique ") |
    (df_ts_2025["REGIONTO"] == "Martinique ")
]

df_prev_2024_hdf = df_prev_2024[df_prev_2024["REGION"] == "Martinique "]
df_prev_2025_hdf = df_prev_2025[df_prev_2025["REGION"] == "Martinique "]

# === SLIDE 8 ===
slide8 = prs.slides[7]

A_2024 = df_immat_2024_hdf["SIREN"].notna().sum()
A_2025 = df_immat_2025_hdf["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_hdf["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_hdf["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_hdf["SIREN"].notna().sum()
C_2025 = df_pcl_2025_hdf["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_hdf["SIREN"].notna().sum()
D_2025 = df_radiation_2025_hdf["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_hdf["SIREN"].notna().sum()
E_2025 = df_ip_2025_hdf["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_hdf["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_hdf["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_hdf["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_hdf["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_hdf, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_hdf, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_hdf, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_hdf, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_hdf, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_hdf, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_hdf, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_hdf, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_hdf, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_hdf, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer dans la slide 8
for shape in slide8.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 9 ===
slide9 = prs.slides[8]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide9 = df_mbe_2024_hdf["SIREN"].notna().sum()
A_2025_slide9 = df_mbe_2025_hdf["SIREN"].notna().sum()
A_val_slide9 = f"{A_2025_slide9:,}".replace(",", " ")
A_pct_slide9 = evol_percent_txt(A_2024_slide9, A_2025_slide9)

# B : DIV
B_2024_slide9 = df_div_2024_hdf["SIREN"].notna().sum()
B_2025_slide9 = df_div_2025_hdf["SIREN"].notna().sum()
B_val_slide9 = f"{B_2025_slide9:,}".replace(",", " ")
B_pct_slide9 = evol_percent_txt(B_2024_slide9, B_2025_slide9)

# C : DPCA
C_2024_slide9 = df_dpca_2024_hdf["SIREN"].notna().sum()
C_2025_slide9 = df_dpca_2025_hdf["SIREN"].notna().sum()
C_val_slide9 = f"{C_2025_slide9:,}".replace(",", " ")
C_pct_slide9 = evol_percent_txt(C_2024_slide9, C_2025_slide9)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Martinique ").sum()
departs = (df_ts_2025["REGIONFROM"] == "Martinique ").sum()
solde = arrivees - departs

D_val_slide9 = f"{arrivees:,}".replace(",", " ")
E_val_slide9 = f"{departs:,}".replace(",", " ")
F_val_slide9 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 9 ===
for shape in slide9.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide9
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide9, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide9, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide9, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide9, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide9, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide9, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide9, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide9, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide9, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# REGION MAYOTTE

# === FILTRAGE Mayotte ===
df_immat_2024_bfc = df_immat_2024[df_immat_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_immat_2025_bfc = df_immat_2025[df_immat_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_pcl_2024_bfc = df_pcl_2024[df_pcl_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_pcl_2025_bfc = df_pcl_2025[df_pcl_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_radiation_2024_bfc = df_radiation_2024[df_radiation_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_radiation_2025_bfc = df_radiation_2025[df_radiation_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_stock_2025_bfc = df_stock_2025[df_stock_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"].reset_index(drop=True)

df_ip_2024_bfc = df_ip_2024[df_ip_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_ip_2025_bfc = df_ip_2025[df_ip_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_sstr_tresor_2024_bfc = df_sstr_tresor_2024[df_sstr_tresor_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_sstr_tresor_2025_bfc = df_sstr_tresor_2025[df_sstr_tresor_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_sstr_ss_2024_bfc = df_sstr_ss_2024[df_sstr_ss_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_sstr_ss_2025_bfc = df_sstr_ss_2025[df_sstr_ss_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_mbe_2024_bfc = df_mbe_2024[df_mbe_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_mbe_2025_bfc = df_mbe_2025[df_mbe_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

greffes_bfc = [
    "MAMOUDZOU"
]

df_div_2024_bfc = df_div_2024[df_div_2024["Nom du greffe"].isin(greffes_bfc)].reset_index(drop=True)
df_div_2025_bfc = df_div_2025[df_div_2025["Nom du greffe"].isin(greffes_bfc)].reset_index(drop=True)

df_dpca_2024_bfc = df_dpca_2024[df_dpca_2024["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]
df_dpca_2025_bfc = df_dpca_2025[df_dpca_2025["ADRESSEFRANCEREGIONLIBELLE"] == "Mayotte"]

df_ts_2024_bfc = df_ts_2024[
    (df_ts_2024["REGIONFROM"] == "Mayotte") |
    (df_ts_2024["REGIONTO"] == "Mayotte")
]

df_ts_2025_bfc = df_ts_2025[
    (df_ts_2025["REGIONFROM"] == "Mayotte") |
    (df_ts_2025["REGIONTO"] == "Mayotte")
]

df_prev_2024_bfc = df_prev_2024[df_prev_2024["REGION"] == "Mayotte"]
df_prev_2025_bfc = df_prev_2025[df_prev_2025["REGION"] == "Mayotte"]

# === SLIDE 10 ===
slide10 = prs.slides[9]

A_2024 = df_immat_2024_bfc["SIREN"].notna().sum()
A_2025 = df_immat_2025_bfc["SIREN"].notna().sum()
var_A = ((A_2025 - A_2024) / A_2024) * 100 if A_2024 > 0 else 0
A_val = f"{A_2025:,}".replace(",", " ")
A_var = f"{var_A:+.1f}%"

stock_2025 = df_stock_2025_bfc["INSCRIPTIONS_PRINCIPALES"].sum() + df_stock_2025_bfc["INSCRIPTIONS_SECONDAIRES"].sum()
B_val = f"{int(stock_2025):,}".replace(",", " ")

C_2024 = df_pcl_2024_bfc["SIREN"].notna().sum()
C_2025 = df_pcl_2025_bfc["SIREN"].notna().sum()
var_C = ((C_2025 - C_2024)/C_2024)*100 if C_2024>0 else 0
C_val = f"{C_2025:,}".replace(",", " ")
C_var = f"{var_C:+.1f}%"

D_2024 = df_radiation_2024_bfc["SIREN"].notna().sum()
D_2025 = df_radiation_2025_bfc["SIREN"].notna().sum()
var_D = ((D_2025 - D_2024)/D_2024)*100 if D_2024>0 else 0
D_val = f"{D_2025:,}".replace(",", " ")
D_var = f"{var_D:+.1f}%"

E_2024 = df_ip_2024_bfc["SIREN"].notna().sum()
E_2025 = df_ip_2025_bfc["SIREN"].notna().sum()
var_E = ((E_2025 - E_2024) / E_2024) * 100 if E_2024>0 else 0
E_val = f"{E_2025:,}".replace(",", " ")
E_var = f"{var_E:+.1f}%"

W_2024 = df_sstr_tresor_2024_bfc["SIREN"].notna().sum()
W_2025 = df_sstr_tresor_2025_bfc["SIREN"].notna().sum()
var_W = ((W_2025 - W_2024) / W_2024) * 100 if W_2024>0 else 0
W_val = f"{W_2025:,}".replace(",", " ")
W_var = f"{var_W:+.1f}%"

G_2024 = df_sstr_ss_2024_bfc["SIREN"].notna().sum()
G_2025 = df_sstr_ss_2025_bfc["SIREN"].notna().sum()
var_G = ((G_2025 - G_2024) / G_2024) * 100 if G_2024>0 else 0
G_val = f"{G_2025:,}".replace(",", " ")
G_var = f"{var_G:+.1f}%"

def sum_prev(df, label):
    mask = df["STATISTIQUE"].astype(str).str.strip().str.lower() == label.lower()
    return df.loc[mask, "NOMBRE"].sum()

H_2024 = sum_prev(df_prev_2024_bfc, "Nombre total de dossiers de prévention étudiés")
H_2025 = sum_prev(df_prev_2025_bfc, "Nombre total de dossiers de prévention étudiés")
H_var = ((H_2025 - H_2024)/H_2024*100) if H_2024>0 else 0
H_val = f"{int(H_2025):,}".replace(",", " ")
H_var_txt = f"{H_var:+.1f}%"

J_2024 = sum_prev(df_prev_2024_bfc, "dont demandes d'entretien spontanées")
J_2025 = sum_prev(df_prev_2025_bfc, "dont demandes d'entretien spontanées")
J_var = ((J_2025 - J_2024)/J_2024*100) if J_2024>0 else 0
J_val = f"{int(J_2025):,}".replace(",", " ")
J_var_txt = f"{J_var:+.1f}%"

K_2024 = sum_prev(df_prev_2024_bfc, "dont convocations suite à informations internes")
K_2025 = sum_prev(df_prev_2025_bfc, "dont convocations suite à informations internes")
K_var = ((K_2025 - K_2024)/K_2024*100) if K_2024>0 else 0
K_val = f"{int(K_2025):,}".replace(",", " ")
K_var_txt = f"{K_var:+.1f}%"

L_2024 = sum_prev(df_prev_2024_bfc, "Concilliations  (ouverture)")
L_2025 = sum_prev(df_prev_2025_bfc, "Concilliations  (ouverture)")
L_var = ((L_2025 - L_2024)/L_2024*100) if L_2024>0 else 0
L_val = f"{int(L_2025):,}".replace(",", " ")
L_var_txt = f"{L_var:+.1f}%"

O_2024 = sum_prev(df_prev_2024_bfc, "mandat ad hoc (ouverture)")
O_2025 = sum_prev(df_prev_2025_bfc, "mandat ad hoc (ouverture)")
O_var = ((O_2025 - O_2024)/O_2024*100) if O_2024>0 else 0
O_val = f"{int(O_2025):,}".replace(",", " ")
O_var_txt = f"{O_var:+.1f}%"

def evol_percent_txt(v2024, v2025):
    if v2024 == 0:
        return "N/A"
    diff = ((v2025 - v2024) / v2024) * 100
    return f"{diff:+.1f}%"

token_map = {
    "A": (A_val, A_var),
    "B": (B_val, ""),
    "C": (C_val, C_var),
    "D": (D_val, D_var),
    "E": (E_val, E_var),
    "W": (W_val, W_var),
    "G": (G_val, G_var),
    "H": (H_val, H_var_txt),
    "J": (J_val, J_var_txt),
    "K": (K_val, K_var_txt),
    "L": (L_val, L_var_txt),
    "O": (O_val, O_var_txt)
}

# Remplacer sur la slide 10
for shape in slide10.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            orig = run.text or ""
            text = orig.strip()
            replaced = False

            if text in token_map:
                run.text = token_map[text][0]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True
            elif text.endswith("%") and text[:-1] in token_map:
                token = text[:-1]
                run.text = token_map[token][1]
                paragraph.alignment = PP_ALIGN.CENTER
                replaced = True

            if not replaced:
                new_text = orig
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}%\b', val_pct, new_text)
                for tok, (val_abs, val_pct) in token_map.items():
                    new_text = re.sub(rf'\b{re.escape(tok)}\b', val_abs, new_text)
                if new_text != orig:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER

# === SLIDE 11 ===
slide11 = prs.slides[10]

# === Recalcul des valeurs ===
# A : MBE
A_2024_slide11 = df_mbe_2024_bfc["SIREN"].notna().sum()
A_2025_slide11 = df_mbe_2025_bfc["SIREN"].notna().sum()
A_val_slide11 = f"{A_2025_slide11:,}".replace(",", " ")
A_pct_slide11 = evol_percent_txt(A_2024_slide11, A_2025_slide11)

# B : DIV
B_2024_slide11 = df_div_2024_bfc["SIREN"].notna().sum()
B_2025_slide11 = df_div_2025_bfc["SIREN"].notna().sum()
B_val_slide11 = f"{B_2025_slide11:,}".replace(",", " ")
B_pct_slide11 = evol_percent_txt(B_2024_slide11, B_2025_slide11)

# C : DPCA
C_2024_slide11 = df_dpca_2024_bfc["SIREN"].notna().sum()
C_2025_slide11 = df_dpca_2025_bfc["SIREN"].notna().sum()
C_val_slide11 = f"{C_2025_slide11:,}".replace(",", " ")
C_pct_slide11 = evol_percent_txt(C_2024_slide11, C_2025_slide11)

# D, E, F : arrivées, départs, solde (inchangés)
arrivees = (df_ts_2025["REGIONTO"] == "Mayotte").sum()
departs = (df_ts_2025["REGIONFROM"] == "Mayotte").sum()
solde = arrivees - departs

D_val_slide11 = f"{arrivees:,}".replace(",", " ")
E_val_slide11 = f"{departs:,}".replace(",", " ")
F_val_slide11 = f"{solde:,}".replace(",", " ")

# === Remplacement du texte dans la slide 11 ===
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = (run.text or "").strip()

            # Remplacement exact
            if text == "A":
                run.text = A_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "A%":
                run.text = A_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "B":
                run.text = B_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "W%":
                run.text = B_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C":
                run.text = C_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "C%":
                run.text = C_pct_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "D":
                run.text = D_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "E":
                run.text = E_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            elif text == "F":
                run.text = F_val_slide11
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Remplacement dans du texte
                new_text = text
                new_text = re.sub(r'\bA%\b', A_pct_slide11, new_text)
                new_text = re.sub(r'\bA\b', A_val_slide11, new_text)
                new_text = re.sub(r'\bW%\b', B_pct_slide11, new_text)
                new_text = re.sub(r'\bB\b', B_val_slide11, new_text)
                new_text = re.sub(r'\bC%\b', C_pct_slide11, new_text)
                new_text = re.sub(r'\bC\b', C_val_slide11, new_text)
                new_text = re.sub(r'\bD\b', D_val_slide11, new_text)
                new_text = re.sub(r'\bE\b', E_val_slide11, new_text)
                new_text = re.sub(r'\bF\b', F_val_slide11, new_text)
                if new_text != text:
                    run.text = new_text
                    paragraph.alignment = PP_ALIGN.CENTER


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace("X", "3").replace("x", "3")
                    run.text = run.text.replace("tribunau3", "tribunaux").replace("TRIBUNAU3", "TRIBUNAUX")
                    run.text = run.text.replace("sociau3", "sociaux").replace("SOCIAU3", "SOCIAUX")
                    run.text = run.text.replace("avril", "juillet")
                    run.text = run.text.replace("juin", "septembre")

#   SAUVEGARDE 
output_filename = os.path.join(base_path, "Stats_DROM_new.pptx")
prs.save(output_filename)
print(f"Fichier PowerPoint modifié enregistré sous : {output_filename}") 